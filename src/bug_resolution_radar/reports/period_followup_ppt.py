"""Template-based fortnight follow-up PPT report."""

from __future__ import annotations

import logging
import re
import unicodedata
from copy import deepcopy
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any, List, Mapping, Optional, Sequence, cast
from uuid import uuid4

import pandas as pd
import plotly.graph_objects as go
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from bug_resolution_radar.analytics.analysis_window import apply_analysis_depth_filter
from bug_resolution_radar.analytics.finalist_discrepancies import (
    apply_effective_finalist_lookup_state,
    apply_effective_finalist_lookup_state_for_scope,
    split_root_cause_evolutive_discrepancies,
)
from bug_resolution_radar.analytics.finalist_discrepancy_lists import (
    FinalistDiscrepancyIssueRow,
    build_finalist_discrepancy_issue_list,
)
from bug_resolution_radar.analytics.insights import (
    build_theme_color_map,
    build_theme_fortnight_trend,
    build_theme_render_order,
    order_theme_labels_by_volume,
)
from bug_resolution_radar.analytics.issues import normalize_text_col, priority_rank
from bug_resolution_radar.analytics.kpis import (
    OPEN_AGE_BUCKET_LABELS,
    build_open_age_priority_payload,
    compute_kpis,
)
from bug_resolution_radar.analytics.period_functionality_followup import (
    FunctionalityIssueRow,
    FunctionalityTopRow,
    FunctionalityZoomSlide,
    PeriodFunctionalityFollowupSummary,
    build_period_functionality_followup_summary,
)
from bug_resolution_radar.analytics.period_risk_issue_lists import (
    PeriodRiskIssueRow,
    build_period_risk_issue_lists,
)
from bug_resolution_radar.analytics.period_summary import (
    OPEN_ISSUES_FOCUS_MODE_MAESTRAS,
    QuincenalDelta,
    QuincenalScopeResult,
    build_country_quincenal_result,
    format_window_label,
    scope_country_sources,
    source_label_map,
)
from bug_resolution_radar.analytics.status_semantics import effective_closed_mask
from bug_resolution_radar.analytics.time_windows import TimeWindowService
from bug_resolution_radar.analytics.trend_charts import ChartContext, build_trends_registry
from bug_resolution_radar.analytics.trend_insights import build_trend_insight_pack
from bug_resolution_radar.common.issue_links import (
    HELIX_ID_RE,
    build_issue_url_maps,
    linkify_issue_references,
)
from bug_resolution_radar.config import Settings, jira_sources, resolve_period_ppt_template_path
from bug_resolution_radar.reports.executive_ppt import _fig_to_png, _kaleido_png_bytes
from bug_resolution_radar.reports.period_followup_layout import (
    PERIOD_FOLLOWUP_LAYOUT,
    KpiRow,
    KpiSideMetric,
    apply_text_frame_margins,
    delta_badge_font_size,
    iter_out_of_viewport_shapes,
    metric_card_typography,
)
from bug_resolution_radar.reports.pptx_native_tables import (
    ellipsize_text,
    native_column_widths,
    populate_native_table,
    rebuild_native_table_shape,
)
from bug_resolution_radar.repositories.issues_store import load_issues_df
from bug_resolution_radar.services.notes import NotesStore
from bug_resolution_radar.theme.design_tokens import (
    BBVA_FONT_HEADLINE_PPT,
    BBVA_FONT_SANS_BOOK_PPT,
    BBVA_FONT_SANS_MEDIUM_PPT,
    BBVA_LIGHT,
    BBVA_REPORT_AMBER_BG,
    BBVA_REPORT_AMBER_BORDER,
    BBVA_REPORT_AMBER_TEXT,
    BBVA_REPORT_RED_BG,
    BBVA_REPORT_RED_BORDER,
    BBVA_REPORT_RED_TEXT,
    EXEC_CHART_AXIS_FONT_PT,
    EXEC_CHART_AXIS_TITLE_FONT_PT,
    EXEC_CHART_EXPORT_HEIGHT,
    EXEC_CHART_EXPORT_WIDTH,
    EXEC_CHART_INSIDE_VALUE_FONT_PT,
    EXEC_CHART_LEGEND_FONT_PT,
    EXEC_CHART_MARGIN,
    EXEC_CHART_TOTAL_FONT_PT,
    EXEC_CHART_TREND_EXPORT_HEIGHT,
    hex_to_rgb,
)
from bug_resolution_radar.theme.semantic_colors import priority_color_map

LOGGER = logging.getLogger(__name__)

_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_EMU_PER_INCH = 914400.0
_FUNCTIONALITY_TEMPLATE_FILENAME = "Seguimiento de incidencias por funcionalidad.pptx"
_TABLE_BODY_FG_RGB = hex_to_rgb(BBVA_LIGHT.core_blue)
_ISSUE_TABLE_ROWS_PER_SLIDE = 5
_ISSUE_TABLE_LEFT = Inches(0.52)
_ISSUE_TABLE_TOP = Inches(1.48)
_ISSUE_TABLE_WIDTH = Inches(8.96)
_ISSUE_TABLE_HEADER_HEIGHT = Inches(0.34)
_ISSUE_TABLE_ROW_HEIGHT = Inches(0.62)
_ISSUE_TABLE_COMMENT_ROW_HEIGHT = int(_ISSUE_TABLE_ROW_HEIGHT * 1.55)
_ISSUE_TABLE_COMMENT_ROW_UNITS = float(_ISSUE_TABLE_COMMENT_ROW_HEIGHT) / float(
    _ISSUE_TABLE_ROW_HEIGHT
)
_ISSUE_COMMENT_CHUNK_CHARS = 240
_ISSUE_TABLE_FONT_NAME = "Arial"
_ISSUE_TABLE_BODY_FONT_SIZE_PT = 9.2
_ISSUE_TABLE_HEADER_FONT_SIZE_PT = 8.4
_ISSUE_TABLE_COLUMN_WEIGHTS: tuple[float, ...] = (15.0, 33.0, 17.0, 13.0, 11.0, 11.0)
_FUNCTIONALITY_ISSUE_TABLE_HEADERS: tuple[str, ...] = (
    "ID",
    "Descripción",
    "Funcionalidad/\nCausa raíz",
    "Estado",
    "Criticidad",
    "Días abierta",
)
_RISK_ASSIGNEE_TABLE_HEADERS: tuple[str, ...] = (
    "ID",
    "Descripción",
    "Responsable",
    "Estado",
    "Criticidad",
    "Días abierta",
)


@dataclass(frozen=True)
class _IssuePageItem:
    issue: Any
    comment_chunks: tuple[str, ...] = ()


@dataclass(frozen=True)
class _FinalistIssuePageItem:
    issue: FinalistDiscrepancyIssueRow
    comment_chunks: tuple[str, ...] = ()


_FUNCTIONALITY_DASHBOARD_TABLE_HEADERS: tuple[str, ...] = (
    "#",
    "Resto incidencias abiertas",
    "Nuevas",
    "Agregadas",
    "Días promedio abiertas",
)
_FUNCTIONALITY_DASHBOARD_TABLE_COLUMN_WEIGHTS: tuple[float, ...] = (0.68, 2.88, 1.02, 1.32, 2.50)
_FUNCTIONALITY_DASHBOARD_TABLE_ROWS = 4
_FUNCTIONALITY_TABLE_GAP_TOP_EMU = 120_000
_FUNCTIONALITY_TABLE_GAP_RIGHT_EMU = 185_000
_FUNCTIONALITY_TABLE_BOTTOM_GAP_EMU = 24_000
_FUNCTIONALITY_TABLE_FONT_NAME = "Arial"
_FUNCTIONALITY_TABLE_BODY_FONT_SIZE_PT = 10.0
_FUNCTIONALITY_TABLE_HEADER_FONT_SIZE_PT = 8.4
_SUMMARY_DELTA_FONT_SIZE_PT = 8.8
_SUMMARY_DELTA_UP_RGB = RGBColor(201, 67, 77)
_SUMMARY_DELTA_DOWN_RGB = RGBColor(62, 133, 64)
_SUMMARY_DELTA_NEUTRAL_RGB = RGBColor(95, 112, 142)
_EXEC_BG_RGB = (7, 36, 96)
_EXEC_TEXT_PRIMARY_RGB = (247, 251, 255)
_EXEC_TEXT_SECONDARY_RGB = (170, 191, 226)
_EXEC_ACCENT_BORDER_RGB = (104, 151, 222)
_EXEC_CARD_BG_RGB = (12, 52, 118)
_EXEC_CARD_TITLE_RGB = (186, 226, 252)
_PPT_FONT_HEAD = BBVA_FONT_HEADLINE_PPT
_PPT_FONT_BODY = BBVA_FONT_SANS_BOOK_PPT
_PPT_FONT_BODY_MEDIUM = BBVA_FONT_SANS_MEDIUM_PPT
_COVER_TITLE_TEXT = "Seguimiento incidencias"
_COVER_REMOVED_SUBTITLE = "KPIs, evolución y análisis del periodo"
_FUNCTIONALITY_TREND_AGGREGATE_TITLE = (
    "Tendencia por funcionalidad : vista agregada ultimo semestre"
)
_RISK_HIGH_PRIORITY_ORDER_NOTE = (
    "Detalle - incidencias CRÍTICAS - ordenado por 1º : Criticidad, 2º: Días abierta y 3º: Estado"
)
_RISK_AGED_ORDER_NOTE = (
    "Detalle - TODAS las incidencias abiertas - ordenado por 1º : Días abierta, "
    "2º: Criticidad y 3º: Estado"
)
_FINALIST_DISCREPANCIES_TITLE = "Incidencias con discrepancias en estado finalista"
_FINALIST_DISCREPANCIES_ORDER_NOTE = (
    "Detalle - Helix en estado finalista y Jira pendiente - ordenado por 1º : Criticidad, "
    "2º: Días abierta, 3º: Estado y 4º: Helix ID"
)
_ROOT_CAUSE_EVOLUTIVES_TITLE = "Evolutivos para solucionar causas raíces"
_ROOT_CAUSE_EVOLUTIVES_ORDER_NOTE = (
    "Detalle - Helix en estado finalista, Jira pendiente y label de causa raíz configurada - "
    "ordenado por 1º : Criticidad, 2º: Días abierta, 3º: Estado y 4º: Helix ID"
)


@dataclass(frozen=True)
class PeriodFollowupReportResult:
    file_name: str
    content: bytes
    slide_count: int
    total_issues: int
    open_issues: int
    closed_issues: int
    country: str
    source_ids: tuple[str, ...]
    applied_filter_summary: str


def _slide_text_blob(slide: Any) -> str:
    """Collect lower-cased text from all text-capable shapes in a slide."""
    chunks: list[str] = []
    for shape in getattr(slide, "shapes", []):
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = str(getattr(shape, "text", "") or "").strip()
        if txt:
            chunks.append(txt.lower())
    return " ".join(chunks).strip()


def _looks_like_explanatory_helper_slide(slide: Any) -> bool:
    """Heuristic to detect optional helper/instruction slide in position 2."""
    blob = _slide_text_blob(slide)
    if not blob:
        return True

    production_tokens = (
        "dashboard",
        "seguimiento de incidencias",
        "seguimiento de kpis",
        "gráficos de evolución",
        "graficos de evolucion",
    )
    if any(tok in blob for tok in production_tokens):
        return False

    helper_tokens = (
        "instrucci",
        "comentario",
        "comentarios",
        "helper",
        "ayuda",
        "ejemplo",
        "plantilla",
        "template",
        "no editar",
        "borrar",
        "delete",
    )
    return any(tok in blob for tok in helper_tokens)


def _ensure_slide_index(prs: Any, *, index: int, role: str) -> None:
    if int(index) < len(prs.slides):
        return
    raise ValueError(
        "Plantilla de seguimiento inválida: falta la slide "
        f"{int(index) + 1} ({role}). Revisa PERIOD_PPT_TEMPLATE_PATH."
    )


def _normalize_period_template(prs: Any) -> None:
    """
    Normalize user-provided template into the 8-slide structure expected by renderer.

    Target structure:
      1) Portada
      2) Dashboard header
      3) Resumen país
      4) Resumen origen A
      5) Resumen origen B
      6) Header evolución
      7) Evolución origen A
      8) Evolución origen B
    """
    if len(prs.slides) < 7:
        raise ValueError(
            "La plantilla de seguimiento debe tener al menos 7 slides base "
            "(portada, dashboard, 3 resúmenes, cabecera evolución y 1 evolución). "
            f"Slides detectadas: {len(prs.slides)}."
        )

    # Optional helper slide used in earlier corporate templates.
    if len(prs.slides) > 1 and _looks_like_explanatory_helper_slide(prs.slides[1]):
        _remove_slide(prs, 1)

    if len(prs.slides) < 7:
        raise ValueError(
            "La plantilla de seguimiento quedó incompleta tras eliminar la slide de ayuda. "
            f"Slides detectadas: {len(prs.slides)}."
        )

    # Ensure we always have two evolution slides based on the same visual template.
    if len(prs.slides) >= 8:
        _copy_slide_content(prs, source_index=6, target_index=7)
    else:
        _append_slide_clone(prs, source_index=6)

    # Keep exactly 8 slides, preserving canonical order.
    while len(prs.slides) > 8:
        _remove_slide(prs, 8)

    required_roles = {
        0: "Portada",
        2: "Resumen país",
        3: "Resumen origen A",
        4: "Resumen origen B",
        6: "Evolución origen A",
        7: "Evolución origen B",
    }
    for idx, role in required_roles.items():
        _ensure_slide_index(prs, index=idx, role=role)


def _slug(value: str) -> str:
    txt = str(value or "").strip().lower()
    txt = re.sub(r"[^a-z0-9]+", "-", txt).strip("-")
    return txt or "scope"


def _parse_bool_flag(value: object, *, default: bool = False) -> bool:
    token = str(value or "").strip().lower()
    if not token:
        return bool(default)
    if token in {"1", "true", "t", "yes", "y", "on"}:
        return True
    if token in {"0", "false", "f", "no", "n", "off"}:
        return False
    return bool(default)


def _fmt_days(value: float | None) -> int:
    if value is None or pd.isna(value):
        return 0
    return max(0, int(round(float(value))))


def _summary_delta_badge(delta: QuincenalDelta) -> tuple[str, RGBColor]:
    tone = str(getattr(delta, "presentation_semantic_tone", "") or "").strip().lower()
    if tone == "flow":
        color = _SUMMARY_DELTA_DOWN_RGB
    elif tone == "risk":
        color = _SUMMARY_DELTA_UP_RGB
    else:
        color = _SUMMARY_DELTA_NEUTRAL_RGB
    return (str(getattr(delta, "presentation_badge_text", "") or "•0%"), color)


def _clean_source_ids(source_ids: Sequence[str]) -> List[str]:
    out: List[str] = []
    for raw in list(source_ids or []):
        sid = str(raw or "").strip()
        if sid and sid not in out:
            out.append(sid)
    return out


def _resolve_template_path(settings: Settings, explicit_path: str | None = None) -> Path:
    return resolve_period_ppt_template_path(settings, explicit_path=explicit_path)


def _safe_emu(value: Any, *, default: int) -> int:
    try:
        if value is None:
            return int(default)
        return int(cast(Any, value))
    except Exception:
        return int(default)


def _resolve_functionality_template_path() -> Path:
    path = (
        Path(__file__).resolve().parent / "templates" / _FUNCTIONALITY_TEMPLATE_FILENAME
    ).resolve()
    if path.exists() and path.is_file():
        return path
    raise FileNotFoundError(
        f"No se encontró la plantilla de seguimiento por funcionalidad. Ruta esperada: {path}"
    )


def _remove_slide(prs: Any, index: int) -> None:
    sld_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(sld_id.rId)
    del prs.slides._sldIdLst[index]


def _copy_slide_content(prs: Any, *, source_index: int, target_index: int) -> None:
    source = prs.slides[source_index]
    target = prs.slides[target_index]

    for shape in list(target.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    for rel in list(target.part.rels.values()):
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype or "comments" in rel.reltype:
            continue
        target.part.drop_rel(rel.rId)

    rid_map: dict[str, str] = {}
    for rel in source.part.rels.values():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype or "comments" in rel.reltype:
            continue
        rel_target = rel.target_ref if rel.is_external else rel._target
        rid_map[rel.rId] = target.part.rels._add_relationship(
            rel.reltype,
            rel_target,
            rel.is_external,
        )

    for shape in source.shapes:
        clone = deepcopy(shape.element)
        for node in clone.iter():
            for attr_name, attr_value in list(node.attrib.items()):
                if attr_name.startswith(_REL_NS) and attr_value in rid_map:
                    node.set(attr_name, rid_map[attr_value])
        target.shapes._spTree.insert_element_before(clone, "p:extLst")


def _append_slide_clone(prs: Any, *, source_index: int) -> None:
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shape in list(dest.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    rid_map: dict[str, str] = {}
    for rel in source.part.rels.values():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype or "comments" in rel.reltype:
            continue
        rel_target = rel.target_ref if rel.is_external else rel._target
        rid_map[rel.rId] = dest.part.rels._add_relationship(
            rel.reltype,
            rel_target,
            rel.is_external,
        )

    for shape in source.shapes:
        clone = deepcopy(shape.element)
        for node in clone.iter():
            for attr_name, attr_value in list(node.attrib.items()):
                if attr_name.startswith(_REL_NS) and attr_value in rid_map:
                    node.set(attr_name, rid_map[attr_value])
        dest.shapes._spTree.insert_element_before(clone, "p:extLst")


def _move_slide(prs: Any, *, from_index: int, to_index: int) -> None:
    if int(from_index) == int(to_index):
        return
    slides = prs.slides._sldIdLst
    sld_id = slides[int(from_index)]
    del slides[int(from_index)]
    slides.insert(int(to_index), sld_id)


def _append_slide_clone_from_source(prs: Any, *, source_slide: Any) -> Any:
    dest = prs.slides.add_slide(prs.slide_layouts[6])

    for shape in list(dest.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    rid_map: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype or "comments" in rel.reltype:
            continue
        if str(rel.reltype or "").endswith("/image") and not rel.is_external:
            try:
                img_blob = bytes(getattr(rel._target, "blob", b"") or b"")
                if img_blob:
                    _, img_rid = dest.part.get_or_add_image_part(BytesIO(img_blob))
                    rid_map[rel.rId] = img_rid
                    continue
            except Exception:
                pass
        rel_target = rel.target_ref if rel.is_external else rel._target
        rid_map[rel.rId] = dest.part.rels._add_relationship(
            rel.reltype,
            rel_target,
            rel.is_external,
        )

    for shape in source_slide.shapes:
        clone = deepcopy(shape.element)
        for node in clone.iter():
            for attr_name, attr_value in list(node.attrib.items()):
                if attr_name.startswith(_REL_NS) and attr_value in rid_map:
                    node.set(attr_name, rid_map[attr_value])
        dest.shapes._spTree.insert_element_before(clone, "p:extLst")
    _apply_effective_background_from_source(dest_slide=dest, source_slide=source_slide)
    return dest


def _solid_background_rgb(shape_container: Any) -> RGBColor | None:
    if shape_container is None:
        return None
    try:
        fill = shape_container.background.fill
    except Exception:
        return None
    try:
        if int(fill.type or 0) != 1:  # SOLID
            return None
    except Exception:
        return None
    try:
        rgb = getattr(fill.fore_color, "rgb", None)
    except Exception:
        rgb = None
    if rgb is None:
        return None
    return cast(RGBColor, rgb)


def _effective_background_rgb(source_slide: Any) -> RGBColor:
    for container in (
        source_slide,
        getattr(source_slide, "slide_layout", None),
        getattr(getattr(source_slide, "slide_layout", None), "slide_master", None),
    ):
        rgb = _solid_background_rgb(container)
        if rgb is not None:
            return rgb
    return RGBColor(247, 248, 248)


def _apply_effective_background_from_source(*, dest_slide: Any, source_slide: Any) -> None:
    rgb = _effective_background_rgb(source_slide)
    try:
        fill = dest_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb
    except Exception:
        return


def _shape_or_none(slide: Any, index_1_based: int) -> Any | None:
    idx = int(index_1_based) - 1
    if idx < 0 or idx >= len(slide.shapes):
        return None
    return slide.shapes[idx]


def _shape_area_in2(shape: Any) -> float:
    return float(shape.width) * float(shape.height) / (_EMU_PER_INCH * _EMU_PER_INCH)


def _picture_candidates(slide: Any, *, min_area_in2: float = 1.0) -> List[Any]:
    out: List[Any] = []
    for shape in slide.shapes:
        try:
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        if _shape_area_in2(shape) < float(min_area_in2):
            continue
        out.append(shape)
    return out


def _remove_shape(shape: Any) -> None:
    try:
        node = shape.element
        node.getparent().remove(node)
    except Exception:
        return


def _set_shape_text_fit(shape: Any) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _set_shape_text(slide: Any, index_1_based: int, text: str) -> None:
    shape = _shape_or_none(slide, index_1_based)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    _set_shape_text_by_shape(shape, text)


def _set_shape_text_by_shape(shape: Any, text: str) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    target_lines = str(text or "").splitlines() or [""]
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    while len(paragraphs) < len(target_lines):
        tf.add_paragraph()
        paragraphs = list(tf.paragraphs)

    for idx, line in enumerate(target_lines):
        p = paragraphs[idx]
        runs = list(p.runs)
        if not runs:
            p.add_run()
            runs = list(p.runs)
        runs[0].text = str(line)
        for run in runs[1:]:
            run.text = ""

    for idx in range(len(target_lines), len(paragraphs)):
        p = paragraphs[idx]
        runs = list(p.runs)
        if not runs:
            p.add_run()
            runs = list(p.runs)
        runs[0].text = ""
        for run in runs[1:]:
            run.text = ""

    _set_shape_text_fit(shape)


def _set_shape_text_strict(slide: Any, index_1_based: int, text: str) -> None:
    shape = _shape_or_none(slide, index_1_based)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    _set_shape_text_strict_by_shape(shape, text)


def _set_shape_text_strict_by_shape(shape: Any, text: str) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    sample_run = None
    try:
        sample_run = tf.paragraphs[0].runs[0]
    except Exception:
        sample_run = None

    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text or "")
    if sample_run is not None:
        try:
            run.font.bold = sample_run.font.bold
        except Exception:
            pass
        try:
            run.font.italic = sample_run.font.italic
        except Exception:
            pass
        try:
            run.font.size = sample_run.font.size
        except Exception:
            pass
        try:
            run.font.name = sample_run.font.name
        except Exception:
            pass
        try:
            rgb = getattr(getattr(sample_run.font, "color", None), "rgb", None)
            if rgb is not None:
                run.font.color.rgb = rgb
        except Exception:
            pass
    _set_shape_text_fit(shape)


def _shape_table_or_none(slide: Any, index_1_based: int) -> Any | None:
    shape = _shape_or_none(slide, index_1_based)
    if shape is not None and getattr(shape, "has_table", False):
        return shape
    table_shapes = [item for item in slide.shapes if getattr(item, "has_table", False)]
    if not table_shapes:
        return None
    return max(table_shapes, key=_shape_area_in2)


def _trim_text(value: object, *, max_chars: int) -> str:
    txt = str(value or "").strip()
    if max_chars <= 0 or len(txt) <= max_chars:
        return txt
    return txt[: max(0, max_chars - 3)].rstrip() + "..."


_HELIX_INC_RE = re.compile(r"\binc\d{5,}[a-z0-9-]*\b", flags=re.IGNORECASE)
_ALPHANUM_KEY_RE = re.compile(r"\b[a-z][a-z0-9]+-\d+[a-z0-9-]*\b", flags=re.IGNORECASE)


def _premium_sentence_case(value: object) -> str:
    raw = re.sub(r"\s+", " ", str(value or "").strip())
    if not raw:
        return ""
    clean = re.sub(r"\s*/\s*", " / ", raw).lower()
    clean = re.sub(
        r"^([\s\"'¿¡\(\[]*)([a-záéíóúñ])",
        lambda match: f"{match.group(1)}{match.group(2).upper()}",
        clean,
        count=1,
    )
    clean = re.sub(
        r"([.!?]\s+)([a-záéíóúñ])",
        lambda match: f"{match.group(1)}{match.group(2).upper()}",
        clean,
    )
    clean = _HELIX_INC_RE.sub(lambda match: str(match.group(0)).upper(), clean)
    clean = _ALPHANUM_KEY_RE.sub(lambda match: str(match.group(0)).upper(), clean)
    return clean


def _set_shape_font_size(
    slide: Any,
    *,
    shape_index: int,
    font_size_pt: float,
    bold: bool | None = None,
    disable_autofit: bool = False,
) -> None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    if disable_autofit:
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        try:
            tf.word_wrap = False
        except Exception:
            pass
    for paragraph in list(tf.paragraphs):
        for run in list(paragraph.runs):
            run.font.size = Pt(float(font_size_pt))
            if bold is not None:
                run.font.bold = bool(bold)


def _set_shape_font_color(
    slide: Any,
    *,
    shape_index: int,
    color_rgb: RGBColor,
) -> None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for paragraph in list(tf.paragraphs):
        runs = list(paragraph.runs)
        if not runs:
            run = paragraph.add_run()
            runs = [run]
        for run in runs:
            try:
                run.font.color.rgb = color_rgb
            except Exception:
                continue


def _set_shape_font_name(
    slide: Any,
    *,
    shape_index: int,
    font_name: str,
) -> None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for paragraph in list(tf.paragraphs):
        runs = list(paragraph.runs)
        if not runs:
            run = paragraph.add_run()
            runs = [run]
        for run in runs:
            try:
                run.font.name = str(font_name or "").strip() or _PPT_FONT_BODY
            except Exception:
                continue


def _remove_shape_indices(slide: Any, *shape_indices: int) -> None:
    shapes: list[Any] = []
    seen: set[int] = set()
    for shape_index in shape_indices:
        shape = _shape_or_none(slide, int(shape_index))
        if shape is None:
            continue
        marker = id(shape)
        if marker in seen:
            continue
        seen.add(marker)
        shapes.append(shape)
    for shape in shapes:
        _remove_shape(shape)


def _remove_slide_number_artifacts(prs: Any) -> None:
    """Drop inherited PowerPoint page-number fields such as ``p. 3``."""

    def _remove_slide_number_nodes(element: Any) -> None:
        try:
            nodes = list(element.iter())
        except Exception:
            nodes = []
        for node in nodes:
            tag = str(getattr(node, "tag", "") or "")
            if not tag.endswith("}sp"):
                continue
            try:
                xml = str(node.xml)
            except Exception:
                xml = ""
            if 'type="slidenum"' not in xml and "p. </a:t>" not in xml:
                continue
            try:
                node.getparent().remove(node)
            except Exception:
                continue

    collections: list[Any] = []
    try:
        collections.extend(list(prs.slides))
    except Exception:
        pass
    try:
        collections.extend(list(prs.slide_layouts))
    except Exception:
        pass
    try:
        collections.extend(list(prs.slide_masters))
    except Exception:
        pass

    for owner in collections:
        for shape in list(getattr(owner, "shapes", [])):
            try:
                xml = str(shape.element.xml)
            except Exception:
                xml = ""
            if 'type="slidenum"' in xml or ">p. <" in xml or "p. </a:t>" in xml:
                _remove_shape(shape)
        _remove_slide_number_nodes(getattr(owner, "element", None))

    try:
        parts = list(prs.part.package.iter_parts())
    except Exception:
        parts = []
    for part in parts:
        partname = str(getattr(part, "partname", "") or "")
        if "/slideLayouts/" not in partname and "/slideMasters/" not in partname:
            continue
        _remove_slide_number_nodes(getattr(part, "_element", None))


def validate_shapes_inside_slide(prs: Any) -> None:
    slide_width = _safe_emu(getattr(prs, "slide_width", None), default=9_144_000)
    slide_height = _safe_emu(getattr(prs, "slide_height", None), default=5_143_500)
    offenders = list(
        iter_out_of_viewport_shapes(
            getattr(prs, "slides", []),
            slide_width=slide_width,
            slide_height=slide_height,
        )
    )
    if not offenders:
        return
    details: list[str] = []
    for shape in offenders[:8]:
        details.append(
            "shape"
            f"(left={int(getattr(shape, 'left', 0) or 0)}, "
            f"top={int(getattr(shape, 'top', 0) or 0)}, "
            f"width={int(getattr(shape, 'width', 0) or 0)}, "
            f"height={int(getattr(shape, 'height', 0) or 0)})"
        )
    raise ValueError("El informe contiene shapes fuera del canvas: " + "; ".join(details))


def _to_roman(value: int) -> str:
    num = max(int(value or 0), 0)
    if num <= 0:
        return ""
    pairs = (
        (1000, "M"),
        (900, "CM"),
        (500, "D"),
        (400, "CD"),
        (100, "C"),
        (90, "XC"),
        (50, "L"),
        (40, "XL"),
        (10, "X"),
        (9, "IX"),
        (5, "V"),
        (4, "IV"),
        (1, "I"),
    )
    out: list[str] = []
    rem = num
    for arabic, roman in pairs:
        while rem >= arabic:
            out.append(roman)
            rem -= arabic
    return "".join(out)


def _chunk_zoom_issues(
    issues: Sequence[FunctionalityIssueRow],
    *,
    rows_per_slide: int,
    notes_by_key: Mapping[str, str] | None = None,
) -> list[tuple[_IssuePageItem, ...]]:
    return _chunk_issue_page_items(
        issues,
        rows_per_slide=rows_per_slide,
        notes_by_key=notes_by_key,
    )


def _issue_note_for_key(
    issue_key: object,
    notes_by_key: Mapping[str, str] | None,
) -> str:
    key = str(issue_key or "").strip().upper()
    if not key or not notes_by_key:
        return ""
    return str(notes_by_key.get(key) or "").strip()


def _split_comment_block(block: str, *, max_chars: int) -> list[str]:
    text = str(block or "").strip()
    if not text:
        return []
    size = max(int(max_chars or 0), 80)
    chunks: list[str] = []
    remaining = text
    while len(remaining) > size:
        candidate = remaining[:size].rstrip()
        min_break = max(int(size * 0.55), 1)
        break_positions = [
            candidate.rfind(separator) for separator in ("\n", ". ", "; ", ", ", " ", "/", "-")
        ]
        best_break = max(break_positions or [-1])
        if best_break < min_break:
            best_break = size
        chunk = remaining[:best_break].strip()
        if chunk:
            chunks.append(chunk)
        remaining = remaining[best_break:].strip()
    if remaining:
        chunks.append(remaining)
    return chunks


def _issue_comment_chunks(comment: object) -> tuple[str, ...]:
    clean = str(comment or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not clean:
        return ()
    clean = re.sub(r"\n{3,}", "\n\n", clean)
    chunks: list[str] = []
    for block in re.split(r"\n\s*\n", clean):
        chunks.extend(_split_comment_block(block, max_chars=_ISSUE_COMMENT_CHUNK_CHARS))
    return tuple(chunk for chunk in chunks if str(chunk or "").strip())


def _issue_page_item(issue: Any, *, notes_by_key: Mapping[str, str] | None) -> _IssuePageItem:
    return _IssuePageItem(
        issue=issue,
        comment_chunks=_issue_comment_chunks(
            _issue_note_for_key(getattr(issue, "key", ""), notes_by_key)
        ),
    )


def _coerce_issue_page_items(
    issues: Sequence[Any],
    *,
    notes_by_key: Mapping[str, str] | None,
) -> list[_IssuePageItem]:
    items: list[_IssuePageItem] = []
    for item in list(issues or []):
        if isinstance(item, _IssuePageItem):
            items.append(item)
        else:
            items.append(_issue_page_item(item, notes_by_key=notes_by_key))
    return items


def _chunk_issue_page_items(
    issues: Sequence[Any],
    *,
    rows_per_slide: int,
    notes_by_key: Mapping[str, str] | None = None,
) -> list[tuple[_IssuePageItem, ...]]:
    size = max(float(rows_per_slide or 0), 1.0)
    items = list(issues or [])
    if not items:
        return [tuple()]

    pages: list[tuple[_IssuePageItem, ...]] = []
    current: list[_IssuePageItem] = []
    current_units = 0.0

    def flush() -> None:
        nonlocal current, current_units
        if current:
            pages.append(tuple(current))
            current = []
            current_units = 0.0

    for issue in items:
        base_item = _issue_page_item(issue, notes_by_key=notes_by_key)
        if not base_item.comment_chunks:
            if current and current_units + 1.0 > size:
                flush()
            current.append(base_item)
            current_units += 1.0
            continue

        remaining = list(base_item.comment_chunks)
        while remaining:
            minimum_units = 1.0 + _ISSUE_TABLE_COMMENT_ROW_UNITS
            if current and current_units + minimum_units > size:
                flush()
                continue

            space_units = max(size - current_units, minimum_units)
            max_comment_chunks = int((space_units - 1.0) // _ISSUE_TABLE_COMMENT_ROW_UNITS)
            max_comment_chunks = max(max_comment_chunks, 1)
            take = min(len(remaining), max_comment_chunks)
            page_item = _IssuePageItem(issue=issue, comment_chunks=tuple(remaining[:take]))
            current.append(page_item)
            current_units += 1.0 + (_ISSUE_TABLE_COMMENT_ROW_UNITS * float(take))
            remaining = remaining[take:]
            if remaining:
                flush()

    flush()
    return pages or [tuple()]


def _issue_comment_row(comment: str) -> list[str]:
    return [
        "",
        "Comentarios registrados:\n" + str(comment or "").strip(),
        "",
        "",
        "",
        "",
    ]


def _shape_text_frame(
    slide: Any,
    *,
    shape_index: int,
) -> Any | None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.word_wrap = True
    except Exception:
        pass
    return tf


def _set_paragraph_single_run(
    paragraph: Any,
    *,
    text: str,
    size_pt: float,
    bold: bool = True,
    italic: bool = False,
    space_before_pt: float = 0.0,
    color_rgb: RGBColor | None = None,
    font_name: str | None = None,
) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.size = Pt(float(size_pt))
    try:
        run.font.name = str(font_name or _PPT_FONT_BODY_MEDIUM)
    except Exception:
        pass
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if color_rgb is not None:
        run.font.color.rgb = color_rgb
    paragraph.space_before = Pt(float(space_before_pt))
    paragraph.space_after = Pt(0)


def _set_paragraph_value_label(
    paragraph: Any,
    *,
    value_text: str,
    label_text: str,
    value_size_pt: float,
    label_size_pt: float,
    color_rgb: RGBColor | None = None,
    value_font_name: str | None = None,
    label_font_name: str | None = None,
) -> None:
    paragraph.clear()
    value_run = paragraph.add_run()
    value_run.text = f"{str(value_text)} "
    value_run.font.size = Pt(float(value_size_pt))
    try:
        value_run.font.name = str(value_font_name or _PPT_FONT_BODY_MEDIUM)
    except Exception:
        pass
    value_run.font.bold = True
    if color_rgb is not None:
        value_run.font.color.rgb = color_rgb
    label_run = paragraph.add_run()
    label_run.text = str(label_text or "")
    label_run.font.size = Pt(float(label_size_pt))
    try:
        label_run.font.name = str(label_font_name or _PPT_FONT_BODY_MEDIUM)
    except Exception:
        pass
    label_run.font.bold = False
    if color_rgb is not None:
        label_run.font.color.rgb = color_rgb
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def _set_paragraph_kpi_side_metric(
    paragraph: Any,
    *,
    metric: KpiSideMetric,
    size_pt: float,
    color_rgb: RGBColor | None = None,
    space_before_pt: float = 0.0,
) -> None:
    paragraph.clear()
    value_run = paragraph.add_run()
    value_run.text = f"{str(metric.value_text or '').strip()} "
    value_run.font.size = Pt(float(size_pt))
    try:
        value_run.font.name = _PPT_FONT_BODY_MEDIUM
    except Exception:
        pass
    value_run.font.bold = True
    if color_rgb is not None:
        value_run.font.color.rgb = color_rgb

    label_run = paragraph.add_run()
    label_run.text = str(metric.label_text or "").strip()
    label_run.font.size = Pt(float(size_pt))
    try:
        label_run.font.name = _PPT_FONT_BODY_MEDIUM
    except Exception:
        pass
    label_run.font.bold = False
    if color_rgb is not None:
        label_run.font.color.rgb = color_rgb
    paragraph.space_before = Pt(float(space_before_pt))
    paragraph.space_after = Pt(0)


def _first_run_color_rgb(shape: Any) -> RGBColor | None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in list(shape.text_frame.paragraphs):
        for run in list(paragraph.runs):
            color = getattr(getattr(run, "font", None), "color", None)
            rgb = getattr(color, "rgb", None)
            if rgb is not None:
                return cast(RGBColor, rgb)
    return None


def _set_shape_fill_and_line(slide: Any, *, shape_index: int, fill_hex: str, line_hex: str) -> None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None:
        return
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_hex))
    except Exception:
        pass
    try:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(line_hex))
        shape.line.width = Pt(1.0)
    except Exception:
        pass


def _style_summary_open_criticity_cards(slide: Any) -> None:
    _set_shape_fill_and_line(
        slide,
        shape_index=5,
        fill_hex=BBVA_REPORT_RED_BG,
        line_hex=BBVA_REPORT_RED_BORDER,
    )
    _set_shape_fill_and_line(
        slide,
        shape_index=6,
        fill_hex=BBVA_REPORT_AMBER_BG,
        line_hex=BBVA_REPORT_AMBER_BORDER,
    )


def _write_open_criticity_card(
    slide: Any,
    *,
    shape_index: int,
    value: int,
    label: str,
    color_rgb: RGBColor | None = None,
) -> None:
    shape = _shape_or_none(slide, shape_index)
    base_color = color_rgb or _first_run_color_rgb(shape)
    tf = _shape_text_frame(slide, shape_index=shape_index)
    if tf is None:
        return
    tf.clear()
    p0 = tf.paragraphs[0]
    _set_paragraph_single_run(
        p0,
        text=str(int(value)),
        size_pt=35.0,
        bold=True,
        color_rgb=base_color,
    )
    p1 = tf.add_paragraph()
    _set_paragraph_single_run(
        p1,
        text=str(label or "").strip(),
        size_pt=9.5,
        bold=True,
        space_before_pt=0.6,
        color_rgb=base_color,
    )


def _write_metric_card(
    slide: Any,
    *,
    shape_index: int,
    value_text: str,
    label_text: str,
    extra_lines: Sequence[tuple[str, float, bool, bool, float]] | None = None,
    value_size_pt: float | None = None,
    label_size_pt: float | None = None,
    text_color_rgb: RGBColor | None = None,
) -> None:
    base_color = text_color_rgb if text_color_rgb is not None else RGBColor(0, 0, 0)
    tf = _shape_text_frame(slide, shape_index=shape_index)
    if tf is None:
        return
    row = KpiRow(value_text=str(value_text or ""), label_text=str(label_text or ""))
    typography = metric_card_typography(row.value_text, row.label_text)
    resolved_value_size = (
        float(value_size_pt) if value_size_pt is not None else typography.value_size_pt
    )
    resolved_label_size = (
        float(label_size_pt) if label_size_pt is not None else typography.label_size_pt
    )
    apply_text_frame_margins(tf, margin_pt=0.0)
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    tf.clear()
    p0 = tf.paragraphs[0]
    _set_paragraph_value_label(
        p0,
        value_text=row.value_text,
        label_text=row.label_text,
        value_size_pt=resolved_value_size,
        label_size_pt=resolved_label_size,
        color_rgb=base_color,
    )

    for text, size_pt, bold, italic, space_before in list(extra_lines or []):
        p = tf.add_paragraph()
        _set_paragraph_single_run(
            p,
            text=str(text or ""),
            size_pt=float(size_pt),
            bold=bool(bold),
            italic=bool(italic),
            space_before_pt=float(space_before),
            color_rgb=base_color,
        )


def _add_metric_split_column(
    slide: Any,
    *,
    card_shape_index: int,
    top_label: str,
    top_value: int,
    bottom_label: str,
    bottom_value: int,
    value_unit: str = "",
    text_color_rgb: RGBColor | None = None,
    detail_size_pt: float | None = None,
) -> None:
    card = _shape_or_none(slide, card_shape_index)
    if card is None:
        return

    base_color = (
        text_color_rgb
        if text_color_rgb is not None
        else (_first_run_color_rgb(_shape_or_none(slide, 16)) or RGBColor(4, 19, 139))
    )

    left = int(card.left)
    top = int(card.top)
    width = int(card.width)
    height = int(card.height)
    if width <= 0 or height <= 0:
        return

    theme = PERIOD_FOLLOWUP_LAYOUT
    divider_left = left + int(width * theme.split_column_ratio)
    divider_top = top + int(height * theme.split_divider_top_ratio)
    divider_height = int(height * theme.split_divider_height_ratio)
    divider_width = max(int(width * theme.split_divider_width_ratio), 1)

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        divider_left,
        divider_top,
        divider_width,
        divider_height,
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = base_color
    divider.line.fill.background()

    text_left = divider_left + int(width * theme.split_column_padding_ratio)
    text_top = top + int(height * theme.split_text_top_ratio)
    text_width = max(
        (left + width) - text_left - int(width * theme.split_column_right_padding_ratio), 1
    )
    text_height = int(height * theme.split_text_height_ratio)
    split_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = split_box.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass
    tf.clear()
    typography = metric_card_typography(top_value, f"{top_label} {bottom_label}")
    resolved_detail_size = (
        float(detail_size_pt) if detail_size_pt is not None else typography.detail_size_pt
    )
    unit = str(value_unit or "").strip()
    top_metric = KpiSideMetric(
        value_text=str(int(top_value)),
        label_text=f"{unit} {str(top_label).strip()}".strip(),
    )
    bottom_metric = KpiSideMetric(
        value_text=str(int(bottom_value)),
        label_text=f"{unit} {str(bottom_label).strip()}".strip(),
    )

    p0 = tf.paragraphs[0]
    _set_paragraph_kpi_side_metric(
        p0,
        metric=top_metric,
        size_pt=resolved_detail_size,
        color_rgb=base_color,
    )
    p1 = tf.add_paragraph()
    _set_paragraph_kpi_side_metric(
        p1,
        metric=bottom_metric,
        size_pt=resolved_detail_size,
        color_rgb=base_color,
        space_before_pt=theme.split_metric_gap_pt,
    )


def _write_created_total_column(
    slide: Any,
    *,
    card_shape_index: int,
    detail_shape_index: int,
    divider_shape_index: int,
    previous_range_label: str,
    previous_value: int,
    total_value: int,
    text_color_rgb: RGBColor | None = None,
) -> None:
    card = _shape_or_none(slide, card_shape_index)
    detail = _shape_or_none(slide, detail_shape_index)
    divider = _shape_or_none(slide, divider_shape_index)
    if card is None or detail is None or not getattr(detail, "has_text_frame", False):
        return

    base_color = (
        text_color_rgb
        if text_color_rgb is not None
        else (_first_run_color_rgb(detail) or RGBColor(4, 19, 139))
    )

    left = int(card.left)
    top = int(card.top)
    width = int(card.width)
    height = int(card.height)
    theme = PERIOD_FOLLOWUP_LAYOUT
    divider_left = left + int(width * theme.split_column_ratio)
    divider_top = top + int(height * theme.split_divider_top_ratio)
    divider_height = int(height * theme.split_divider_height_ratio)
    divider_width = max(int(width * theme.split_divider_width_ratio), 1)

    if divider is not None:
        divider.left = divider_left
        divider.top = divider_top
        divider.width = divider_width
        divider.height = divider_height
        try:
            divider.line.color.rgb = base_color
        except Exception:
            pass

    text_left = divider_left + int(width * theme.split_column_padding_ratio)
    detail.left = text_left
    detail.top = top + int(height * theme.split_text_top_ratio)
    detail.width = max(
        (left + width) - text_left - int(width * theme.split_column_right_padding_ratio), 1
    )
    detail.height = int(height * theme.split_text_height_ratio)
    tf = detail.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        tf.word_wrap = True
    except Exception:
        pass
    apply_text_frame_margins(tf, margin_pt=0.0)
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass
    tf.clear()
    typography = metric_card_typography(total_value, previous_range_label)
    detail_size = typography.detail_size_pt

    p0 = tf.paragraphs[0]
    _set_paragraph_kpi_side_metric(
        p0,
        metric=KpiSideMetric(
            value_text=str(int(previous_value)),
            label_text=f"del {str(previous_range_label).strip()}",
        ),
        size_pt=detail_size,
        color_rgb=base_color,
    )
    p1 = tf.add_paragraph()
    _set_paragraph_kpi_side_metric(
        p1,
        metric=KpiSideMetric(value_text=str(int(total_value)), label_text="en TOTAL"),
        size_pt=detail_size,
        color_rgb=base_color,
        space_before_pt=theme.split_metric_gap_pt,
    )


def _remove_summary_legacy_artifacts(slide: Any) -> None:
    # Legacy template markers, detail links, and stray arrows no longer render.
    _remove_shape_indices(slide, 18, 14, 11, 8, 7)


def _configure_summary_delta_badge(
    slide: Any,
    *,
    shape_index: int,
    card_shape_index: int | None = None,
    font_size_pt: float = _SUMMARY_DELTA_FONT_SIZE_PT,
) -> None:
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return

    if card_shape_index is not None:
        card = _shape_or_none(slide, int(card_shape_index))
        if card is not None:
            try:
                theme = PERIOD_FOLLOWUP_LAYOUT
                divider_left = int(card.left) + int(card.width * theme.split_column_ratio)
                badge_width = max(int(card.width * theme.delta_badge_width_ratio), 1)
                badge_height = max(int(card.height * theme.delta_badge_height_ratio), 1)
                badge_right = divider_left - int(card.width * theme.delta_badge_right_gap_ratio)
                shape.width = badge_width
                shape.height = badge_height
                min_left = int(card.left) + int(card.width * theme.delta_badge_min_left_ratio)
                shape.left = max(badge_right - badge_width, min_left)
                shape.top = int(card.top) + int(card.height * theme.delta_badge_top_ratio)
            except Exception:
                pass

    tf = shape.text_frame
    try:
        tf.word_wrap = False
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass
    try:
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass

    for paragraph in list(tf.paragraphs):
        try:
            paragraph.alignment = PP_ALIGN.CENTER
        except Exception:
            pass
        for run in list(paragraph.runs):
            txt = str(getattr(run, "text", "") or "")
            if not txt.strip():
                continue
            try:
                run.font.bold = True
            except Exception:
                pass
            try:
                run.font.size = Pt(float(font_size_pt))
            except Exception:
                pass


def _overlay_picture(
    slide: Any,
    *,
    payload: bytes,
    anchor_shape: Any | None = None,
    anchor_shape_index: int | None = None,
    replace_anchor: bool = False,
    preserve_aspect: bool = False,
) -> Any | None:
    anchor = anchor_shape
    if anchor is None and anchor_shape_index is not None:
        anchor = _shape_or_none(slide, anchor_shape_index)
    if anchor is None:
        return None
    if not payload:
        return None
    if preserve_aspect:
        rendered = _overlay_picture_contain(
            slide,
            payload=payload,
            frame_left=int(anchor.left),
            frame_top=int(anchor.top),
            frame_width=int(anchor.width),
            frame_height=int(anchor.height),
        )
        if replace_anchor:
            _remove_shape(anchor)
        return rendered
    rendered = slide.shapes.add_picture(
        BytesIO(payload),
        anchor.left,
        anchor.top,
        width=anchor.width,
        height=anchor.height,
    )
    if replace_anchor:
        _remove_shape(anchor)
    return rendered


def _resolve_summary_chart_anchor(slide: Any) -> Any | None:
    # Prefer the main chart placeholder (largest picture in summary slide).
    picture_shapes = _picture_candidates(slide, min_area_in2=1.0)
    if picture_shapes:
        return max(picture_shapes, key=_shape_area_in2)
    # Backward compatibility with canonical corporate template.
    return _shape_or_none(slide, 20)


def _summary_chart_export_size(anchor: Any | None) -> tuple[int, int]:
    export_width = 1200
    if anchor is None:
        return export_width, int(round(export_width / 1.45))
    width = max(int(getattr(anchor, "width", 0) or 0), 1)
    height = max(int(getattr(anchor, "height", 0) or 0), 1)
    ratio = max(float(width) / float(height), 0.8)
    return export_width, max(int(round(float(export_width) / ratio)), 520)


def _chart_png(
    settings: Settings,
    *,
    dff: pd.DataFrame,
    open_df: pd.DataFrame,
    chart_id: str,
    width: int | None = None,
    height: int | None = None,
    slide_optimized: bool = False,
) -> bytes:
    registry = build_trends_registry()
    spec = registry.get(chart_id)
    if spec is None:
        return b""

    kpis = compute_kpis(dff, settings=settings, include_timeseries_chart=(chart_id == "timeseries"))
    fig = spec.render(ChartContext(dff=dff, open_df=open_df, kpis=kpis))
    if fig is None:
        return b""
    if chart_id == "timeseries":
        export_width = int(width or EXEC_CHART_EXPORT_WIDTH)
        export_height = int(height or EXEC_CHART_TREND_EXPORT_HEIGHT)
        for trace in list(getattr(fig, "data", ())):
            trace_type = str(getattr(trace, "type", "") or "").lower()
            if trace_type in {"scatter", "scattergl"}:
                try:
                    trace.mode = "lines+markers"
                except Exception:
                    pass
                try:
                    base_width = 6.4 if slide_optimized else 4.2
                    token = str(getattr(trace, "name", "") or "").strip().lower()
                    if "backlog" in token or "abierto" in token:
                        base_width = 7.0 if slide_optimized else 4.8
                    trace.line.width = base_width
                except Exception:
                    pass
                try:
                    trace.line.dash = "solid"
                except Exception:
                    pass
                try:
                    trace.marker.size = 11.0 if slide_optimized else 8.0
                except Exception:
                    pass
                try:
                    trace.marker.symbol = "circle"
                except Exception:
                    pass
                try:
                    trace.connectgaps = True
                except Exception:
                    pass
                try:
                    trace.opacity = 0.98
                except Exception:
                    pass
        _apply_executive_timeseries_chart_layout(
            fig,
            width=export_width,
            height=export_height,
            font_scale=1.42 if slide_optimized else 1.0,
            x_nticks=5 if slide_optimized else 8,
            tickangle=-34 if slide_optimized else -24,
        )
        payload = _fig_to_png_exact(
            fig,
            width=export_width,
            height=export_height,
            scale=1.0,
        )
        return payload or b""
    payload = _fig_to_png_exact(fig, width=3400, height=760)
    return payload or b""


def _normalize_lookup_token(value: object) -> str:
    txt = str(value or "").strip().lower()
    txt = unicodedata.normalize("NFKD", txt)
    txt = "".join(ch for ch in txt if not unicodedata.combining(ch))
    txt = re.sub(r"[^a-z0-9]+", " ", txt)
    return re.sub(r"\s+", " ", txt).strip()


def _boardroom_snippet(text: str, *, max_chars: int) -> str:
    clean = re.sub(r"\s+", " ", str(text or "").strip())
    if not clean:
        return ""
    first = clean.split(". ")[0].strip()
    if first and len(first) <= max_chars:
        return first if first.endswith(".") else f"{first}."
    return _trim_text(clean, max_chars=max_chars)


def _single_line_ellipsis(text: object, *, max_chars: int) -> str:
    clean = re.sub(r"\s+", " ", str(text or "").strip())
    return ellipsize_text(clean, max_chars=max_chars)


def _priority_order_key(value: object) -> tuple[int, str]:
    label = str(value or "").strip()
    return (int(priority_rank(label)), label.lower())


def _format_quincena_axis_ym(start_value: object, end_value: object) -> str:
    start = pd.to_datetime(start_value, errors="coerce")
    end = pd.to_datetime(end_value, errors="coerce")
    if pd.isna(start) or pd.isna(end):
        return str(start_value or "")
    return f"{int(start.month):02d} |\n{int(start.day)}-{int(end.day)}"


def _clear_slide_shapes(slide: Any) -> None:
    for shape in list(getattr(slide, "shapes", [])):
        _remove_shape(shape)


def _fig_to_png_exact(
    fig: Optional[go.Figure], *, width: int, height: int, scale: float = 2.0
) -> bytes:
    if fig is None:
        return b""
    safe_scale = max(float(scale or 0.0), 0.5)
    try:
        return _kaleido_png_bytes(
            fig_obj=fig,
            scale=safe_scale,
            export_width=max(int(width), 640),
            export_height=max(int(height), 360),
        )
    except Exception:
        payload = _fig_to_png(fig)
        return payload or b""


def _exec_chart_margin(**overrides: int) -> dict[str, int]:
    margin = dict(EXEC_CHART_MARGIN)
    margin.update({str(key): int(value) for key, value in overrides.items()})
    return margin


def _chart_total_offset(max_total: float) -> float:
    safe_max = max(float(max_total or 0.0), 0.0)
    return max(safe_max * 0.075, 1.0)


def _inside_label_min_value(max_total: float) -> float:
    safe_max = max(float(max_total or 0.0), 0.0)
    if safe_max <= 4.0:
        return 1.0
    return max(safe_max * 0.07, 2.0)


def _safe_inside_bar_text(values: Sequence[float], *, min_value: float) -> list[str]:
    labels: list[str] = []
    threshold = max(float(min_value or 0.0), 0.0)
    for value in list(values or []):
        try:
            numeric = float(value)
        except Exception:
            numeric = 0.0
        labels.append(str(int(round(numeric))) if numeric >= threshold else "")
    return labels


def _add_stacked_bar_totals(
    fig: go.Figure,
    *,
    axis_labels: Sequence[str],
    totals: Sequence[float],
    max_total: float,
    color: str,
) -> None:
    labels = [str(label) for label in list(axis_labels or [])]
    values = [float(value or 0.0) for value in list(totals or [])]
    if not labels or not values:
        return
    offset = _chart_total_offset(max_total)
    fig.add_trace(
        go.Scatter(
            x=labels,
            y=[value + offset for value in values],
            mode="text",
            text=[str(int(round(value))) if value > 0 else "" for value in values],
            textposition="top center",
            textfont=dict(size=EXEC_CHART_TOTAL_FONT_PT, color=color),
            hoverinfo="skip",
            showlegend=False,
            cliponaxis=False,
        )
    )
    fig.update_yaxes(range=[0, max(float(max_total) + (offset * 2.9), 1.0)])


def _apply_executive_chart_layout(
    fig: go.Figure,
    *,
    kind: str,
    show_legend: bool = True,
    x_title: str | None = None,
    y_title: str | None = None,
    width: int | None = None,
    height: int | None = None,
    margin: Mapping[str, int] | None = None,
    font_scale: float = 1.0,
    x_nticks: int | None = None,
    tickangle: int | None = None,
) -> None:
    kind_token = str(kind or "").strip().lower()
    export_width = int(width or EXEC_CHART_EXPORT_WIDTH)
    export_height = int(height or EXEC_CHART_EXPORT_HEIGHT)
    legend_y = -0.36 if kind_token in {"trend", "timeseries"} else -0.25
    x_tick_angle = (
        int(tickangle)
        if tickangle is not None
        else (-24 if kind_token in {"trend", "timeseries"} else 0)
    )
    safe_font_scale = max(float(font_scale or 1.0), 0.6)
    axis_font_pt = int(round(EXEC_CHART_AXIS_FONT_PT * safe_font_scale))
    axis_title_font_pt = int(round(EXEC_CHART_AXIS_TITLE_FONT_PT * safe_font_scale))
    legend_font_pt = int(round(EXEC_CHART_LEGEND_FONT_PT * safe_font_scale))
    fig.update_layout(
        width=export_width,
        height=export_height,
        xaxis_title=str(x_title or ""),
        yaxis_title=str(y_title or ""),
        showlegend=bool(show_legend),
        hovermode="x",
        uniformtext=dict(minsize=max(EXEC_CHART_INSIDE_VALUE_FONT_PT - 6, 18), mode="hide"),
        margin=dict(margin or EXEC_CHART_MARGIN),
        plot_bgcolor="#F6F8FC",
        paper_bgcolor="#F6F8FC",
    )
    fig.update_xaxes(
        tickangle=x_tick_angle,
        tickfont=dict(size=axis_font_pt, color="#1E2C46"),
        title_font=dict(size=axis_title_font_pt, color="#17253F"),
        automargin=True,
        nticks=(
            x_nticks
            if x_nticks is not None
            else (8 if kind_token in {"trend", "timeseries"} else None)
        ),
        gridcolor="rgba(155, 169, 196, 0.20)",
        zeroline=False,
    )
    fig.update_yaxes(
        tickfont=dict(size=axis_font_pt, color="#1E2C46"),
        title_font=dict(size=axis_title_font_pt, color="#17253F"),
        automargin=True,
        nticks=6,
        gridcolor="rgba(155, 169, 196, 0.24)",
        zeroline=False,
    )
    if show_legend:
        fig.update_layout(
            legend=dict(
                title=dict(text=""),
                orientation="h",
                xanchor="center",
                x=0.5,
                yanchor="top",
                y=legend_y,
                font=dict(size=legend_font_pt, color="#1A2740"),
                bgcolor="rgba(255,255,255,0.96)",
                bordercolor="rgba(188,198,216,0.95)",
                borderwidth=1,
                traceorder="normal",
            )
        )


def _axis_title_text(fig: go.Figure, axis_name: str) -> str:
    axis = getattr(getattr(fig, "layout", None), axis_name, None)
    title = getattr(axis, "title", None)
    return str(getattr(title, "text", "") or "")


def _apply_executive_timeseries_chart_layout(
    fig: go.Figure,
    *,
    width: int = EXEC_CHART_EXPORT_WIDTH,
    height: int = EXEC_CHART_TREND_EXPORT_HEIGHT,
    font_scale: float = 1.0,
    x_nticks: int = 8,
    tickangle: int = -24,
) -> None:
    axis_font_pt = int(round(EXEC_CHART_AXIS_FONT_PT * max(float(font_scale or 1.0), 0.6)))
    axis_title_font_pt = int(
        round(EXEC_CHART_AXIS_TITLE_FONT_PT * max(float(font_scale or 1.0), 0.6))
    )
    margin_scale = max(float(font_scale or 1.0), 1.0)
    _apply_executive_chart_layout(
        fig,
        kind="timeseries",
        show_legend=True,
        x_title=_axis_title_text(fig, "xaxis"),
        y_title=_axis_title_text(fig, "yaxis"),
        width=int(width),
        height=int(height),
        margin=_exec_chart_margin(
            l=int(round(82 * margin_scale)),
            r=int(round(56 * margin_scale)),
            t=int(round(56 * margin_scale)),
            b=int(round(190 * margin_scale)),
        ),
        font_scale=font_scale,
        x_nticks=int(x_nticks),
        tickangle=int(tickangle),
    )
    fig.update_layout(
        plot_bgcolor="#FFFFFF",
        paper_bgcolor="#FFFFFF",
        font=dict(size=axis_font_pt, color="#132A7B"),
    )
    fig.update_xaxes(
        showgrid=False,
        showline=True,
        linecolor="#C7D1E6",
        linewidth=1.2,
        tickfont=dict(size=axis_font_pt, color="#213A8F"),
        title_font=dict(size=axis_title_font_pt, color="#17253F"),
        automargin=True,
        nticks=int(x_nticks),
        tickangle=int(tickangle),
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor="#E7EDF8",
        gridwidth=1.0,
        showline=True,
        linecolor="#C7D1E6",
        linewidth=1.2,
        tickfont=dict(size=axis_font_pt, color="#213A8F"),
        title_font=dict(size=axis_title_font_pt, color="#17253F"),
        automargin=True,
        nticks=6,
    )


def _overlay_picture_contain(
    slide: Any,
    *,
    payload: bytes,
    frame_left: int,
    frame_top: int,
    frame_width: int,
    frame_height: int,
) -> Any | None:
    if not payload:
        return None
    try:
        img = Image.open(BytesIO(payload))
        src_w = float(max(int(getattr(img, "width", 1) or 1), 1))
        src_h = float(max(int(getattr(img, "height", 1) or 1), 1))
    except Exception:
        src_w, src_h = 16.0, 9.0

    frame_w = float(max(int(frame_width or 1), 1))
    frame_h = float(max(int(frame_height or 1), 1))
    src_ratio = src_w / src_h
    frame_ratio = frame_w / frame_h

    if frame_ratio >= src_ratio:
        pic_h = frame_h
        pic_w = pic_h * src_ratio
    else:
        pic_w = frame_w
        pic_h = pic_w / src_ratio

    left = int(round(float(frame_left) + (frame_w - pic_w) / 2.0))
    top = int(round(float(frame_top) + (frame_h - pic_h) / 2.0))
    return slide.shapes.add_picture(
        BytesIO(payload),
        left,
        top,
        width=int(round(pic_w)),
        height=int(round(pic_h)),
    )


def _add_exec_textbox(
    slide: Any,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size_pt: float,
    color_rgb: RGBColor,
    font_name: str | None = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.clear()
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(float(font_size_pt))
    try:
        run.font.name = str(font_name).strip() if str(font_name or "").strip() else _PPT_FONT_BODY
    except Exception:
        pass
    run.font.bold = bool(bold)
    run.font.color.rgb = color_rgb
    return box


def _write_exec_metric_block(
    slide: Any,
    *,
    left: int,
    top: int,
    width: int,
    kicker: str,
    value: str,
) -> None:
    _add_exec_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=200_000,
        text=str(kicker or "").upper(),
        font_size_pt=11.4,
        color_rgb=RGBColor(*_EXEC_TEXT_SECONDARY_RGB),
        bold=True,
    )
    _add_exec_textbox(
        slide,
        left=left,
        top=top + 150_000,
        width=width,
        height=290_000,
        text=str(value or "—"),
        font_size_pt=32.0,
        color_rgb=RGBColor(*_EXEC_TEXT_PRIMARY_RGB),
        bold=True,
    )


def _extract_resolution_story_values(
    dff: pd.DataFrame, open_df: pd.DataFrame
) -> tuple[str, str, str]:
    pack = build_trend_insight_pack("resolution_hist", dff=dff, open_df=open_df)
    metric_by_label = {
        _normalize_lookup_token(getattr(metric, "label", "")): str(getattr(metric, "value", "—"))
        for metric in list(getattr(pack, "metrics", []) or [])
    }
    habitual = metric_by_label.get(_normalize_lookup_token("Antigüedad habitual"), "—")
    stalled = metric_by_label.get(_normalize_lookup_token("Casos más atascados"), "—")
    over_30 = metric_by_label.get(_normalize_lookup_token(">30d abiertas"), "—")
    return habitual, stalled, over_30


def _resolution_cards_by_title(dff: pd.DataFrame, open_df: pd.DataFrame) -> Mapping[str, str]:
    pack = build_trend_insight_pack("resolution_hist", dff=dff, open_df=open_df)
    cards = list(getattr(pack, "cards", []) or [])
    cards_by_token: dict[str, str] = {
        _normalize_lookup_token(getattr(card, "title", "")): str(getattr(card, "body", "")).strip()
        for card in cards
    }
    ordered_titles = (
        "Incidencias críticas envejecidas",
        "Brecha por prioridad",
        "Riesgo real de envejecimiento",
        "Cola extrema de antigüedad",
    )
    resolved: dict[str, str] = {}
    for title in ordered_titles:
        token = _normalize_lookup_token(title)
        body = str(cards_by_token.get(token, "") or "").strip()
        if body:
            resolved[title] = body
            continue
        fallback = next(
            (
                str(getattr(card, "body", "") or "").strip()
                for card in cards
                if str(getattr(card, "body", "") or "").strip()
                and str(getattr(card, "title", "") or "").strip() not in resolved
            ),
            "",
        )
        resolved[title] = (
            fallback or "Sin datos suficientes para este insight en el scope seleccionado."
        )
    return resolved


def _resolution_chart_png_executive(
    settings: Settings,
    *,
    dff: pd.DataFrame,
    open_df: pd.DataFrame,
    reference_now: pd.Timestamp | None = None,
) -> bytes:
    _ = (settings, open_df)
    age_payload = build_open_age_priority_payload(dff, reference_now=reference_now)
    grouped = age_payload.get("grouped") if isinstance(age_payload, dict) else None
    if not isinstance(grouped, pd.DataFrame) or grouped.empty:
        return b""

    work = grouped.copy(deep=False)
    work["age_bucket"] = work["age_bucket"].astype(str)
    work["priority"] = work["priority"].fillna("(sin priority)").astype(str)
    work["count"] = pd.to_numeric(work["count"], errors="coerce").fillna(0).astype(int)
    axis_labels = [str(label) for label in OPEN_AGE_BUCKET_LABELS]
    positive = work.loc[work["count"].gt(0)]
    priorities = sorted(
        positive["priority"].dropna().astype(str).unique().tolist(), key=_priority_order_key
    )
    if not priorities:
        return b""

    by_bucket_priority = {
        (str(row.age_bucket), str(row.priority)): int(row.count)
        for row in work.itertuples(index=False)
    }
    totals = (
        work.groupby("age_bucket", dropna=False)["count"]
        .sum()
        .reindex(axis_labels)
        .fillna(0)
        .astype(int)
    )
    max_total = float(totals.max()) if not totals.empty else 0.0
    label_min_value = _inside_label_min_value(max_total)
    colors = priority_color_map()
    neutral = colors.get("(sin priority)", "#7E8EA7")

    fig = go.Figure()
    for priority in reversed(priorities):
        values = [int(by_bucket_priority.get((bucket, str(priority)), 0)) for bucket in axis_labels]
        fig.add_trace(
            go.Bar(
                x=axis_labels,
                y=values,
                name=str(priority),
                marker=dict(
                    color=colors.get(str(priority), neutral), line=dict(color="#0A2E72", width=1.0)
                ),
                text=_safe_inside_bar_text(
                    [float(value) for value in values], min_value=label_min_value
                ),
                textposition="inside",
                insidetextanchor="middle",
                textfont=dict(size=EXEC_CHART_INSIDE_VALUE_FONT_PT, color="#FFFFFF"),
                customdata=[[int(totals.get(label, 0))] for label in axis_labels],
                cliponaxis=False,
                hovertemplate=(
                    "Rango: %{x}<br>"
                    "Prioridad: %{fullData.name}<br>"
                    "Incidencias abiertas: %{y}<br>"
                    "Total columna: %{customdata[0]}<extra></extra>"
                ),
            )
        )

    _add_stacked_bar_totals(
        fig,
        axis_labels=axis_labels,
        totals=[float(totals.get(label, 0)) for label in axis_labels],
        max_total=max_total,
        color="#0B3E76",
    )
    _apply_executive_chart_layout(
        fig,
        kind="stacked",
        show_legend=True,
        x_title="Rango en días",
        y_title="Incidencias abiertas",
        height=EXEC_CHART_EXPORT_HEIGHT,
        margin=_exec_chart_margin(t=88, b=154),
    )
    fig.update_layout(barmode="stack", bargap=0.14)
    fig.update_xaxes(type="category", categoryorder="array", categoryarray=axis_labels)
    chart_payload = _fig_to_png_exact(
        fig,
        width=EXEC_CHART_EXPORT_WIDTH,
        height=EXEC_CHART_EXPORT_HEIGHT,
        scale=1.0,
    )
    return chart_payload or b""


def _add_exec_insight_card(
    slide: Any,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    title: str,
    body: str,
    title_font_size_pt: float = 17.0,
    body_font_size_pt: float = 12.8,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(left),
        int(top),
        int(width),
        int(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*_EXEC_CARD_BG_RGB)
    card.line.color.rgb = RGBColor(*_EXEC_ACCENT_BORDER_RGB)
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.clear()
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.margin_left = 82_000
        tf.margin_right = 78_000
        tf.margin_top = 48_000
        tf.margin_bottom = 42_000
    except Exception:
        pass

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.space_before = Pt(0)
    p0.space_after = Pt(0)
    title_run = p0.add_run()
    title_run.text = str(title or "").strip()
    title_run.font.size = Pt(float(title_font_size_pt))
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(*_EXEC_CARD_TITLE_RGB)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(2.2)
    p1.space_after = Pt(0)
    body_run = p1.add_run()
    body_run.text = str(body or "").strip()
    body_run.font.size = Pt(float(body_font_size_pt))
    body_run.font.bold = False
    body_run.font.color.rgb = RGBColor(*_EXEC_TEXT_PRIMARY_RGB)
    try:
        p1.line_spacing = 1.14
    except Exception:
        pass


def _populate_open_aging_executive_slide(
    slide: Any,
    *,
    settings: Settings,
    scope_result: QuincenalScopeResult,
    slide_width: int,
    slide_height: int,
) -> None:
    _clear_slide_shapes(slide)

    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_EXEC_BG_RGB)
    except Exception:
        pass

    slide_w = int(slide_width or 9_144_000)
    slide_h = int(slide_height or 5_143_500)
    margin_x = int(slide_w * 0.044)
    content_w = max(slide_w - (2 * margin_x), 1)

    _add_exec_textbox(
        slide,
        left=margin_x,
        top=int(slide_h * 0.036),
        width=content_w,
        height=int(slide_h * 0.103),
        text="Visión agregada de incidencias abiertas : rango de días por prioridad",
        font_size_pt=22.0,
        color_rgb=RGBColor(*_EXEC_TEXT_PRIMARY_RGB),
        bold=True,
    )

    habitual, stalled, over_30 = _extract_resolution_story_values(
        scope_result.dff,
        scope_result.open_df,
    )

    metric_top = int(slide_h * 0.183)
    metric_gap = int(slide_w * 0.018)
    metric_w = int((content_w - (2 * metric_gap)) / 3)
    _write_exec_metric_block(
        slide,
        left=margin_x,
        top=metric_top,
        width=metric_w,
        kicker="Antigüedad habitual",
        value=habitual,
    )
    _write_exec_metric_block(
        slide,
        left=margin_x + metric_w + metric_gap,
        top=metric_top,
        width=metric_w,
        kicker="Casos más atascados",
        value=stalled,
    )
    _write_exec_metric_block(
        slide,
        left=margin_x + (2 * (metric_w + metric_gap)),
        top=metric_top,
        width=metric_w,
        kicker=">30d abiertas",
        value=over_30,
    )

    chart_frame_left = margin_x
    chart_frame_top = int(slide_h * 0.305)
    chart_frame_width = content_w
    chart_frame_height = int(slide_h * 0.322)
    chart_frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        chart_frame_left,
        chart_frame_top,
        chart_frame_width,
        chart_frame_height,
    )
    chart_frame.fill.solid()
    chart_frame.fill.fore_color.rgb = RGBColor(247, 248, 252)
    chart_frame.line.color.rgb = RGBColor(*_EXEC_ACCENT_BORDER_RGB)
    chart_frame.line.width = Pt(1.2)

    chart_png = _resolution_chart_png_executive(
        settings,
        dff=scope_result.dff,
        open_df=scope_result.open_df,
        reference_now=pd.Timestamp(scope_result.summary.window.current_end),
    )
    if chart_png:
        _overlay_picture_contain(
            slide,
            payload=chart_png,
            frame_left=chart_frame_left + 190_000,
            frame_top=chart_frame_top + 26_000,
            frame_width=chart_frame_width - 380_000,
            frame_height=chart_frame_height - 52_000,
        )
    else:
        _add_exec_textbox(
            slide,
            left=chart_frame_left + 55_000,
            top=chart_frame_top + int(chart_frame_height * 0.42),
            width=chart_frame_width - 110_000,
            height=300_000,
            text="No hay datos suficientes para renderizar el gráfico de antigüedad por prioridad.",
            font_size_pt=17.0,
            color_rgb=RGBColor(*_EXEC_TEXT_PRIMARY_RGB),
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    insight_text = _resolution_cards_by_title(scope_result.dff, scope_result.open_df)
    cards = [
        (
            "Incidencias críticas envejecidas",
            insight_text.get("Incidencias críticas envejecidas", ""),
        ),
        ("Brecha por prioridad", insight_text.get("Brecha por prioridad", "")),
        ("Riesgo real de envejecimiento", insight_text.get("Riesgo real de envejecimiento", "")),
        ("Cola extrema de antigüedad", insight_text.get("Cola extrema de antigüedad", "")),
    ]

    cards_top = int(slide_h * 0.643)
    cards_gap_x = int(slide_w * 0.021)
    cards_gap_y = int(slide_h * 0.017)
    card_w = int((content_w - cards_gap_x) / 2)
    card_h = int((slide_h - cards_top - cards_gap_y - int(slide_h * 0.018)) / 2)
    card_h = max(card_h, int(slide_h * 0.13))
    coords = [
        (margin_x, cards_top),
        (margin_x + card_w + cards_gap_x, cards_top),
        (margin_x, cards_top + card_h + cards_gap_y),
        (margin_x + card_w + cards_gap_x, cards_top + card_h + cards_gap_y),
    ]
    for (left, top), (title, body) in zip(coords, cards):
        _add_exec_insight_card(
            slide,
            left=left,
            top=top,
            width=card_w,
            height=card_h,
            title=title,
            body=_boardroom_snippet(body, max_chars=118),
            title_font_size_pt=12.6,
            body_font_size_pt=9.6,
        )


def _priority_cards_by_title(dff: pd.DataFrame, open_df: pd.DataFrame) -> Mapping[str, str]:
    pack = build_trend_insight_pack("open_priority_pie", dff=dff, open_df=open_df)
    cards = list(getattr(pack, "cards", []) or [])
    cards_by_token: dict[str, str] = {
        _normalize_lookup_token(getattr(card, "title", "")): str(getattr(card, "body", "")).strip()
        for card in cards
    }
    ordered_titles = (
        "Inflación de prioridades altas",
        "Concentración de prioridad",
        "Incidencias de mayor impacto con antigüedad elevada",
        "Incidencias de mayor impacto sin arrancar",
        "Incidencias de mayor impacto sin movimiento reciente",
    )
    resolved: dict[str, str] = {}
    for title in ordered_titles:
        token = _normalize_lookup_token(title)
        body = str(cards_by_token.get(token, "") or "").strip()
        if body:
            resolved[title] = body
            continue
        fallback = next(
            (
                str(getattr(card, "body", "") or "").strip()
                for card in cards
                if str(getattr(card, "body", "") or "").strip()
                and str(getattr(card, "title", "") or "").strip() not in resolved
            ),
            "",
        )
        resolved[title] = (
            fallback or "Sin datos suficientes para este insight en el scope seleccionado."
        )
    return resolved


def _priority_chart_png_executive(
    settings: Settings,
    *,
    dff: pd.DataFrame,
    open_df: pd.DataFrame,
) -> bytes:
    _ = settings
    safe_open = open_df if isinstance(open_df, pd.DataFrame) else pd.DataFrame()
    if safe_open.empty:
        safe_open = dff if isinstance(dff, pd.DataFrame) else pd.DataFrame()
    if safe_open.empty:
        return b""

    if "priority" not in safe_open.columns:
        return b""

    work = safe_open.copy(deep=False)
    work["priority"] = normalize_text_col(work["priority"], "(sin priority)")
    aggregated = (
        work.groupby("priority", dropna=False)
        .size()
        .reset_index(name="count")
        .sort_values("priority", key=lambda col: col.map(_priority_order_key), kind="mergesort")
    )
    aggregated["count"] = pd.to_numeric(aggregated["count"], errors="coerce").fillna(0).astype(int)
    aggregated = aggregated.loc[aggregated["count"].gt(0)].copy(deep=False)
    if aggregated.empty:
        return b""
    labels = [str(value) for value in aggregated["priority"].tolist()]
    values = [int(value) for value in aggregated["count"].tolist()]
    total = max(sum(values), 1)
    aggregated["pct"] = (aggregated["count"].astype(float) / float(total)) * 100.0
    color_map = {
        "supone un impedimento": "#8B0000",
        "highest": "#B51F29",
        "high": "#D64C4C",
        "medium": "#F2A529",
        "low": "#2FA84F",
        "lowest": "#1E8C45",
        "(sin priority)": "#7E8EA7",
    }
    fig = go.Figure()
    max_value = max(values) if values else 1
    label_min_value = _inside_label_min_value(float(max_value))
    for label, value in zip(labels, values):
        color = color_map.get(str(label).strip().lower(), "#4A7BD1")
        fig.add_trace(
            go.Bar(
                x=[label],
                y=[value],
                marker=dict(color=color, line=dict(color="#0A2E72", width=1)),
                text=_safe_inside_bar_text([float(value)], min_value=label_min_value),
                textposition="inside",
                insidetextanchor="middle",
                textfont=dict(size=EXEC_CHART_INSIDE_VALUE_FONT_PT, color="#FFFFFF"),
                cliponaxis=False,
                hovertemplate="Prioridad: %{x}<br>Incidencias: %{y}<extra></extra>",
                name=str(label),
                showlegend=False,
            )
        )
    top_offset = _chart_total_offset(float(max_value))
    percentages = [float(value) for value in aggregated["pct"].tolist()]
    fig.add_trace(
        go.Scatter(
            x=labels,
            y=[float(v) + top_offset for v in values],
            mode="text",
            text=[f"{pct:.1f}%" if val > 0 else "" for pct, val in zip(percentages, values)],
            textposition="top center",
            textfont=dict(size=EXEC_CHART_TOTAL_FONT_PT, color="#0C376E"),
            hoverinfo="skip",
            showlegend=False,
            cliponaxis=False,
        )
    )
    _apply_executive_chart_layout(
        fig,
        kind="priority",
        show_legend=False,
        x_title="",
        y_title="",
        height=EXEC_CHART_EXPORT_HEIGHT,
        margin=_exec_chart_margin(l=58, r=36, t=82, b=92),
    )
    fig.update_layout(
        bargap=0.38,
        uniformtext=dict(minsize=max(EXEC_CHART_INSIDE_VALUE_FONT_PT - 4, 20), mode="hide"),
    )
    fig.update_xaxes(type="category", categoryorder="array", categoryarray=labels)
    fig.update_yaxes(range=[0, max((float(max_value) + top_offset) * 1.28, 1.0)])
    payload = _fig_to_png_exact(
        fig,
        width=EXEC_CHART_EXPORT_WIDTH,
        height=EXEC_CHART_EXPORT_HEIGHT,
        scale=1.0,
    )
    return payload or b""


def _populate_open_priority_executive_slide(
    slide: Any,
    *,
    settings: Settings,
    scope_result: QuincenalScopeResult,
    slide_width: int,
    slide_height: int,
) -> None:
    _clear_slide_shapes(slide)
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_EXEC_BG_RGB)
    except Exception:
        pass

    slide_w = int(slide_width or 9_144_000)
    slide_h = int(slide_height or 5_143_500)
    margin_x = int(slide_w * 0.044)
    content_w = max(slide_w - (2 * margin_x), 1)

    _add_exec_textbox(
        slide,
        left=margin_x,
        top=int(slide_h * 0.031),
        width=content_w,
        height=int(slide_h * 0.068),
        text="Visión agregada de incidencias abiertas por prioridad",
        font_size_pt=24.0,
        color_rgb=RGBColor(*_EXEC_TEXT_PRIMARY_RGB),
        bold=True,
    )

    chart_frame_left = margin_x
    chart_frame_top = int(slide_h * 0.105)
    chart_frame_width = content_w
    chart_frame_height = int(slide_h * 0.43)
    chart_frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        chart_frame_left,
        chart_frame_top,
        chart_frame_width,
        chart_frame_height,
    )
    chart_frame.fill.solid()
    chart_frame.fill.fore_color.rgb = RGBColor(247, 248, 252)
    chart_frame.line.color.rgb = RGBColor(*_EXEC_ACCENT_BORDER_RGB)
    chart_frame.line.width = Pt(1.2)

    chart_png = _priority_chart_png_executive(
        settings,
        dff=scope_result.dff,
        open_df=scope_result.open_df,
    )
    if chart_png:
        _overlay_picture_contain(
            slide,
            payload=chart_png,
            frame_left=chart_frame_left + 190_000,
            frame_top=chart_frame_top + 26_000,
            frame_width=chart_frame_width - 380_000,
            frame_height=chart_frame_height - 52_000,
        )
    else:
        _add_exec_textbox(
            slide,
            left=chart_frame_left + 55_000,
            top=chart_frame_top + int(chart_frame_height * 0.42),
            width=chart_frame_width - 110_000,
            height=300_000,
            text="No hay datos suficientes para renderizar la distribución de prioridad.",
            font_size_pt=17.0,
            color_rgb=RGBColor(*_EXEC_TEXT_PRIMARY_RGB),
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    insight_text = _priority_cards_by_title(scope_result.dff, scope_result.open_df)
    cards = [
        ("Inflación de prioridades altas", insight_text.get("Inflación de prioridades altas", "")),
        ("Concentración de prioridad", insight_text.get("Concentración de prioridad", "")),
        (
            "Incidencias de mayor impacto con antigüedad elevada",
            insight_text.get("Incidencias de mayor impacto con antigüedad elevada", ""),
        ),
        (
            "Incidencias de mayor impacto sin movimiento reciente",
            insight_text.get("Incidencias de mayor impacto sin movimiento reciente", ""),
        ),
        (
            "Incidencias de mayor impacto sin arrancar",
            insight_text.get("Incidencias de mayor impacto sin arrancar", ""),
        ),
    ]

    cards_top = int(slide_h * 0.547)
    cards_gap_x = int(slide_w * 0.021)
    cards_gap_y = int(slide_h * 0.010)
    card_w = int((content_w - cards_gap_x) / 2)
    card_h = int(slide_h * 0.135)
    _add_exec_insight_card(
        slide,
        left=margin_x,
        top=cards_top,
        width=card_w,
        height=card_h,
        title=cards[0][0],
        body=re.sub(r"\s+", " ", str(cards[0][1] or "").strip()),
        title_font_size_pt=12.0,
        body_font_size_pt=8.35,
    )
    _add_exec_insight_card(
        slide,
        left=margin_x + card_w + cards_gap_x,
        top=cards_top,
        width=card_w,
        height=card_h,
        title=cards[2][0],
        body=re.sub(r"\s+", " ", str(cards[2][1] or "").strip()),
        title_font_size_pt=11.8,
        body_font_size_pt=8.2,
    )
    _add_exec_insight_card(
        slide,
        left=margin_x,
        top=cards_top + card_h + cards_gap_y,
        width=card_w,
        height=card_h,
        title=cards[1][0],
        body=re.sub(r"\s+", " ", str(cards[1][1] or "").strip()),
        title_font_size_pt=12.0,
        body_font_size_pt=8.35,
    )
    _add_exec_insight_card(
        slide,
        left=margin_x + card_w + cards_gap_x,
        top=cards_top + card_h + cards_gap_y,
        width=card_w,
        height=card_h,
        title=cards[3][0],
        body=re.sub(r"\s+", " ", str(cards[3][1] or "").strip()),
        title_font_size_pt=11.8,
        body_font_size_pt=8.2,
    )
    full_card_top = cards_top + (2 * (card_h + cards_gap_y))
    full_card_h = max(slide_h - full_card_top - int(slide_h * 0.022), int(slide_h * 0.075))
    _add_exec_insight_card(
        slide,
        left=margin_x,
        top=full_card_top,
        width=content_w,
        height=full_card_h,
        title=cards[4][0],
        body=re.sub(r"\s+", " ", str(cards[4][1] or "").strip()),
        title_font_size_pt=12.0,
        body_font_size_pt=8.35,
    )


def _populate_summary_slide(slide: Any, *, title: str, scope_result: QuincenalScopeResult) -> None:
    summary = scope_result.summary
    summary_metric_color = _first_run_color_rgb(_shape_or_none(slide, 16)) or RGBColor(4, 19, 139)
    window_service = TimeWindowService()
    created_label = window_service.format_current_created_label(
        summary.window,
        singular=int(summary.new_now) == 1,
    )
    closed_label = window_service.format_current_closed_label(
        summary.window,
        singular=int(summary.closed_now) == 1,
    )
    previous_created_label = window_service.format_previous_range_label(summary.window)
    _style_summary_open_criticity_cards(slide)
    _set_shape_text(slide, 3, title)
    delta_badges = {
        10: _summary_delta_badge(summary.closed_delta),
        13: _summary_delta_badge(summary.resolution_delta),
        19: _summary_delta_badge(summary.new_delta),
    }
    for shape_idx, (badge_text, badge_color) in delta_badges.items():
        _set_shape_text(slide, shape_idx, badge_text)
        _set_shape_font_color(slide, shape_index=shape_idx, color_rgb=badge_color)
    focus_side_label = (
        "MAESTRAS"
        if str(summary.open_group_mode or "").strip() == OPEN_ISSUES_FOCUS_MODE_MAESTRAS
        else "CRITICIDADES ALTAS"
    )
    _write_open_criticity_card(
        slide,
        shape_index=4,
        value=int(summary.open_total),
        label="INCIDENCIAS ABIERTAS EN TOTAL",
        color_rgb=RGBColor(255, 255, 255),
    )
    _write_open_criticity_card(
        slide,
        shape_index=5,
        value=int(summary.open_focus_total),
        label=str(summary.open_focus_report_label),
        color_rgb=RGBColor(*hex_to_rgb(BBVA_REPORT_RED_TEXT)),
    )
    _write_open_criticity_card(
        slide,
        shape_index=6,
        value=int(summary.open_other_total),
        label=str(summary.open_other_report_label),
        color_rgb=RGBColor(*hex_to_rgb(BBVA_REPORT_AMBER_TEXT)),
    )

    _write_metric_card(
        slide,
        shape_index=15,
        value_text=str(int(summary.new_now)),
        label_text=created_label,
        text_color_rgb=summary_metric_color,
    )
    _write_created_total_column(
        slide,
        card_shape_index=15,
        detail_shape_index=16,
        divider_shape_index=17,
        previous_range_label=previous_created_label,
        previous_value=int(summary.new_before),
        total_value=int(summary.new_accumulated),
        text_color_rgb=summary_metric_color,
    )
    _write_metric_card(
        slide,
        shape_index=9,
        value_text=str(int(summary.closed_now)),
        label_text=closed_label,
        text_color_rgb=summary_metric_color,
    )
    _add_metric_split_column(
        slide,
        card_shape_index=9,
        top_label=str(focus_side_label),
        top_value=int(summary.closed_focus_now),
        bottom_label="RESTO",
        bottom_value=int(summary.closed_other_now),
    )
    _write_metric_card(
        slide,
        shape_index=12,
        value_text=str(_fmt_days(summary.resolution_days_now)),
        label_text="DÍAS DE RESOLUCIÓN (EN PROMEDIO)",
        extra_lines=(
            [("SIN DATOS", 8.1, True, False, 0.25)] if summary.resolution_days_now is None else None
        ),
        text_color_rgb=summary_metric_color,
    )
    _add_metric_split_column(
        slide,
        card_shape_index=12,
        top_label="MAX",
        top_value=int(_fmt_days(summary.resolution_days_max_now)),
        bottom_label="MIN",
        bottom_value=int(_fmt_days(summary.resolution_days_min_now)),
        value_unit="días",
    )

    _configure_summary_delta_badge(
        slide,
        shape_index=10,
        card_shape_index=9,
        font_size_pt=delta_badge_font_size(
            delta_badges[10][0],
            base_size_pt=_SUMMARY_DELTA_FONT_SIZE_PT,
        ),
    )
    _configure_summary_delta_badge(
        slide,
        shape_index=13,
        card_shape_index=12,
        font_size_pt=delta_badge_font_size(
            delta_badges[13][0],
            base_size_pt=_SUMMARY_DELTA_FONT_SIZE_PT,
        ),
    )
    _configure_summary_delta_badge(
        slide,
        shape_index=19,
        card_shape_index=15,
        font_size_pt=delta_badge_font_size(
            delta_badges[19][0],
            base_size_pt=_SUMMARY_DELTA_FONT_SIZE_PT,
        ),
    )

    _set_shape_font_name(slide, shape_index=3, font_name=_PPT_FONT_HEAD)
    for idx in (2, 4, 5, 6, 9, 10, 12, 13, 15, 16, 19):
        _set_shape_font_name(slide, shape_index=idx, font_name=_PPT_FONT_BODY_MEDIUM)

    _remove_summary_legacy_artifacts(slide)


def _update_cover_period(slide: Any, *, period_label: str) -> None:
    candidates: list[tuple[int, Any]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = str(getattr(shape, "text", "") or "")
        lower = text.lower()
        if "periodo" not in lower:
            continue

        score = 0
        # Main template placeholder: "Periodo dd/mm - dd/mm yyyy".
        if "dd/mm" in lower:
            score += 100
        if lower.strip().startswith("periodo"):
            score += 20

        # Corporate cover period ribbon is yellow and sits in lower area.
        try:
            fill = shape.fill
            if int(fill.type or 0) == 1:
                rgb = getattr(fill.fore_color, "rgb", None)
                if rgb is not None:
                    score += 40 if int(getattr(rgb, "blue", 255)) < 170 else 0
        except Exception:
            pass
        try:
            score += int(int(shape.top) / 100000)
        except Exception:
            pass

        # Avoid replacing explanatory subtitle "...análisis del periodo".
        if "kpi" in lower or "análisis" in lower or "analisis" in lower:
            score -= 30
        candidates.append((score, shape))

    if not candidates:
        return
    target = max(candidates, key=lambda item: item[0])[1]
    if not getattr(target, "has_text_frame", False):
        return

    tf = target.text_frame
    sample_run = None
    try:
        sample_run = tf.paragraphs[0].runs[0]
    except Exception:
        sample_run = None

    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = _single_line_ellipsis(period_label, max_chars=42)

    if sample_run is not None:
        try:
            run.font.bold = sample_run.font.bold
        except Exception:
            pass
        try:
            run.font.italic = sample_run.font.italic
        except Exception:
            pass
        try:
            run.font.name = sample_run.font.name
        except Exception:
            pass
        try:
            rgb = getattr(getattr(sample_run.font, "color", None), "rgb", None)
            if rgb is not None:
                run.font.color.rgb = rgb
        except Exception:
            pass

    # Keep the period text readable and centered inside the yellow ribbon.
    try:
        run.font.name = _PPT_FONT_BODY_MEDIUM
    except Exception:
        pass
    run.font.size = Pt(13.5)
    run.font.bold = True
    try:
        if getattr(getattr(run.font, "color", None), "rgb", None) is None:
            run.font.color.rgb = RGBColor(*hex_to_rgb(BBVA_LIGHT.midnight))
    except Exception:
        pass
    try:
        tf.margin_left = 20_000
        tf.margin_right = 20_000
        tf.margin_top = 4_000
        tf.margin_bottom = 4_000
    except Exception:
        pass
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass
    try:
        tf.word_wrap = False
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def _cover_title_shape_or_none(slide: Any) -> Any | None:
    candidates: list[tuple[int, Any]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = str(getattr(shape, "text", "") or "").strip()
        if not text:
            continue
        lower = text.lower()
        score = 0
        if "quincena" in lower or "periodo" in lower or "dd/mm" in lower:
            score -= 1_000
        if any(
            token in lower
            for token in ("incidencia", "abiertas", "critic", "funcionalidad", "días", "dias")
        ):
            score += 80
        score += int(_shape_area_in2(shape) * 10.0)
        try:
            score -= int(int(shape.top) / 200000)
        except Exception:
            pass
        candidates.append((score, shape))
    if not candidates:
        return None
    return max(candidates, key=lambda item: item[0])[1]


def _remove_cover_obsolete_subtitle(slide: Any) -> None:
    obsolete = _normalize_lookup_token(_COVER_REMOVED_SUBTITLE)
    for shape in list(getattr(slide, "shapes", [])):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = str(getattr(shape, "text", "") or "")
        if obsolete and obsolete in _normalize_lookup_token(text):
            _remove_shape(shape)


def _update_followup_cover(slide: Any, *, period_label: str) -> None:
    _remove_cover_obsolete_subtitle(slide)
    title_shape = _cover_title_shape_or_none(slide)
    if title_shape is not None:
        _set_shape_text_strict_by_shape(title_shape, _COVER_TITLE_TEXT)
    _update_cover_period(slide, period_label=period_label)


def _update_risk_issue_cover(slide: Any, *, title: str, period_label: str) -> None:
    title_shape = _cover_title_shape_or_none(slide)
    if title_shape is not None:
        _set_shape_text_strict_by_shape(title_shape, str(title or "").strip())
    _update_cover_period(slide, period_label=period_label)


def _fmt_avg_days(value: float) -> str:
    if pd.isna(value):
        return "0"
    safe = max(float(value or 0.0), 0.0)
    return str(int(round(safe)))


def _write_functionality_total_open_badge(
    slide: Any,
    *,
    card_shape_index: int,
    shape_index: int,
    total_open: int,
    critical_wording: bool,
) -> None:
    card_shape = _shape_or_none(slide, card_shape_index)
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return

    if card_shape is not None:
        try:
            horizontal_pad = int(card_shape.width * 0.045)
            vertical_pad = int(card_shape.height * 0.058)
            shape.left = int(card_shape.left) + horizontal_pad
            shape.top = int(card_shape.top) + vertical_pad
            shape.width = max(int(card_shape.width) - (2 * horizontal_pad), 1)
            shape.height = max(int(card_shape.height) - (2 * vertical_pad), 1)
        except Exception:
            pass

    tf = shape.text_frame
    tf.clear()
    try:
        tf.word_wrap = False
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass

    number_size_pt = 40.0 if int(abs(total_open)) >= 100 else 44.0
    label_size_pt = 14.0 if critical_wording else 16.0
    label_lines = (
        ("INCIDENCIAS CRÍTICAS", "ABIERTAS") if critical_wording else ("INCIDENCIAS ABIERTAS",)
    )

    p0 = tf.paragraphs[0]
    _set_paragraph_single_run(
        p0,
        text=str(int(total_open)),
        size_pt=number_size_pt,
        bold=True,
        color_rgb=RGBColor(255, 255, 255),
    )
    try:
        p0.alignment = PP_ALIGN.LEFT
    except Exception:
        pass

    for line in label_lines:
        p = tf.add_paragraph()
        _set_paragraph_single_run(
            p,
            text=line,
            size_pt=label_size_pt,
            bold=True,
            color_rgb=RGBColor(255, 255, 255),
        )
        try:
            p.alignment = PP_ALIGN.LEFT
        except Exception:
            pass


def _top_row_line(row: FunctionalityTopRow) -> str:
    count = int(getattr(row, "new_count", 0) or 0)
    count_txt = "incidencia nueva" if count == 1 else "incidencias nuevas"
    functionality = str(getattr(row, "functionality", "") or "").strip()
    open_total = int(getattr(row, "open_total", 0) or 0)
    avg_days = _fmt_avg_days(float(getattr(row, "avg_open_days", 0.0) or 0.0))
    line = (
        f"{count} {count_txt} en {functionality} ({open_total} en total - {avg_days} días promedio)"
    )
    return _single_line_ellipsis(line, max_chars=118)


def _mitigation_status_line(
    *,
    label: str,
    count: int,
    avg_open_days: float,
    in_fortnight: bool = False,
    rest_open: bool = False,
    max_chars: int = 104,
) -> str:
    days = _fmt_avg_days(avg_open_days)
    if rest_open:
        text = f"{label}: {int(count)} incidencias con {days} días de promedio abiertas"
    elif in_fortnight:
        text = (
            f"{label}: {int(count)} incidencias, en la quincena, "
            f"con {days} días de promedio en el estado"
        )
    else:
        text = f"{label}: {int(count)} incidencias con {days} días de promedio en el estado"
    return _single_line_ellipsis(text, max_chars=max_chars)


def _write_mitigation_status_line(slide: Any, shape_index: int, text: str) -> None:
    _set_shape_text_strict(slide, shape_index, text)
    _set_shape_font_size(
        slide,
        shape_index=shape_index,
        font_size_pt=8.2,
        bold=True,
        disable_autofit=True,
    )
    _set_shape_font_name(slide, shape_index=shape_index, font_name=_PPT_FONT_BODY_MEDIUM)
    shape = _shape_or_none(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.margin_left = 12_000
        tf.margin_right = 12_000
        tf.margin_top = 4_000
        tf.margin_bottom = 4_000
    except Exception:
        pass


def _root_cause_caption(zoom: FunctionalityZoomSlide, *, critical_wording: bool) -> str:
    issue_count = int(zoom.current_open_critical_count or 0)
    if issue_count <= 0:
        if critical_wording:
            return "Sin incidencias críticas abiertas de la quincena para esta funcionalidad."
        return "Sin incidencias abiertas de la quincena para esta funcionalidad."

    roots = [
        item
        for item in list(zoom.root_causes or [])
        if str(getattr(item, "label", "") or "").strip().lower() != "sin detalle suficiente"
    ]
    if not roots:
        return f"{issue_count} incidencias sin señal suficiente de causa raíz predominante."

    if len(roots) == 1:
        item = roots[0]
        cause = str(getattr(item, "label", "") or "").strip() or "Sin detalle"
        qty = int(getattr(item, "count", 0) or 0)
        noun = "incidencia" if qty == 1 else "incidencias"
        verb = "fue causada" if qty == 1 else "fueron causadas"
        return f"{qty} {noun} {verb} por {cause}."

    chunks: list[str] = []
    for item in roots:
        count = int(getattr(item, "count", 0) or 0)
        cause = str(getattr(item, "label", "") or "").strip() or "Sin detalle"
        chunks.append(f"{count} por {cause}")
    if len(chunks) == 2:
        detail = f"{chunks[0]} y {chunks[1]}"
    else:
        detail = ", ".join(chunks[:-1]) + f" y {chunks[-1]}"
    return f"Causas raíz detectadas: {detail}."


def _functionality_dashboard_table_target_geometry(
    slide: Any,
    *,
    table_shape_index: int,
) -> tuple[int, int, int, int] | None:
    table_shape = _shape_table_or_none(slide, table_shape_index)
    if table_shape is None:
        return None

    left = int(table_shape.left)
    top = int(table_shape.top)
    width = int(table_shape.width)
    height = int(table_shape.height)
    if width <= 0 or height <= 0:
        return None

    top_anchor_indexes = (6, 7, 8, 9, 10, 11)
    top_anchor_bottom = 0
    for idx in top_anchor_indexes:
        anchor = _shape_or_none(slide, idx)
        if anchor is None:
            continue
        top_anchor_bottom = max(top_anchor_bottom, int(anchor.top) + int(anchor.height))

    new_top = max(top_anchor_bottom + _FUNCTIONALITY_TABLE_GAP_TOP_EMU, top)

    mitigation_panel = _shape_or_none(slide, 12)
    if mitigation_panel is not None:
        right_limit = int(mitigation_panel.left) - _FUNCTIONALITY_TABLE_GAP_RIGHT_EMU
        if right_limit > left + 150_000:
            width = right_limit - left
        bottom_limit = (
            int(mitigation_panel.top)
            + int(mitigation_panel.height)
            - _FUNCTIONALITY_TABLE_BOTTOM_GAP_EMU
        )
    else:
        bottom_limit = int(top + height)

    if bottom_limit <= new_top:
        return None
    new_height = max(bottom_limit - new_top, 1)
    return (left, new_top, width, new_height)


def _native_table_shape(
    slide: Any,
    *,
    table_shape_index: int,
    row_count: int,
    col_count: int,
    geometry: tuple[int, int, int, int],
) -> Any:
    placeholder = _shape_table_or_none(slide, table_shape_index)
    return rebuild_native_table_shape(
        slide,
        placeholder,
        rows=max(int(row_count or 0), 1),
        cols=max(int(col_count or 0), 1),
        geometry=geometry,
    )


def _issue_table_geometry(*, data_row_count: int) -> tuple[int, int, int, int]:
    rows = max(min(int(data_row_count or 0), _ISSUE_TABLE_ROWS_PER_SLIDE), 1)
    table_height = int(_ISSUE_TABLE_HEADER_HEIGHT) + int(_ISSUE_TABLE_ROW_HEIGHT) * rows
    return (
        int(_ISSUE_TABLE_LEFT),
        int(_ISSUE_TABLE_TOP),
        int(_ISSUE_TABLE_WIDTH),
        int(table_height),
    )


def _is_issue_comment_row(row: Sequence[str]) -> bool:
    values = list(row or [])
    if len(values) < 2:
        return False
    return not str(values[0] or "").strip() and str(values[1] or "").startswith(
        "Comentarios registrados:"
    )


def _populate_issue_native_table(
    slide: Any,
    *,
    table_shape_index: int,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    hyperlink_by_row: Mapping[int, str] | None = None,
) -> Any:
    table_headers = tuple(headers or ())
    if not table_headers:
        raise ValueError("Issue table headers are required.")
    data_rows = list(rows or [])
    if not data_rows:
        filler = [""] * max(len(table_headers) - 2, 0)
        data_rows = [["", "Sin incidencias para este criterio.", *filler]]
    base_row_height = int(_ISSUE_TABLE_ROW_HEIGHT)
    row_heights: list[int] = []
    for row in data_rows:
        if _is_issue_comment_row(row):
            row_heights.append(int(_ISSUE_TABLE_COMMENT_ROW_HEIGHT))
        elif any("\n(" in str(cell or "") for cell in list(row or [])):
            row_heights.append(int(base_row_height * 1.18))
        else:
            row_heights.append(base_row_height)
    geometry = _issue_table_geometry(data_row_count=len(data_rows))
    geometry = (
        geometry[0],
        geometry[1],
        geometry[2],
        int(_ISSUE_TABLE_HEADER_HEIGHT) + sum(row_heights or [base_row_height]),
    )
    table_shape = _native_table_shape(
        slide,
        table_shape_index=table_shape_index,
        row_count=len(data_rows) + 1,
        col_count=len(table_headers),
        geometry=geometry,
    )
    populate_native_table(
        table_shape,
        headers=table_headers,
        rows=data_rows,
        column_widths=native_column_widths(geometry[2], _ISSUE_TABLE_COLUMN_WEIGHTS),
        row_height=base_row_height,
        header_height=int(_ISSUE_TABLE_HEADER_HEIGHT),
        font_name=_ISSUE_TABLE_FONT_NAME,
        body_font_size_pt=_ISSUE_TABLE_BODY_FONT_SIZE_PT,
        header_font_size_pt=_ISSUE_TABLE_HEADER_FONT_SIZE_PT,
        left_align_cols=(0, 1, 2),
        center_align_cols=(3, 4, 5),
        hyperlink_by_row=hyperlink_by_row,
        zebra=True,
    )
    try:
        for data_row_idx, height in enumerate(row_heights, start=1):
            table_shape.table.rows[data_row_idx].height = max(int(height), 1)
    except Exception:
        pass
    return table_shape


def _populate_functionality_dashboard_native_table(
    slide: Any,
    *,
    rows: Sequence[Sequence[str]],
    target_geometry: tuple[int, int, int, int] | None,
) -> None:
    table_rows = list(rows or [])
    if not table_rows:
        table_rows = [["", "Sin incidencias abiertas para el criterio actual.", "", "", ""]]
    target = target_geometry or (
        int(Inches(0.28)),
        int(Inches(2.18)),
        int(Inches(5.45)),
        int(Inches(2.6)),
    )
    data_row_count = max(len(table_rows), 1)
    max_rows = max(_FUNCTIONALITY_DASHBOARD_TABLE_ROWS, data_row_count)
    available_body_h = max(int(target[3]) - int(_ISSUE_TABLE_HEADER_HEIGHT), 1)
    row_height = max(int(available_body_h / max(max_rows, 1)), int(Inches(0.48)))
    compact_height = int(_ISSUE_TABLE_HEADER_HEIGHT) + row_height * data_row_count
    geometry = (int(target[0]), int(target[1]), int(target[2]), min(int(target[3]), compact_height))
    table_shape = _native_table_shape(
        slide,
        table_shape_index=1,
        row_count=data_row_count + 1,
        col_count=len(_FUNCTIONALITY_DASHBOARD_TABLE_HEADERS),
        geometry=geometry,
    )
    populate_native_table(
        table_shape,
        headers=_FUNCTIONALITY_DASHBOARD_TABLE_HEADERS,
        rows=table_rows,
        column_widths=native_column_widths(
            geometry[2],
            _FUNCTIONALITY_DASHBOARD_TABLE_COLUMN_WEIGHTS,
        ),
        row_height=row_height,
        header_height=int(_ISSUE_TABLE_HEADER_HEIGHT),
        font_name=_FUNCTIONALITY_TABLE_FONT_NAME,
        body_font_size_pt=_FUNCTIONALITY_TABLE_BODY_FONT_SIZE_PT,
        header_font_size_pt=_FUNCTIONALITY_TABLE_HEADER_FONT_SIZE_PT,
        left_align_cols=(1,),
        center_align_cols=(0, 2, 3, 4),
        zebra=True,
    )


def _populate_functionality_dashboard_slide(
    slide: Any,
    *,
    summary: PeriodFunctionalityFollowupSummary,
) -> None:
    critical_wording = bool(getattr(summary, "is_critical_focus", False))
    table_rows: list[list[str]] = []
    for row in list(summary.tail_rows or [])[:_FUNCTIONALITY_DASHBOARD_TABLE_ROWS]:
        table_rows.append(
            [
                str(int(row.rank)),
                ellipsize_text(str(row.functionality or ""), max_chars=70),
                str(int(row.new_count or 0)),
                str(int(row.open_total or 0)),
                _fmt_avg_days(float(row.avg_open_days)),
            ]
        )

    _set_shape_text(
        slide,
        2,
        (
            "Top tres de las incidencias por funcionalidad identificadas en la "
            f"{str(summary.period_label or '').replace('Quincena ', 'quincena ')}"
        ),
    )
    _set_shape_font_name(slide, shape_index=2, font_name=_PPT_FONT_BODY_MEDIUM)
    _write_functionality_total_open_badge(
        slide,
        card_shape_index=4,
        shape_index=5,
        total_open=int(summary.total_open_critical),
        critical_wording=critical_wording,
    )

    top_rows = list(summary.top_rows or [])
    top_shapes = (6, 8, 10)
    for idx, shape_idx in enumerate(top_shapes):
        if idx < len(top_rows):
            _set_shape_text(slide, shape_idx, _top_row_line(top_rows[idx]))
        else:
            _set_shape_text(slide, shape_idx, "Sin incidencias nuevas para esta posición.")
        _set_shape_font_size(
            slide,
            shape_index=shape_idx,
            font_size_pt=11.6,
            bold=True,
            disable_autofit=True,
        )
        _set_shape_font_name(slide, shape_index=shape_idx, font_name=_PPT_FONT_BODY_MEDIUM)

    _write_mitigation_status_line(
        slide,
        13,
        _mitigation_status_line(
            label="Estado Ready to Verify",
            count=int(summary.mitigation_ready_to_verify.count),
            avg_open_days=summary.mitigation_ready_to_verify.avg_open_days,
            in_fortnight=True,
        ),
    )
    _write_mitigation_status_line(
        slide,
        19,
        _mitigation_status_line(
            label="Estado New",
            count=int(summary.mitigation_new.count),
            avg_open_days=summary.mitigation_new.avg_open_days,
        ),
    )
    _write_mitigation_status_line(
        slide,
        20,
        _mitigation_status_line(
            label="Estado bloqueadas",
            count=int(summary.mitigation_blocked.count),
            avg_open_days=summary.mitigation_blocked.avg_open_days,
        ),
    )
    _write_mitigation_status_line(
        slide,
        21,
        _mitigation_status_line(
            label="Resto",
            count=int(summary.mitigation_non_critical.count),
            avg_open_days=summary.mitigation_non_critical.avg_open_days,
            rest_open=True,
        ),
    )
    for idx in (5, 13, 19, 20, 21):
        _set_shape_font_name(slide, shape_index=idx, font_name=_PPT_FONT_BODY_MEDIUM)
    _set_shape_font_color(slide, shape_index=18, color_rgb=RGBColor(255, 255, 255))
    table_target_geometry = _functionality_dashboard_table_target_geometry(
        slide,
        table_shape_index=1,
    )
    _populate_functionality_dashboard_native_table(
        slide,
        rows=table_rows,
        target_geometry=table_target_geometry,
    )


def _populate_functionality_zoom_slide(
    slide: Any,
    *,
    zoom: FunctionalityZoomSlide,
    critical_wording: bool,
    issues_page: Sequence[_IssuePageItem] | None = None,
    page_number: int = 1,
    total_pages: int = 1,
    notes_by_key: Mapping[str, str] | None = None,
) -> None:
    functionality = str(zoom.functionality or "").strip() or "Sin funcionalidad"
    page_suffix = ""
    if int(total_pages or 0) > 1:
        roman = _to_roman(int(page_number or 1))
        page_suffix = f" ({roman})" if roman else f" ({int(page_number or 1)})"
    _set_shape_text(
        slide,
        1,
        f"Incidencias, en {functionality}, abiertas en la quincena{page_suffix}",
    )
    _set_shape_font_name(slide, shape_index=1, font_name=_PPT_FONT_BODY_MEDIUM)
    _set_shape_text(
        slide,
        3,
        _trim_text(_root_cause_caption(zoom, critical_wording=critical_wording), max_chars=200),
    )
    _set_shape_font_color(slide, shape_index=3, color_rgb=RGBColor(*_TABLE_BODY_FG_RGB))
    _set_shape_font_name(slide, shape_index=3, font_name=_PPT_FONT_BODY)
    _set_shape_text(
        slide,
        4,
        (
            "Zoom de incidencias críticas del periodo:"
            if critical_wording
            else "Zoom de incidencias del periodo:"
        ),
    )
    _set_shape_font_name(slide, shape_index=4, font_name=_PPT_FONT_BODY_MEDIUM)

    page_issues = _coerce_issue_page_items(
        list(issues_page if issues_page is not None else zoom.issues or []),
        notes_by_key=notes_by_key,
    )
    rows: list[list[str]] = []
    row_links: dict[int, str] = {}
    comment_by_row: dict[int, str] = {}
    for page_item in page_issues:
        issue = cast(FunctionalityIssueRow, page_item.issue)
        issue_key = str(issue.key or "").strip().upper()
        issue_summary = _premium_sentence_case(str(issue.summary or ""))
        issue_root_cause = _premium_sentence_case(str(issue.root_cause or ""))
        main_row_idx = len(rows)
        rows.append(
            [
                issue_key,
                ellipsize_text(str(issue_summary or "").replace("/", " / "), max_chars=125),
                ellipsize_text(str(issue_root_cause or "").replace("/", " / "), max_chars=65),
                ellipsize_text(str(issue.status or ""), max_chars=28),
                ellipsize_text(str(issue.priority or ""), max_chars=18),
                f"{int(issue.open_days or 0)} días",
            ]
        )
        if str(issue.url or "").strip():
            row_links[main_row_idx] = str(issue.url).strip()
        for comment in page_item.comment_chunks:
            comment_row_idx = len(rows)
            rows.append(_issue_comment_row(comment))
            comment_by_row[comment_row_idx] = comment
    table_shape = _populate_issue_native_table(
        slide,
        table_shape_index=2,
        headers=_FUNCTIONALITY_ISSUE_TABLE_HEADERS,
        rows=rows,
        hyperlink_by_row=row_links,
    )
    _style_issue_comment_rows(table_shape, comment_by_row=comment_by_row)


def _chunk_risk_issues(
    issues: Sequence[PeriodRiskIssueRow],
    *,
    rows_per_slide: int,
    notes_by_key: Mapping[str, str] | None = None,
) -> list[tuple[_IssuePageItem, ...]]:
    return _chunk_issue_page_items(
        issues,
        rows_per_slide=rows_per_slide,
        notes_by_key=notes_by_key,
    )


def _assignee_with_po_text(
    assignee: object,
    po_team_leader: object,
    *,
    assignee_max_chars: int = 48,
    po_max_chars: int = 44,
) -> str:
    assignee_text = str(assignee or "").strip() or "(sin asignar)"
    po_text = str(po_team_leader or "").strip()
    assignee_text = ellipsize_text(assignee_text, max_chars=assignee_max_chars)
    if not po_text:
        return assignee_text
    po_text = ellipsize_text(po_text, max_chars=po_max_chars)
    return f"{assignee_text}\n({po_text})"


def _enrich_po_team_leader_from_sources(df: pd.DataFrame, settings: Settings) -> pd.DataFrame:
    safe = df if isinstance(df, pd.DataFrame) else pd.DataFrame()
    if safe.empty or "source_id" not in safe.columns:
        return safe
    source_po = {
        str(source.get("source_id") or "").strip(): str(source.get("po_team_leader") or "").strip()
        for source in jira_sources(settings)
        if str(source.get("source_id") or "").strip()
        and str(source.get("po_team_leader") or "").strip()
    }
    if not source_po:
        return safe
    out = safe.copy(deep=False)
    if "po_team_leader" not in out.columns:
        out["po_team_leader"] = ""
    current_po = out["po_team_leader"].fillna("").astype(str).str.strip()
    source_ids = out["source_id"].fillna("").astype(str).str.strip()
    out["po_team_leader"] = [
        po or source_po.get(source_id, "")
        for po, source_id in zip(current_po.tolist(), source_ids.tolist())
    ]
    return out


def _risk_issue_rows_for_table(
    issues: Sequence[PeriodRiskIssueRow | _IssuePageItem],
    *,
    empty_message: str,
    notes_by_key: Mapping[str, str] | None = None,
) -> tuple[list[list[str]], dict[int, str], dict[int, str]]:
    rows: list[list[str]] = []
    row_links: dict[int, str] = {}
    comment_by_row: dict[int, str] = {}
    page_items = _coerce_issue_page_items(list(issues or []), notes_by_key=notes_by_key)
    for page_item in page_items:
        issue = cast(PeriodRiskIssueRow, page_item.issue)
        issue_key = str(issue.key or "").strip().upper()
        main_row_idx = len(rows)
        rows.append(
            [
                issue_key,
                ellipsize_text(
                    _premium_sentence_case(str(issue.summary or "")),
                    max_chars=125,
                ),
                _assignee_with_po_text(issue.assignee, issue.po_team_leader),
                ellipsize_text(str(issue.status or ""), max_chars=28),
                ellipsize_text(str(issue.priority or ""), max_chars=18),
                f"{int(issue.open_days or 0)} días",
            ]
        )
        if str(issue.url or "").strip():
            row_links[main_row_idx] = str(issue.url or "").strip()
        for comment in page_item.comment_chunks:
            comment_row_idx = len(rows)
            rows.append(_issue_comment_row(comment))
            comment_by_row[comment_row_idx] = comment
    if not rows:
        rows.append(
            ["", str(empty_message or "Sin incidencias para este criterio."), "", "", "", ""]
        )
    return rows, row_links, comment_by_row


def _populate_risk_issue_list_slide(
    slide: Any,
    *,
    title: str,
    order_note: str,
    issues_page: Sequence[PeriodRiskIssueRow | _IssuePageItem],
    empty_message: str,
    page_number: int,
    total_pages: int,
    notes_by_key: Mapping[str, str] | None = None,
) -> None:
    _ = total_pages
    roman = _to_roman(int(page_number or 1))
    page_suffix = f" ({roman})" if roman else f" ({int(page_number or 1)})"
    _set_shape_text(slide, 1, f"{str(title or '').strip()}{page_suffix}")
    _set_shape_font_name(slide, shape_index=1, font_name=_PPT_FONT_BODY_MEDIUM)
    _set_shape_text(
        slide,
        3,
        str(order_note or "").strip(),
    )
    _set_shape_font_color(slide, shape_index=3, color_rgb=RGBColor(*_TABLE_BODY_FG_RGB))
    _set_shape_font_name(slide, shape_index=3, font_name=_PPT_FONT_BODY)
    _set_shape_text(slide, 4, "")
    _set_shape_font_name(slide, shape_index=4, font_name=_PPT_FONT_BODY_MEDIUM)

    rows, row_links, comment_by_row = _risk_issue_rows_for_table(
        issues_page,
        empty_message=empty_message,
        notes_by_key=notes_by_key,
    )
    table_shape = _populate_issue_native_table(
        slide,
        table_shape_index=2,
        headers=_RISK_ASSIGNEE_TABLE_HEADERS,
        rows=rows,
        hyperlink_by_row=row_links,
    )
    _style_issue_comment_rows(table_shape, comment_by_row=comment_by_row)
    LOGGER.info(
        "period_followup_slide_rows",
        extra={
            "run_id": uuid4().hex[:12],
            "slide_name": str(title or "").strip(),
            "rows_generated": int(len(issues_page or ())),
        },
    )


def _append_period_risk_issue_cover(
    prs: Any,
    *,
    cover_template_slide: Any,
    title: str,
    period_label: str,
) -> None:
    slide = _append_slide_clone_from_source(prs, source_slide=cover_template_slide)
    _update_risk_issue_cover(slide, title=title, period_label=period_label)


def _append_period_risk_issue_section(
    prs: Any,
    *,
    title: str,
    order_note: str,
    issues: Sequence[PeriodRiskIssueRow],
    empty_message: str,
    period_label: str,
    cover_template_slide: Any,
    zoom_template_slide: Any,
    notes_by_key: Mapping[str, str] | None = None,
) -> None:
    _append_period_risk_issue_cover(
        prs,
        cover_template_slide=cover_template_slide,
        title=title,
        period_label=period_label,
    )
    pages = _chunk_risk_issues(
        tuple(issues or ()),
        rows_per_slide=_ISSUE_TABLE_ROWS_PER_SLIDE,
        notes_by_key=notes_by_key,
    )
    total_pages = len(pages)
    for page_idx, page_rows in enumerate(pages, start=1):
        slide = _append_slide_clone_from_source(prs, source_slide=zoom_template_slide)
        _populate_risk_issue_list_slide(
            slide,
            title=title,
            order_note=order_note,
            issues_page=page_rows,
            empty_message=empty_message,
            page_number=page_idx,
            total_pages=total_pages,
            notes_by_key=notes_by_key,
        )


def _append_period_risk_issue_sections(
    prs: Any,
    *,
    period_label: str,
    high_priority_issues: Sequence[PeriodRiskIssueRow],
    aged_issues: Sequence[PeriodRiskIssueRow],
    notes_by_key: Mapping[str, str] | None = None,
) -> None:
    template_path = _resolve_functionality_template_path()
    template_prs = Presentation(str(template_path))
    if len(template_prs.slides) < 3:
        raise ValueError("La plantilla de funcionalidad debe contener la slide de zoom.")
    zoom_template_slide = template_prs.slides[2]
    if len(prs.slides) < 2:
        raise ValueError("La plantilla de periodo debe contener la slide de portada de sección.")
    cover_template_slide = prs.slides[1]
    specs = (
        (
            "Incidencias abiertas por criticidad alta",
            _RISK_HIGH_PRIORITY_ORDER_NOTE,
            tuple(high_priority_issues or ()),
            "Sin incidencias abiertas de criticidad alta en el scope actual.",
        ),
        (
            "Incidencias abiertas con más de 30 días",
            _RISK_AGED_ORDER_NOTE,
            tuple(aged_issues or ()),
            "Sin incidencias abiertas con más de 30 días en el scope actual.",
        ),
    )
    for title, order_note, issues, empty_message in specs:
        _append_period_risk_issue_section(
            prs,
            title=title,
            order_note=order_note,
            issues=issues,
            empty_message=empty_message,
            period_label=period_label,
            cover_template_slide=cover_template_slide,
            zoom_template_slide=zoom_template_slide,
            notes_by_key=notes_by_key,
        )


def _chunk_finalist_discrepancy_issues(
    issues: Sequence[FinalistDiscrepancyIssueRow],
    *,
    rows_per_slide: int,
) -> list[tuple[_FinalistIssuePageItem, ...]]:
    size = max(float(rows_per_slide or 0), 1.0)
    items = list(issues or [])
    if not items:
        return [tuple()]

    chunks: list[tuple[_FinalistIssuePageItem, ...]] = []
    current: list[_FinalistIssuePageItem] = []
    current_units = 0.0

    def flush() -> None:
        nonlocal current, current_units
        if current:
            chunks.append(tuple(current))
            current = []
            current_units = 0.0

    for issue in items:
        comment_chunks = _issue_comment_chunks(issue.comment)
        if not comment_chunks:
            if current and current_units + 1.0 > size:
                flush()
            current.append(_FinalistIssuePageItem(issue=issue))
            current_units += 1.0
            continue

        remaining = list(comment_chunks)
        while remaining:
            minimum_units = 1.0 + _ISSUE_TABLE_COMMENT_ROW_UNITS
            if current and current_units + minimum_units > size:
                flush()
                continue

            space_units = max(size - current_units, minimum_units)
            max_comment_chunks = int((space_units - 1.0) // _ISSUE_TABLE_COMMENT_ROW_UNITS)
            max_comment_chunks = max(max_comment_chunks, 1)
            take = min(len(remaining), max_comment_chunks)
            current.append(
                _FinalistIssuePageItem(issue=issue, comment_chunks=tuple(remaining[:take]))
            )
            current_units += 1.0 + (_ISSUE_TABLE_COMMENT_ROW_UNITS * float(take))
            remaining = remaining[take:]
            if remaining:
                flush()

    flush()
    return chunks or [tuple()]


def _finalist_discrepancy_rows_for_table(
    issues: Sequence[FinalistDiscrepancyIssueRow | _FinalistIssuePageItem],
    *,
    empty_message: str,
) -> tuple[
    list[list[str]],
    dict[int, str],
    dict[int, str],
    dict[int, FinalistDiscrepancyIssueRow],
    dict[int, str],
]:
    rows: list[list[str]] = []
    row_links: dict[int, str] = {}
    description_by_row: dict[int, str] = {}
    issue_by_row: dict[int, FinalistDiscrepancyIssueRow] = {}
    comment_by_row: dict[int, str] = {}
    page_items: list[_FinalistIssuePageItem] = []
    for item in list(issues or []):
        if isinstance(item, _FinalistIssuePageItem):
            page_items.append(item)
        else:
            page_items.append(
                _FinalistIssuePageItem(
                    issue=item,
                    comment_chunks=_issue_comment_chunks(getattr(item, "comment", "")),
                )
            )
    for page_item in page_items:
        issue = page_item.issue
        jira_key = str(issue.jira_key or "").strip().upper()
        description_text = (
            f"JIRA: {str(issue.jira_summary or '').strip() or 'Sin título JIRA'}\n"
            f"Helix: {str(issue.helix_text or '').strip() or 'Sin descripción Helix'}"
        )
        main_row_idx = len(rows)
        rows.append(
            [
                jira_key,
                ellipsize_text(description_text, max_chars=150),
                _assignee_with_po_text(issue.jira_assignee, issue.po_team_leader),
                (
                    f"JIRA: {ellipsize_text(str(issue.jira_status or ''), max_chars=22)}\n"
                    f"Helix: {ellipsize_text(str(issue.helix_status or ''), max_chars=22)}"
                ),
                ellipsize_text(str(issue.jira_priority or ""), max_chars=18),
                f"{int(issue.jira_open_days or 0)} días",
            ]
        )
        if str(issue.jira_url or "").strip():
            row_links[main_row_idx] = str(issue.jira_url or "").strip()
        description_by_row[main_row_idx] = rows[-1][1]
        issue_by_row[main_row_idx] = issue
        for comment in page_item.comment_chunks:
            comment_row_idx = len(rows)
            rows.append(_issue_comment_row(comment))
            comment_by_row[comment_row_idx] = comment
    if not rows:
        rows.append(
            ["", str(empty_message or "Sin incidencias para este criterio."), "", "", "", ""]
        )
    return rows, row_links, description_by_row, issue_by_row, comment_by_row


def _write_linkified_issue_cell(
    cell: Any,
    text: str,
    *,
    jira_url: str,
    helix_id: str,
    helix_url: str,
) -> None:
    tf = getattr(cell, "text_frame", None)
    if tf is None:
        return
    try:
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
    except Exception:
        pass
    try:
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass
    jira_key = ""
    if jira_url:
        match = re.search(_ALPHANUM_KEY_RE, str(jira_url or ""))
        jira_key = str(match.group(0)).upper() if match else ""
    segments_by_line = [
        linkify_issue_references(
            line,
            jira_urls={jira_key: jira_url} if jira_key and jira_url else {},
            helix_urls={str(helix_id or "").upper(): helix_url} if helix_url else {},
        )
        for line in str(text or "").splitlines()
    ]
    if not segments_by_line:
        segments_by_line = [()]
    body_rgb = RGBColor(*_TABLE_BODY_FG_RGB)
    for line_idx, segments in enumerate(segments_by_line):
        paragraph = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if not segments:
            run = paragraph.add_run()
            run.text = ""
            run.font.size = Pt(_ISSUE_TABLE_BODY_FONT_SIZE_PT)
            run.font.name = _ISSUE_TABLE_FONT_NAME
            continue
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment.text
            run.font.size = Pt(_ISSUE_TABLE_BODY_FONT_SIZE_PT)
            run.font.name = _ISSUE_TABLE_FONT_NAME
            try:
                run.font.color.rgb = body_rgb
            except Exception:
                pass
            if segment.url:
                try:
                    run.hyperlink.address = segment.url
                    run.font.underline = True
                except Exception:
                    pass


def _linkify_finalist_description_cells(
    table_shape: Any,
    *,
    issue_by_row: Mapping[int, FinalistDiscrepancyIssueRow],
    description_by_row: Mapping[int, str],
) -> None:
    if table_shape is None or not getattr(table_shape, "has_table", False):
        return
    table = table_shape.table
    for idx, issue in dict(issue_by_row or {}).items():
        if idx not in description_by_row:
            continue
        try:
            cell = table.cell(idx + 1, 1)
        except Exception:
            continue
        _write_linkified_issue_cell(
            cell,
            description_by_row[idx],
            jira_url=str(issue.jira_url or "").strip(),
            helix_id=str(issue.helix_id or "").strip().upper(),
            helix_url=str(issue.helix_url or "").strip(),
        )


def _frame_records(df: pd.DataFrame | None) -> list[Mapping[str, object]]:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return []
    return cast(list[Mapping[str, object]], df.to_dict(orient="records"))


def _load_report_notes_by_key(settings: Settings) -> dict[str, str]:
    try:
        store = NotesStore(Path(settings.NOTES_PATH))
        store.load()
        return {
            str(key or "").strip().upper(): str(note or "").strip()
            for key, note in store.latest_items()
            if str(key or "").strip() and str(note or "").strip()
        }
    except Exception:
        return {}


def _concat_discrepancy_frames(*frames: pd.DataFrame | None) -> pd.DataFrame:
    safe_frames = [frame for frame in frames if isinstance(frame, pd.DataFrame) and not frame.empty]
    if not safe_frames:
        return pd.DataFrame()
    return pd.concat(safe_frames, ignore_index=True, sort=False).copy(deep=False)


def _build_report_issue_url_maps(
    *,
    all_df: pd.DataFrame,
    finalist_discrepancies: pd.DataFrame | None,
    settings: Settings,
) -> tuple[dict[str, str], dict[str, str]]:
    rows: list[Mapping[str, object]] = []
    rows.extend(_frame_records(all_df))
    rows.extend(_frame_records(finalist_discrepancies))
    return build_issue_url_maps(
        rows,
        jira_base_url=str(getattr(settings, "JIRA_BASE_URL", "") or "").strip(),
        helix_base_url=str(
            getattr(settings, "HELIX_ARSQL_DASHBOARD_URL", "")
            or getattr(settings, "HELIX_DASHBOARD_URL", "")
            or ""
        ).strip(),
    )


def _text_frame_link_targets(tf: Any) -> set[str]:
    out: set[str] = set()
    for paragraph in list(getattr(tf, "paragraphs", []) or []):
        for run in list(getattr(paragraph, "runs", []) or []):
            try:
                url = str(run.hyperlink.address or "").strip()
            except Exception:
                url = ""
            if url:
                out.add(url)
    return out


def _write_linkified_issue_reference_cell(
    cell: Any,
    text: str,
    *,
    jira_urls: Mapping[str, str],
    helix_urls: Mapping[str, str],
) -> None:
    tf = getattr(cell, "text_frame", None)
    if tf is None:
        return
    try:
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
    except Exception:
        pass
    try:
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass

    body_rgb = RGBColor(*_TABLE_BODY_FG_RGB)
    lines = str(text or "").splitlines() or [""]
    for line_idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        segments = linkify_issue_references(
            line,
            jira_urls=dict(jira_urls),
            helix_urls=dict(helix_urls),
        )
        if not segments:
            run = paragraph.add_run()
            run.text = ""
            run.font.size = Pt(_ISSUE_TABLE_BODY_FONT_SIZE_PT)
            run.font.name = _ISSUE_TABLE_FONT_NAME
            continue
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment.text
            run.font.size = Pt(_ISSUE_TABLE_BODY_FONT_SIZE_PT)
            run.font.name = _ISSUE_TABLE_FONT_NAME
            try:
                run.font.color.rgb = body_rgb
            except Exception:
                pass
            if segment.url:
                try:
                    run.hyperlink.address = segment.url
                    run.font.underline = True
                except Exception:
                    pass


def _linkify_helix_references_in_tables(
    prs: Any,
    *,
    jira_urls: Mapping[str, str],
    helix_urls: Mapping[str, str],
) -> None:
    if not helix_urls:
        return
    normalized_helix_urls = {
        str(key or "").strip().upper(): str(value or "").strip()
        for key, value in dict(helix_urls).items()
        if str(key or "").strip() and str(value or "").strip()
    }
    if not normalized_helix_urls:
        return
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    text = str(cell.text or "")
                    if not text or not HELIX_ID_RE.search(text):
                        continue
                    ids_in_cell = {
                        str(match.group(0) or "").strip().upper()
                        for match in HELIX_ID_RE.finditer(text)
                    }
                    cell_helix_urls = {
                        inc_id: normalized_helix_urls[inc_id]
                        for inc_id in ids_in_cell
                        if inc_id in normalized_helix_urls
                    }
                    if not cell_helix_urls:
                        continue
                    tf = getattr(cell, "text_frame", None)
                    linked_targets = _text_frame_link_targets(tf)
                    if set(cell_helix_urls.values()).issubset(linked_targets):
                        continue
                    _write_linkified_issue_reference_cell(
                        cell,
                        text,
                        jira_urls=jira_urls,
                        helix_urls=cell_helix_urls,
                    )


def _write_issue_comment_cell(cell: Any, comment: str) -> None:
    tf = getattr(cell, "text_frame", None)
    if tf is None:
        return
    try:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(246, 249, 255)
    except Exception:
        pass
    try:
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.025)
        tf.margin_bottom = Inches(0.025)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass
    body_rgb = RGBColor(*_TABLE_BODY_FG_RGB)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.space_before = Pt(0)
    p0.space_after = Pt(0)
    run0 = p0.add_run()
    run0.text = "Comentarios registrados"
    run0.font.size = Pt(7.6)
    run0.font.name = _ISSUE_TABLE_FONT_NAME
    run0.font.bold = True
    run0.font.color.rgb = body_rgb
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(0)
    p1.space_after = Pt(0)
    run1 = p1.add_run()
    run1.text = str(comment or "").strip()
    run1.font.size = Pt(7.8)
    run1.font.name = _ISSUE_TABLE_FONT_NAME
    run1.font.color.rgb = body_rgb


def _style_issue_comment_rows(
    table_shape: Any,
    *,
    comment_by_row: Mapping[int, str],
    last_column_idx: int = 5,
) -> None:
    if table_shape is None or not getattr(table_shape, "has_table", False):
        return
    table = table_shape.table
    comment_rows = {int(row_idx) for row_idx in dict(comment_by_row or {})}
    for data_row_idx, comment in dict(comment_by_row or {}).items():
        ppt_row_idx = int(data_row_idx) + 1
        main_row_idx = max(ppt_row_idx - 1, 1)
        while main_row_idx > 1 and (main_row_idx - 1) in comment_rows:
            main_row_idx -= 1
        try:
            table.cell(main_row_idx, 0).merge(table.cell(ppt_row_idx, 0))
        except Exception:
            pass
        try:
            merged = table.cell(ppt_row_idx, 1)
            merged.merge(table.cell(ppt_row_idx, max(int(last_column_idx or 0), 1)))
        except Exception:
            try:
                merged = table.cell(ppt_row_idx, 1)
            except Exception:
                continue
        _write_issue_comment_cell(merged, str(comment or "").strip())


def _style_finalist_status_cells(
    table_shape: Any,
    *,
    data_row_indices: Sequence[int],
) -> None:
    if table_shape is None or not getattr(table_shape, "has_table", False):
        return
    table = table_shape.table
    jira_rgb = RGBColor(*hex_to_rgb(BBVA_REPORT_RED_TEXT))
    helix_rgb = RGBColor(*hex_to_rgb(BBVA_REPORT_AMBER_TEXT))
    try:
        helix_rgb = RGBColor(34, 139, 74)
    except Exception:
        pass
    for data_row_idx in list(data_row_indices or []):
        ridx = int(data_row_idx) + 1
        if ridx <= 0 or ridx >= len(table.rows):
            continue
        try:
            cell = table.cell(ridx, 3)
        except Exception:
            continue
        original = str(cell.text or "")
        parts = [part for part in original.splitlines() if part.strip()]
        if len(parts) < 2:
            continue
        tf = cell.text_frame
        tf.clear()
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        except Exception:
            pass
        for idx, (line, color) in enumerate(((parts[0], jira_rgb), (parts[1], helix_rgb))):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(8.5)
            run.font.name = _ISSUE_TABLE_FONT_NAME
            run.font.bold = True
            run.font.color.rgb = color


def _populate_finalist_discrepancy_list_slide(
    slide: Any,
    *,
    title: str,
    order_note: str,
    issues_page: Sequence[FinalistDiscrepancyIssueRow | _FinalistIssuePageItem],
    empty_message: str,
    page_number: int,
    total_pages: int,
) -> None:
    _ = total_pages
    roman = _to_roman(int(page_number or 1))
    page_suffix = f" ({roman})" if roman else f" ({int(page_number or 1)})"
    _set_shape_text(slide, 1, f"{str(title or _FINALIST_DISCREPANCIES_TITLE)}{page_suffix}")
    _set_shape_font_name(slide, shape_index=1, font_name=_PPT_FONT_BODY_MEDIUM)
    _set_shape_text(slide, 3, str(order_note or _FINALIST_DISCREPANCIES_ORDER_NOTE))
    _set_shape_font_color(slide, shape_index=3, color_rgb=RGBColor(*_TABLE_BODY_FG_RGB))
    _set_shape_font_name(slide, shape_index=3, font_name=_PPT_FONT_BODY)
    _set_shape_text(slide, 4, "")
    _set_shape_font_name(slide, shape_index=4, font_name=_PPT_FONT_BODY_MEDIUM)

    rows, row_links, description_by_row, issue_by_row, comment_by_row = (
        _finalist_discrepancy_rows_for_table(
            issues_page,
            empty_message=empty_message,
        )
    )
    table_shape = _populate_issue_native_table(
        slide,
        table_shape_index=2,
        headers=_RISK_ASSIGNEE_TABLE_HEADERS,
        rows=rows,
        hyperlink_by_row=row_links,
    )
    _linkify_finalist_description_cells(
        table_shape,
        issue_by_row=issue_by_row,
        description_by_row=description_by_row,
    )
    _style_issue_comment_rows(table_shape, comment_by_row=comment_by_row)
    _style_finalist_status_cells(table_shape, data_row_indices=tuple(issue_by_row.keys()))
    LOGGER.info(
        "period_followup_slide_rows",
        extra={
            "run_id": uuid4().hex[:12],
            "slide_name": str(title or _FINALIST_DISCREPANCIES_TITLE),
            "rows_generated": int(len(issues_page or ())),
        },
    )


def _append_finalist_discrepancy_section(
    prs: Any,
    *,
    title: str = _FINALIST_DISCREPANCIES_TITLE,
    order_note: str = _FINALIST_DISCREPANCIES_ORDER_NOTE,
    empty_message: str = "Sin incidencias con discrepancias en estado finalista en el scope actual.",
    period_label: str,
    issues: Sequence[FinalistDiscrepancyIssueRow],
) -> None:
    template_path = _resolve_functionality_template_path()
    template_prs = Presentation(str(template_path))
    if len(template_prs.slides) < 3:
        raise ValueError("La plantilla de funcionalidad debe contener la slide de zoom.")
    zoom_template_slide = template_prs.slides[2]
    if len(prs.slides) < 2:
        raise ValueError("La plantilla de periodo debe contener la slide de portada de sección.")
    cover_template_slide = prs.slides[1]
    _append_period_risk_issue_cover(
        prs,
        cover_template_slide=cover_template_slide,
        title=str(title or _FINALIST_DISCREPANCIES_TITLE),
        period_label=period_label,
    )
    pages = _chunk_finalist_discrepancy_issues(
        tuple(issues or ()),
        rows_per_slide=_ISSUE_TABLE_ROWS_PER_SLIDE,
    )
    total_pages = len(pages)
    for page_idx, page_rows in enumerate(pages, start=1):
        slide = _append_slide_clone_from_source(prs, source_slide=zoom_template_slide)
        _populate_finalist_discrepancy_list_slide(
            slide,
            title=str(title or _FINALIST_DISCREPANCIES_TITLE),
            order_note=str(order_note or _FINALIST_DISCREPANCIES_ORDER_NOTE),
            issues_page=page_rows,
            empty_message=str(empty_message or "Sin incidencias para este criterio."),
            page_number=page_idx,
            total_pages=total_pages,
        )


def _filter_last_six_months_trend(trend: pd.DataFrame) -> pd.DataFrame:
    safe = trend if isinstance(trend, pd.DataFrame) else pd.DataFrame()
    if safe.empty or "quincena_start" not in safe.columns:
        return safe.copy(deep=False)

    starts = pd.to_datetime(safe["quincena_start"], errors="coerce")
    if "quincena_end" in safe.columns:
        ends = pd.to_datetime(safe["quincena_end"], errors="coerce")
    else:
        ends = starts
    max_candidates = [series.max() for series in (starts, ends) if bool(series.notna().any())]
    if not max_candidates:
        return safe.copy(deep=False)

    max_date = max(pd.Timestamp(value).normalize() for value in max_candidates)
    window_start = (max_date.to_period("M") - 5).to_timestamp()
    effective_start = starts.where(starts.notna(), ends)
    filtered = safe.loc[effective_start.ge(window_start)].copy(deep=False)
    if filtered.empty:
        return filtered
    return filtered.sort_values(
        ["quincena_start", "tema"], ascending=[True, True], kind="mergesort"
    )


def _functionality_fortnight_trend_png(*, open_df: pd.DataFrame) -> bytes:
    safe_open = open_df if isinstance(open_df, pd.DataFrame) else pd.DataFrame()
    if safe_open.empty:
        return b""

    trend = build_theme_fortnight_trend(safe_open, cumulative=True)
    if not isinstance(trend, pd.DataFrame) or trend.empty:
        return b""
    trend = _filter_last_six_months_trend(trend)
    if trend.empty:
        return b""

    raw_axis_labels = trend["quincena_label"].dropna().astype(str).drop_duplicates().tolist()
    if not raw_axis_labels:
        return b""
    axis_meta = (
        trend.loc[:, ["quincena_label", "quincena_start", "quincena_end"]]
        .drop_duplicates(subset=["quincena_label"])
        .copy(deep=False)
    )
    axis_meta["axis_label"] = [
        _format_quincena_axis_ym(start, end)
        for start, end in zip(axis_meta["quincena_start"], axis_meta["quincena_end"])
    ]
    axis_label_map = {
        str(raw): str(lbl) for raw, lbl in zip(axis_meta["quincena_label"], axis_meta["axis_label"])
    }
    axis_labels = list(dict.fromkeys(axis_label_map.get(lbl, lbl) for lbl in raw_axis_labels))

    theme_totals = (
        trend.groupby("tema", dropna=False)["issues_value"]
        .max()
        .sort_values(ascending=False)
        .fillna(0)
    )
    ordered_themes = order_theme_labels_by_volume(
        theme_totals.index.tolist(),
        counts_by_label=theme_totals,
        others_last=True,
    )
    ordering = build_theme_render_order(
        ordered_themes,
        counts_by_label=theme_totals,
        others_last=True,
        others_at_x_axis=True,
    )
    legend_order = list(ordering.display_order)
    stack_order = list(ordering.stack_order_bottom_to_top)
    if not legend_order or not stack_order:
        return b""
    # Keep stack geometry untouched (bottom->top), but show legend in inverse
    # order so the series closest to X axis appears last ("Otros").
    legend_order = list(reversed(stack_order))
    theme_color_map = build_theme_color_map(
        theme_order=list(ordering.display_order),
        dark_mode=False,
    )

    trend_local = trend.copy(deep=False)
    trend_local["axis_label"] = [
        axis_label_map.get(str(lbl), str(lbl)) for lbl in trend_local["quincena_label"].tolist()
    ]
    totals = (
        trend_local.groupby("axis_label", dropna=False)["issues_value"]
        .sum()
        .reindex(axis_labels)
        .fillna(0)
        .astype(int)
    )
    max_total = float(totals.max()) if not totals.empty else 0.0
    label_min_value = _inside_label_min_value(max_total)
    fig = go.Figure()
    legend_rank = {theme: idx for idx, theme in enumerate(legend_order)}
    for theme in stack_order:
        sub = trend_local.loc[trend_local["tema"].eq(theme)].copy(deep=False)
        values_series = (
            pd.to_numeric(sub.get("issues_value"), errors="coerce")
            .fillna(0.0)
            .groupby(sub.get("axis_label"))
            .sum()
            .reindex(axis_labels)
            .fillna(0.0)
        )
        values = values_series.astype(float).tolist()
        value_text = _safe_inside_bar_text(values, min_value=label_min_value)
        color_hex = str(theme_color_map.get(theme) or "#7784A0")
        text_color = "#FFFFFF"
        if _normalize_lookup_token(theme) in {
            _normalize_lookup_token("Monetarias"),
            _normalize_lookup_token("Transferencias"),
            _normalize_lookup_token("Softoken"),
        }:
            text_color = "#0B1F3B"
        fig.add_trace(
            go.Bar(
                x=axis_labels,
                y=values,
                name=str(theme),
                marker=dict(color=color_hex, line=dict(color="#F2F5FA", width=0.8)),
                text=value_text,
                textposition="inside",
                textfont=dict(size=EXEC_CHART_INSIDE_VALUE_FONT_PT, color=text_color),
                legendrank=int(legend_rank.get(theme, len(legend_rank))),
                customdata=[[int(totals.get(lbl, 0))] for lbl in axis_labels],
                hovertemplate=(
                    "Tema: %{fullData.name}<br>Quincena: %{x}<br>"
                    "Incidencias abiertas acumuladas: %{y}<br>"
                    "Total columna: %{customdata[0]}<extra></extra>"
                ),
            )
        )

    _add_stacked_bar_totals(
        fig,
        axis_labels=axis_labels,
        totals=[float(v) for v in totals.tolist()],
        max_total=max_total,
        color="#0B3E76",
    )
    _apply_executive_chart_layout(
        fig,
        kind="trend",
        show_legend=True,
        x_title="Quincena",
        y_title="Incidencias abiertas acumuladas",
        height=EXEC_CHART_TREND_EXPORT_HEIGHT,
        margin=_exec_chart_margin(l=78, r=50, t=76, b=198),
    )
    fig.update_layout(
        barmode="stack",
        bargap=0.19,
    )
    fig.update_xaxes(
        type="category",
        categoryorder="array",
        categoryarray=axis_labels,
    )

    payload = _fig_to_png_exact(
        fig,
        width=EXEC_CHART_EXPORT_WIDTH,
        height=EXEC_CHART_TREND_EXPORT_HEIGHT,
        scale=1.0,
    )
    return payload or b""


def _populate_functionality_trend_aggregate_slide(
    slide: Any,
    *,
    open_df: pd.DataFrame,
    slide_width: int,
    slide_height: int,
) -> None:
    _clear_slide_shapes(slide)

    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 248, 248)
    except Exception:
        pass

    slide_w = int(slide_width or 9_144_000)
    slide_h = int(slide_height or 5_143_500)
    margin_x = int(slide_w * 0.032)
    content_w = max(slide_w - (2 * margin_x), 1)

    title_box = _add_exec_textbox(
        slide,
        left=margin_x,
        top=int(slide_h * 0.03),
        width=content_w,
        height=int(slide_h * 0.065),
        text=_FUNCTIONALITY_TREND_AGGREGATE_TITLE,
        font_size_pt=20.5,
        color_rgb=RGBColor(4, 19, 139),
        font_name="Source Serif 4",
        bold=True,
    )
    try:
        title_box.text_frame.word_wrap = False
    except Exception:
        pass

    frame_left = margin_x
    frame_top = int(slide_h * 0.11)
    frame_width = content_w
    frame_height = max(slide_h - frame_top - int(slide_h * 0.03), int(slide_h * 0.80))
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        frame_left,
        frame_top,
        frame_width,
        frame_height,
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(247, 248, 252)
    frame.line.color.rgb = RGBColor(214, 220, 232)
    frame.line.width = Pt(1.0)

    chart_png = _functionality_fortnight_trend_png(open_df=open_df)
    if chart_png:
        _overlay_picture_contain(
            slide,
            payload=chart_png,
            frame_left=frame_left + 120_000,
            frame_top=frame_top + 82_000,
            frame_width=frame_width - 240_000,
            frame_height=frame_height - 164_000,
        )
    else:
        _add_exec_textbox(
            slide,
            left=frame_left + 60_000,
            top=frame_top + 2_450_000,
            width=frame_width - 120_000,
            height=290_000,
            text="No hay datos suficientes para construir la tendencia acumulada por funcionalidad.",
            font_size_pt=16.0,
            color_rgb=RGBColor(56, 67, 92),
            bold=False,
            align=PP_ALIGN.CENTER,
        )


def _append_functionality_followup_slides(
    prs: Any,
    *,
    summary: PeriodFunctionalityFollowupSummary,
    period_label: str,
    open_df: pd.DataFrame,
    slide_width: int,
    slide_height: int,
) -> None:
    critical_wording = bool(getattr(summary, "is_critical_focus", False))
    template_path = _resolve_functionality_template_path()
    template_prs = Presentation(str(template_path))
    if len(template_prs.slides) < 5:
        raise ValueError(
            "La plantilla de funcionalidad debe contener 5 slides (cabecera + dashboard + 3 zoom)."
        )

    header_slide = _append_slide_clone_from_source(prs, source_slide=template_prs.slides[0])
    trend_slide = _append_slide_clone_from_source(prs, source_slide=template_prs.slides[1])
    dashboard_slide = _append_slide_clone_from_source(prs, source_slide=template_prs.slides[1])

    # Slide 1 (cabecera funcionalidad)
    _set_shape_text(header_slide, 3, str(period_label or "").strip())
    _set_shape_text(
        header_slide,
        2,
        (
            "Detalle, de las incidencias críticas, abiertas por funcionalidad"
            if critical_wording
            else "Detalle, de las incidencias, abiertas por funcionalidad"
        ),
    )

    # Slide 2 (tendencia funcionalidad agregada)
    _populate_functionality_trend_aggregate_slide(
        trend_slide,
        open_df=open_df,
        slide_width=slide_width,
        slide_height=slide_height,
    )

    # Slide 3 (dashboard funcionalidad)
    _set_shape_text(
        dashboard_slide,
        3,
        (
            "Seguimiento de KPIs - Incidencias críticas abiertas por funcionalidad"
            if critical_wording
            else "Seguimiento de KPIs - Incidencias abiertas por funcionalidad"
        ),
    )
    _populate_functionality_dashboard_slide(dashboard_slide, summary=summary)


def _append_functionality_zoom_slides(
    prs: Any,
    *,
    summary: PeriodFunctionalityFollowupSummary,
    notes_by_key: Mapping[str, str] | None = None,
) -> None:
    critical_wording = bool(getattr(summary, "is_critical_focus", False))
    template_path = _resolve_functionality_template_path()
    template_prs = Presentation(str(template_path))
    if len(template_prs.slides) < 3:
        raise ValueError("La plantilla de funcionalidad debe contener la slide de zoom.")
    zoom_template_slide = template_prs.slides[2]
    # Zoom de top 3 funcionalidades con paginado por overflow.
    zooms = list(summary.zoom_slides or [])
    while len(zooms) < 3:
        zooms.append(
            FunctionalityZoomSlide(
                functionality=f"Sin funcionalidad {len(zooms) + 1}",
                current_open_critical_count=0,
                root_causes=(),
                issues=(),
            )
        )
    zooms = zooms[:3]

    zoom_page_specs: list[tuple[FunctionalityZoomSlide, tuple[_IssuePageItem, ...], int, int]] = []
    for zoom in zooms:
        pages = _chunk_zoom_issues(
            tuple(getattr(zoom, "issues", ()) or ()),
            rows_per_slide=_ISSUE_TABLE_ROWS_PER_SLIDE,
            notes_by_key=notes_by_key,
        )
        total_pages = len(pages)
        for page_idx, page_rows in enumerate(pages, start=1):
            zoom_page_specs.append((zoom, page_rows, page_idx, total_pages))

    zoom_target_slides = [
        _append_slide_clone_from_source(prs, source_slide=zoom_template_slide)
        for _ in zoom_page_specs
    ]
    for target_slide, spec in zip(zoom_target_slides, zoom_page_specs):
        zoom, page_rows, page_idx, total_pages = spec
        _populate_functionality_zoom_slide(
            target_slide,
            zoom=zoom,
            critical_wording=critical_wording,
            issues_page=page_rows,
            page_number=page_idx,
            total_pages=total_pages,
            notes_by_key=notes_by_key,
        )


def _load_or_scope_data(
    settings: Settings,
    *,
    country: str,
    source_ids: Sequence[str],
    dff_override: pd.DataFrame | None,
    open_df_override: pd.DataFrame | None,
    all_df_override: pd.DataFrame | None,
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    if dff_override is not None:
        dff = dff_override.copy(deep=False)
        all_df = (
            all_df_override.copy(deep=False)
            if isinstance(all_df_override, pd.DataFrame) and not all_df_override.empty
            else dff
        )
    else:
        base_df = load_issues_df(settings.DATA_PATH)
        all_df = base_df
        scoped = scope_country_sources(base_df, country=country, source_ids=source_ids)
        dff = apply_analysis_depth_filter(scoped, settings=settings)

    if open_df_override is not None:
        open_df = open_df_override.copy(deep=False)
    else:
        closed_mask = effective_closed_mask(dff)
        open_df = dff.loc[~closed_mask].copy(deep=False)
    return dff, open_df, all_df


def generate_country_period_followup_ppt(
    settings: Settings,
    *,
    country: str,
    source_ids: Sequence[str],
    dff_override: pd.DataFrame | None = None,
    open_df_override: pd.DataFrame | None = None,
    all_df_override: pd.DataFrame | None = None,
    finalist_discrepancies_override: pd.DataFrame | None = None,
    template_path: str | None = None,
    applied_filter_summary: str = "",
    functionality_status_filters: Sequence[str] | None = None,
    functionality_priority_filters: Sequence[str] | None = None,
    functionality_filters: Sequence[str] | None = None,
    reference_day: pd.Timestamp | str | None = None,
) -> PeriodFollowupReportResult:
    clean_source_ids = _clean_source_ids(source_ids)
    country_txt = str(country or "").strip()
    dff, _open_df, all_df = _load_or_scope_data(
        settings,
        country=country_txt,
        source_ids=clean_source_ids,
        dff_override=dff_override,
        open_df_override=open_df_override,
        all_df_override=all_df_override,
    )
    effective_finalist_discrepancies = (
        finalist_discrepancies_override
        if isinstance(finalist_discrepancies_override, pd.DataFrame)
        else pd.DataFrame()
    )
    if not effective_finalist_discrepancies.empty:
        dff = apply_effective_finalist_lookup_state(
            dff,
            discrepancies=effective_finalist_discrepancies,
            reference_window=reference_day,
        )
    elif isinstance(all_df, pd.DataFrame) and not all_df.empty:
        dff, effective_finalist_discrepancies = apply_effective_finalist_lookup_state_for_scope(
            dff,
            history_df=all_df,
            settings=settings,
            country=country_txt,
            source_ids=clean_source_ids,
            reference_day=reference_day,
        )
    root_cause_evolutives, effective_finalist_discrepancies = (
        split_root_cause_evolutive_discrepancies(
            effective_finalist_discrepancies,
            settings=settings,
            country=country_txt,
        )
    )
    if dff.empty:
        raise ValueError("No hay incidencias para generar el informe de seguimiento.")

    labels = source_label_map(settings, country=country_txt, source_ids=clean_source_ids)
    quincenal = build_country_quincenal_result(
        df=dff,
        settings=settings,
        country=country_txt,
        source_ids=clean_source_ids,
        source_label_by_id=labels,
        reference_day=pd.Timestamp(reference_day) if reference_day is not None else None,
    )
    template = _resolve_template_path(settings, explicit_path=template_path)
    prs = Presentation(str(template))
    slide_width_emu = _safe_emu(getattr(prs, "slide_width", None), default=9_144_000)
    slide_height_emu = _safe_emu(getattr(prs, "slide_height", None), default=5_143_500)

    # Normalize user template into canonical 8-slide structure.
    _normalize_period_template(prs)

    aggregate = quincenal.aggregate

    _update_followup_cover(
        prs.slides[0], period_label=format_window_label(aggregate.summary.window)
    )

    if clean_source_ids:
        insert_index = 5
        for _ in range(max(len(clean_source_ids) - 2, 0)):
            source_template = prs.slides[4] if len(prs.slides) > 4 else prs.slides[3]
            _append_slide_clone_from_source(prs, source_slide=source_template)
            _move_slide(prs, from_index=len(prs.slides) - 1, to_index=insert_index)
            insert_index += 1

        _populate_summary_slide(
            prs.slides[2],
            title=f"Seguimiento de incidencias - {country_txt.upper()} (vista agregada)",
            scope_result=aggregate,
        )
        aggregate_chart_anchor = _resolve_summary_chart_anchor(prs.slides[2])
        aggregate_chart_width, aggregate_chart_height = _summary_chart_export_size(
            aggregate_chart_anchor
        )
        _overlay_picture(
            prs.slides[2],
            anchor_shape=aggregate_chart_anchor,
            payload=_chart_png(
                settings,
                dff=aggregate.dff,
                open_df=aggregate.open_df,
                chart_id="timeseries",
                width=aggregate_chart_width,
                height=aggregate_chart_height,
                slide_optimized=True,
            ),
            replace_anchor=True,
            preserve_aspect=True,
        )

        for offset, source_id in enumerate(clean_source_ids):
            slide_index = 3 + offset
            source_scope = quincenal.by_source.get(source_id)
            if source_scope is None:
                continue
            source_label = labels.get(source_id, source_id).split("·")[0].strip().upper()
            _populate_summary_slide(
                prs.slides[slide_index],
                title=f"Seguimiento de incidencias - {source_label}",
                scope_result=source_scope,
            )
            source_chart_anchor = _resolve_summary_chart_anchor(prs.slides[slide_index])
            source_chart_width, source_chart_height = _summary_chart_export_size(
                source_chart_anchor
            )
            _overlay_picture(
                prs.slides[slide_index],
                anchor_shape=source_chart_anchor,
                payload=_chart_png(
                    settings,
                    dff=source_scope.dff,
                    open_df=source_scope.open_df,
                    chart_id="timeseries",
                    width=source_chart_width,
                    height=source_chart_height,
                    slide_optimized=True,
                ),
                replace_anchor=True,
                preserve_aspect=True,
            )

        extra_template_source_count = 2
        if len(clean_source_ids) < extra_template_source_count:
            for remove_idx in range(4, 3 + len(clean_source_ids), -1):
                if remove_idx < len(prs.slides):
                    _remove_slide(prs, remove_idx)
    else:
        LOGGER.info(
            "period_followup_no_rollups",
            extra={
                "run_id": uuid4().hex[:12],
                "country": country_txt,
                "message": (
                    "Sin orígenes agregados configurados; se omiten slides de vista agregada "
                    "y detalle por fuente."
                ),
            },
        )
        for remove_idx in (4, 3, 2):
            if remove_idx < len(prs.slides):
                _remove_slide(prs, remove_idx)

    aging_slide_index = 4 + len(clean_source_ids) if clean_source_ids else 3
    priority_slide_index = 5 + len(clean_source_ids) if clean_source_ids else 4

    _populate_open_aging_executive_slide(
        prs.slides[aging_slide_index],
        settings=settings,
        scope_result=aggregate,
        slide_width=slide_width_emu,
        slide_height=slide_height_emu,
    )
    _populate_open_priority_executive_slide(
        prs.slides[priority_slide_index],
        settings=settings,
        scope_result=aggregate,
        slide_width=slide_width_emu,
        slide_height=slide_height_emu,
    )

    risk_lists = build_period_risk_issue_lists(
        _enrich_po_team_leader_from_sources(aggregate.dff, settings),
        fallback_analysis_day=pd.Timestamp(aggregate.summary.window.current_end),
    )
    notes_by_key = _load_report_notes_by_key(settings)
    root_cause_evolutive_rows = build_finalist_discrepancy_issue_list(
        root_cause_evolutives,
        notes_by_key=notes_by_key,
    )
    finalist_discrepancy_rows = build_finalist_discrepancy_issue_list(
        effective_finalist_discrepancies,
        notes_by_key=notes_by_key,
    )
    functionality_followup = build_period_functionality_followup_summary(
        scope_result=aggregate,
        jira_base_url=str(getattr(settings, "JIRA_BASE_URL", "") or "").strip(),
        status_filters=list(functionality_status_filters or []),
        priority_filters=list(functionality_priority_filters or []),
        functionality_filters=list(functionality_filters or []),
        apply_default_status_when_empty=True,
        top_n=3,
        top_root_causes=3,
    )
    _append_functionality_followup_slides(
        prs,
        summary=functionality_followup,
        period_label=functionality_followup.period_label,
        open_df=aggregate.open_df,
        slide_width=slide_width_emu,
        slide_height=slide_height_emu,
    )
    _append_period_risk_issue_sections(
        prs,
        period_label=functionality_followup.period_label,
        high_priority_issues=risk_lists.high_priority,
        aged_issues=risk_lists.aged,
        notes_by_key=notes_by_key,
    )
    if _parse_bool_flag(
        getattr(settings, "PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED", "false"),
        default=False,
    ):
        _append_functionality_zoom_slides(
            prs,
            summary=functionality_followup,
            notes_by_key=notes_by_key,
        )
    if root_cause_evolutive_rows:
        _append_finalist_discrepancy_section(
            prs,
            title=_ROOT_CAUSE_EVOLUTIVES_TITLE,
            order_note=_ROOT_CAUSE_EVOLUTIVES_ORDER_NOTE,
            empty_message="Sin evolutivos de causa raíz en el scope actual.",
            period_label=functionality_followup.period_label,
            issues=root_cause_evolutive_rows,
        )
    if finalist_discrepancy_rows:
        _append_finalist_discrepancy_section(
            prs,
            period_label=functionality_followup.period_label,
            issues=finalist_discrepancy_rows,
        )

    all_finalist_link_discrepancies = _concat_discrepancy_frames(
        root_cause_evolutives,
        effective_finalist_discrepancies,
    )
    jira_url_map, helix_url_map = _build_report_issue_url_maps(
        all_df=all_df,
        finalist_discrepancies=all_finalist_link_discrepancies,
        settings=settings,
    )
    _linkify_helix_references_in_tables(
        prs,
        jira_urls=jira_url_map,
        helix_urls=helix_url_map,
    )

    _remove_slide_number_artifacts(prs)
    validate_shapes_inside_slide(prs)

    buff = BytesIO()
    prs.save(buff)
    content = buff.getvalue()
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M")
    file_name = (
        f"seguimiento-{_slug(country_txt)}-"
        f"{_slug('rollups' if clean_source_ids else 'sin-agregados')}-{stamp}.pptx"
    )
    total_issues = int(len(aggregate.dff))
    open_issues = int(len(aggregate.open_df))
    closed_issues = max(total_issues - open_issues, 0)
    return PeriodFollowupReportResult(
        file_name=file_name,
        content=content,
        slide_count=len(prs.slides),
        total_issues=total_issues,
        open_issues=open_issues,
        closed_issues=closed_issues,
        country=country_txt,
        source_ids=tuple(clean_source_ids),
        applied_filter_summary=str(applied_filter_summary or "").strip(),
    )
