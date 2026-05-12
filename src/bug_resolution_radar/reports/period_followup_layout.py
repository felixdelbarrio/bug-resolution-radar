"""Layout tokens and guards for the quincenal follow-up PPT renderer."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Iterable

from pptx.util import Pt


@dataclass(frozen=True)
class MetricTypography:
    value_size_pt: float
    label_size_pt: float
    detail_size_pt: float


@dataclass(frozen=True)
class KpiRow:
    value_text: str
    label_text: str


@dataclass(frozen=True)
class KpiSideMetric:
    value_text: str
    label_text: str


@dataclass(frozen=True)
class PeriodFollowupLayoutTheme:
    metric_value_base_pt: float = 25.0
    metric_value_min_pt: float = 17.0
    metric_label_base_pt: float = 12.0
    metric_label_min_pt: float = 8.6
    metric_detail_base_pt: float = 12.0
    metric_detail_min_pt: float = 8.8
    metric_extra_small_pt: float = 8.1
    split_column_ratio: float = 0.636
    split_column_padding_ratio: float = 0.012
    split_column_right_padding_ratio: float = 0.040
    split_divider_top_ratio: float = 0.094
    split_divider_height_ratio: float = 0.811
    split_divider_width_ratio: float = 0.0018
    split_text_top_ratio: float = 0.078
    split_text_height_ratio: float = 0.845
    split_metric_gap_pt: float = 0.85
    delta_badge_width_ratio: float = 0.106
    delta_badge_height_ratio: float = 0.152
    delta_badge_right_gap_ratio: float = 0.014
    delta_badge_min_left_ratio: float = 0.418
    delta_badge_top_ratio: float = 0.088


PERIOD_FOLLOWUP_LAYOUT = PeriodFollowupLayoutTheme()


def metric_card_typography(
    value_text: object,
    label_text: object,
    *,
    theme: PeriodFollowupLayoutTheme = PERIOD_FOLLOWUP_LAYOUT,
) -> MetricTypography:
    value_len = len(str(value_text or "").strip())
    label_len = len(str(label_text or "").strip())
    value_penalty = max(value_len - 2, 0) * 2.2 + max(label_len - 28, 0) * 0.2
    label_penalty = max(label_len - 18, 0) * 0.55 + max(value_len - 2, 0) * 0.25
    value_size = max(theme.metric_value_base_pt - value_penalty, theme.metric_value_min_pt)
    label_size = max(theme.metric_label_base_pt - label_penalty, theme.metric_label_min_pt)
    detail_size = max(
        theme.metric_detail_base_pt - max(label_len - 22, 0) * 0.12, theme.metric_detail_min_pt
    )
    return MetricTypography(
        value_size_pt=round(value_size, 2),
        label_size_pt=round(label_size, 2),
        detail_size_pt=round(detail_size, 2),
    )


def apply_text_frame_margins(text_frame: Any, *, margin_pt: float = 0.0) -> None:
    try:
        margin = Pt(float(margin_pt))
        text_frame.margin_left = margin
        text_frame.margin_right = margin
        text_frame.margin_top = margin
        text_frame.margin_bottom = margin
    except Exception:
        return


def shape_is_inside_slide(shape: Any, *, slide_width: int, slide_height: int) -> bool:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return left >= 0 and top >= 0 and left + width <= slide_width and top + height <= slide_height


def iter_out_of_viewport_shapes(
    slides: Iterable[Any],
    *,
    slide_width: int,
    slide_height: int,
) -> Iterable[Any]:
    for slide in slides:
        for shape in getattr(slide, "shapes", []):
            if not shape_is_inside_slide(shape, slide_width=slide_width, slide_height=slide_height):
                yield shape
