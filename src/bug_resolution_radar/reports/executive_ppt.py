"""Executive PowerPoint generation aligned to dashboard scope, charts and insights."""

from __future__ import annotations

import math
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from typing import Any, Dict, List, Optional, Sequence, Tuple, cast

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from bug_resolution_radar.analysis_window import (
    apply_analysis_depth_filter,
    effective_analysis_lookback_months,
    max_available_backlog_months,
)
from bug_resolution_radar.config import Settings, all_configured_sources
from bug_resolution_radar.design_tokens import (
    BBVA_FONT_HEADLINE_PPT,
    BBVA_FONT_SANS_BOOK_PPT,
    BBVA_FONT_SANS_MEDIUM_PPT,
    BBVA_FONT_SANS_PPT,
    BBVA_LIGHT,
)
from bug_resolution_radar.kpis import compute_kpis
from bug_resolution_radar.status_semantics import effective_closed_mask, is_finalist_status
from bug_resolution_radar.ui.common import load_issues_df, normalize_text_col
from bug_resolution_radar.ui.dashboard.registry import ChartContext, build_trends_registry
from bug_resolution_radar.ui.insights.engine import (
    ActionInsight,
    TrendInsightPack,
    build_trend_insight_pack,
)

SLIDE_WIDTH = int(Inches(13.333))
SLIDE_HEIGHT = int(Inches(7.5))

FONT_HEAD = BBVA_FONT_HEADLINE_PPT
FONT_BODY = BBVA_FONT_SANS_PPT
FONT_BODY_BOOK = BBVA_FONT_SANS_BOOK_PPT
FONT_BODY_MEDIUM = BBVA_FONT_SANS_MEDIUM_PPT

PALETTE: Dict[str, str] = {
    "navy": BBVA_LIGHT.midnight.lstrip("#"),
    "blue": BBVA_LIGHT.core_blue.lstrip("#"),
    "sky": BBVA_LIGHT.serene_dark_blue.lstrip("#"),
    "teal": BBVA_LIGHT.aqua.lstrip("#"),
    "green": "38761D",
    "amber": "F5B942",
    "red": "D64550",
    "ink": BBVA_LIGHT.ink.lstrip("#"),
    "muted": BBVA_LIGHT.ink_muted.lstrip("#"),
    "panel": BBVA_LIGHT.white.lstrip("#"),
    "bg": BBVA_LIGHT.bg_light.lstrip("#"),
    "line": "D3D8E1",
    "mist": "EEF3FB",
}

ORTHO_REPLACEMENTS: Tuple[Tuple[str, str], ...] = (
    (r"\bAccion\b", "Acción"),
    (r"\baccion\b", "acción"),
    (r"\bAnalisis\b", "Análisis"),
    (r"\banalisis\b", "análisis"),
    (r"\bAntiguedad\b", "Antigüedad"),
    (r"\bantiguedad\b", "antigüedad"),
    (r"\bContexto del grafico\b", "Contexto del gráfico"),
    (r"\bcontexto del grafico\b", "contexto del gráfico"),
    (r"\bDias\b", "Días"),
    (r"\bdias\b", "días"),
    (r"\bDia\b", "Día"),
    (r"\bdia\b", "día"),
    (r"\bGrafico\b", "Gráfico"),
    (r"\bgrafico\b", "gráfico"),
    (r"\bMas\b", "Más"),
    (r"\bmas\b", "más"),
    (r"\bResolucion\b", "Resolución"),
    (r"\bresolucion\b", "resolución"),
    (r"\bSenal\b", "Señal"),
    (r"\bsenal\b", "señal"),
    (r"\bSintesis\b", "Síntesis"),
    (r"\bsintesis\b", "síntesis"),
    (r"\bTecnica\b", "Técnica"),
    (r"\btecnica\b", "técnica"),
)

SOFT_TONE_REPLACEMENTS: Tuple[Tuple[str, str], ...] = (
    (r"\bcriticidad\b", "priorización"),
    (r"\bCriticidad\b", "Priorización"),
    (r"\bcr[ií]ticas\b", "Must"),
    (r"\bCr[ií]ticas\b", "Must"),
    (r"\bcr[ií]tica\b", "Must"),
    (r"\bCr[ií]tica\b", "Must"),
    (r"\bcr[ií]tico\b", "Must"),
    (r"\bCr[ií]tico\b", "Must"),
    (r"\burgencia\b", "prioridad de acción"),
    (r"\bUrgencia\b", "Prioridad de acción"),
    (r"\burgente\b", "prioritario"),
    (r"\bUrgente\b", "Prioritario"),
)

BUSINESS_PLAIN_REPLACEMENTS: Tuple[Tuple[str, str], ...] = (
    (r"\bSLA\b", "tiempo objetivo"),
    (r"\bslas\b", "tiempos objetivo"),
    (r"\bP95\b", "tramo del 5% más lento"),
    (r"\bpercentil\s*95\b", "tramo del 5% más lento"),
    (r"\bmediana\b", "valor típico"),
    (r"\bthroughput\b", "ritmo de salida"),
)

CHART_ORDER: Tuple[str, ...] = (
    "open_status_bar",
    "open_priority_pie",
    "age_buckets",
    "timeseries",
    "resolution_hist",
)

CHART_THEMES: Dict[str, str] = {
    "open_status_bar": "Situación actual",
    "open_priority_pie": "Riesgo y exposición",
    "age_buckets": "Riesgo y exposición",
    "timeseries": "Ritmo del flujo",
    "resolution_hist": "Capacidad de cierre",
}

CHART_BUSINESS_SUBTITLES: Dict[str, str] = {
    "open_status_bar": "Dónde se concentra hoy el trabajo pendiente",
    "open_priority_pie": "Qué parte del backlog demanda más atención",
    "age_buckets": "Qué volumen acumula mayor antigüedad",
    "timeseries": "Cómo evoluciona la carga semana a semana",
    "resolution_hist": "Dónde se alarga el cierre de incidencias",
}


@dataclass(frozen=True)
class ExecutiveReportResult:
    """Generated report payload and metadata."""

    file_name: str
    content: bytes
    slide_count: int
    total_issues: int
    open_issues: int
    closed_issues: int
    country: str
    source_id: str
    source_label: str
    applied_filter_summary: str


@dataclass(frozen=True)
class _FilterSnapshot:
    status: Tuple[str, ...]
    priority: Tuple[str, ...]
    assignee: Tuple[str, ...]
    analysis_lookback_months: int = 0
    analysis_available_months: int = 0

    def is_active(self) -> bool:
        return bool(self.status or self.priority or self.assignee)

    def _window_summary(self) -> str:
        lookback = max(0, int(self.analysis_lookback_months or 0))
        available = max(0, int(self.analysis_available_months or 0))
        if lookback <= 0 and available <= 0:
            return ""
        if available <= 0:
            available = max(1, lookback)
        if lookback <= 0:
            lookback = available
        if lookback >= available:
            return f"Ventana temporal: histórico completo ({available} meses)"
        return f"Ventana temporal: últimos {lookback} de {available} meses"

    @staticmethod
    def _values_summary(values: Tuple[str, ...]) -> str:
        if not values:
            return "Todos"
        vals = [str(v or "").strip() for v in values if str(v or "").strip()]
        if len(vals) <= 3:
            return ", ".join(vals)
        return ", ".join(vals[:3]) + f" (+{len(vals) - 3})"

    def summary(self) -> str:
        parts: List[str] = []
        win = self._window_summary()
        if win:
            parts.append(win)
        parts.append(f"Estado: {self._values_summary(self.status)}")
        parts.append(f"Prioridad: {self._values_summary(self.priority)}")
        parts.append(f"Responsable: {self._values_summary(self.assignee)}")
        return " · ".join(parts)


@dataclass(frozen=True)
class _ChartSection:
    chart_id: str
    theme: str
    title: str
    subtitle: str
    figure: Optional[go.Figure]
    insight_pack: TrendInsightPack


@dataclass(frozen=True)
class _ScopeContext:
    country: str
    source_id: str
    source_label: str
    filters: _FilterSnapshot
    dff: pd.DataFrame
    open_df: pd.DataFrame
    closed_df: pd.DataFrame
    generated_at: datetime
    sections: List[_ChartSection]


def _rgb(hex_code: str) -> RGBColor:
    code = str(hex_code or "").strip().lstrip("#")
    if len(code) != 6:
        code = "000000"
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


def _slug(value: str) -> str:
    txt = str(value or "").strip().lower()
    txt = re.sub(r"[^a-z0-9]+", "-", txt).strip("-")
    return txt or "scope"


def _to_dt_naive(series: pd.Series | None) -> pd.Series:
    if series is None:
        return pd.Series([], dtype="datetime64[ns]")
    out = pd.to_datetime(series, errors="coerce")
    try:
        if hasattr(out.dt, "tz") and out.dt.tz is not None:
            out = out.dt.tz_localize(None)
    except Exception:
        try:
            out = out.dt.tz_localize(None)
        except Exception:
            pass
    return out


def _is_finalist_status(value: object) -> bool:
    return is_finalist_status(value)


def _dominant_operational_status(open_df: pd.DataFrame) -> str:
    if open_df is None or open_df.empty or "status" not in open_df.columns:
        return "—"
    statuses = normalize_text_col(open_df["status"], "(sin estado)").astype(str)
    counts = statuses.value_counts()
    for status_name in counts.index.tolist():
        if _is_finalist_status(status_name):
            continue
        return str(status_name)
    return "—"


def _normalize_filter_values(values: Sequence[str] | None) -> Tuple[str, ...]:
    normalized = sorted({str(v or "").strip() for v in list(values or []) if str(v or "").strip()})
    return tuple(normalized)


def _resolve_source_label(settings: Settings, source_id: str) -> str:
    sid = str(source_id or "").strip()
    if not sid:
        return "Origen no especificado"

    for src in all_configured_sources(settings):
        if str(src.get("source_id") or "").strip() != sid:
            continue
        alias = str(src.get("alias") or "").strip() or sid
        source_type = str(src.get("source_type") or "").strip().upper() or "SOURCE"
        country = str(src.get("country") or "").strip()
        return f"{country} · {alias} ({source_type})" if country else f"{alias} ({source_type})"
    return sid


def _scope_df(df: pd.DataFrame, *, country: str, source_id: str) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()

    mask = pd.Series(True, index=df.index)
    country_txt = str(country or "").strip()
    source_txt = str(source_id or "").strip()

    if country_txt and "country" in df.columns:
        mask &= df["country"].fillna("").astype(str).eq(country_txt)
    if source_txt and "source_id" in df.columns:
        mask &= df["source_id"].fillna("").astype(str).eq(source_txt)
    return df.loc[mask].copy(deep=False)


def _apply_report_filters(df: pd.DataFrame, filters: _FilterSnapshot) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()

    mask = pd.Series(True, index=df.index)

    if filters.status and "status" in df.columns:
        status_norm = normalize_text_col(df["status"], "(sin estado)")
        mask &= status_norm.isin(list(filters.status))

    if filters.priority and "priority" in df.columns:
        prio_norm = normalize_text_col(df["priority"], "(sin priority)")
        mask &= prio_norm.isin(list(filters.priority))

    if filters.assignee and "assignee" in df.columns:
        assignee_norm = normalize_text_col(df["assignee"], "(sin asignar)")
        mask &= assignee_norm.isin(list(filters.assignee))

    return df.loc[mask].copy(deep=False)


def _open_closed(df: pd.DataFrame) -> Tuple[pd.DataFrame, pd.DataFrame]:
    if df.empty:
        return pd.DataFrame(), pd.DataFrame()
    closed_mask = effective_closed_mask(df)
    open_df = df.loc[~closed_mask].copy(deep=False)
    closed_df = df.loc[closed_mask].copy(deep=False)
    return open_df, closed_df


def _clip_text(text: str, *, max_len: int = 180) -> str:
    clean = re.sub(r"\s+", " ", str(text or "").strip())
    clean = _fix_spanish_orthography(clean)
    clean = clean.replace("…", " ")
    clean = re.sub(r"\.{3,}", ". ", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    if len(clean) <= max_len:
        return clean
    return clean[:max_len].rstrip()


def _clip_chart_text(text: str, *, max_len: int = 180) -> str:
    clean = _clip_text(text, max_len=max_len)
    if len(clean) < len(re.sub(r"\s+", " ", str(text or "").strip())):
        return clean[: max_len - 1].rstrip() + "…"
    return clean


def _fix_spanish_orthography(text: str) -> str:
    out = str(text or "")
    for pattern, replacement in ORTHO_REPLACEMENTS:
        out = re.sub(pattern, replacement, out)
    return out


def _soften_insight_tone(text: str) -> str:
    out = _fix_spanish_orthography(str(text or ""))
    for pattern, replacement in SOFT_TONE_REPLACEMENTS:
        out = re.sub(pattern, replacement, out)
    for pattern, replacement in BUSINESS_PLAIN_REPLACEMENTS:
        out = re.sub(pattern, replacement, out)
    out = out.replace("…", " ")
    out = re.sub(r"\.{3,}", ". ", out)
    return re.sub(r"\s+", " ", out).strip()


def _fmt_int(value: int) -> str:
    return f"{int(value):,}".replace(",", ".")


def _legend_profile(fig: go.Figure, trace_count: int) -> Tuple[float, int]:
    types = {
        str(getattr(trace, "type", "") or "").strip().lower()
        for trace in list(getattr(fig, "data", []) or [])
    }
    if "pie" in types:
        font_size, max_len = 10.4, 20
    elif "bar" in types:
        font_size, max_len = (10.2, 22) if trace_count <= 5 else (9.4, 18)
    elif "scatter" in types:
        font_size, max_len = 11.2, 24
    else:
        font_size, max_len = 10.4, 21

    if trace_count >= 8:
        font_size -= 0.8
        max_len -= 2
    elif trace_count >= 6:
        font_size -= 0.5
        max_len -= 1
    return max(8.0, font_size), max(12, max_len)


def _estimate_legend_rows(names: Sequence[str], *, font_size: float, width_px: int = 1600) -> int:
    usable = max(900, int(width_px - 72))
    rows = 1
    cursor = 0.0
    for name in [str(n or "").strip() for n in list(names or []) if str(n or "").strip()]:
        # Approx label width = marker + spacing + text width.
        item_w = 36.0 + (len(name) * font_size * 0.62)
        if cursor and (cursor + item_w) > usable:
            rows += 1
            cursor = item_w
        else:
            cursor += item_w
    return max(1, rows)


def _set_rich_text(
    paragraph: Any,
    text: str,
    *,
    font_name: str,
    font_size: float,
    color_hex: str,
    default_bold: bool = False,
) -> None:
    paragraph.clear()
    normalized = _fix_spanish_orthography(str(text or ""))
    chunks = re.split(r"(\*\*.*?\*\*)", normalized)
    for chunk in chunks:
        if not chunk:
            continue
        is_bold = chunk.startswith("**") and chunk.endswith("**") and len(chunk) >= 4
        txt = chunk[2:-2] if is_bold else chunk
        if not txt:
            continue
        run = paragraph.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bool(default_bold or is_bold)
        run.font.color.rgb = _rgb(color_hex)


def _tone_style(tone: str) -> Tuple[str, str, str]:
    key = str(tone or "").strip().lower()
    styles: Dict[str, Tuple[str, str, str]] = {
        "blue": ("EAF2FF", "B8CCE8", "0B3A75"),
        "sky": ("E8F7FF", "9DDCFB", "0B4A6F"),
        "teal": ("E6F9F7", "9EDFD9", "0E5C5C"),
        "amber": ("FFF4DE", "F3D89B", "7A5A12"),
        "green": ("EAF6EC", "B8DDBF", "1F5B2E"),
        "red": ("FDEBEC", "E3A5AA", "8B1D26"),
        "urgency_must": ("FDEBEC", "E3A5AA", "8B1D26"),
        "urgency_should": ("FFF4DE", "F3D89B", "7A5A12"),
        "urgency_nice": ("EAF6EC", "B8DDBF", "1F5B2E"),
    }
    return styles.get(key, ("EEF3FB", "C8D6E8", PALETTE["navy"]))


def _urgency_from_score(score: float) -> Tuple[str, str]:
    val = float(score or 0.0)
    if val >= 18.0:
        return "Must", "urgency_must"
    if val >= 9.0:
        return "Should", "urgency_should"
    return "Nice to have", "urgency_nice"


def _add_mini_kpi_ribbons(
    slide: Any,
    *,
    items: Sequence[Tuple[str, str, str]],
    y: float,
    x: float = 0.62,
    total_width: float = 12.0,
    height: float = 0.34,
    gap: float = 0.12,
) -> None:
    normalized: List[Tuple[str, str, str]] = [
        (str(lbl or "").strip(), str(val or "").strip(), str(tone or "").strip())
        for lbl, val, tone in items
        if str(lbl or "").strip() and str(val or "").strip()
    ]
    if not normalized:
        return
    n = len(normalized)
    if n == 1:
        width = total_width
    else:
        width = max(1.2, (total_width - (gap * (n - 1))) / float(n))

    for idx, (label, value, tone) in enumerate(normalized):
        bg, border, txt = _tone_style(tone)
        xpos = x + idx * (width + gap)
        ribbon = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(xpos),
            Inches(y),
            Inches(width),
            Inches(height),
        )
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = _rgb(bg)
        ribbon.line.color.rgb = _rgb(border)

        tf = ribbon.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_rich_text(
            p,
            f"{_clip_text(label, max_len=26)}: **{_clip_text(value, max_len=24)}**",
            font_name=FONT_BODY_MEDIUM,
            font_size=10.8,
            color_hex=txt,
        )


def _build_sections(
    settings: Settings, *, dff: pd.DataFrame, open_df: pd.DataFrame
) -> List[_ChartSection]:
    registry = build_trends_registry()
    kpis = compute_kpis(dff, settings=settings)
    chart_ctx = ChartContext(dff=dff, open_df=open_df, kpis=kpis)

    sections: List[_ChartSection] = []
    for chart_id in CHART_ORDER:
        spec = registry.get(chart_id)
        if spec is None:
            continue
        insight_pack = build_trend_insight_pack(chart_id, dff=dff, open_df=open_df)
        figure = spec.render(chart_ctx)
        if figure is None:
            # Skip chart section when scope/filter has insufficient data.
            continue
        sections.append(
            _ChartSection(
                chart_id=chart_id,
                theme=CHART_THEMES.get(chart_id, "Señales operativas"),
                title=spec.title,
                subtitle=spec.subtitle,
                figure=figure,
                insight_pack=insight_pack,
            )
        )
    return sections


def _build_context(
    settings: Settings,
    *,
    country: str,
    source_id: str,
    filters: _FilterSnapshot,
    dff_override: pd.DataFrame | None,
    open_df_override: pd.DataFrame | None,
    scoped_source_df: pd.DataFrame | None,
) -> _ScopeContext:
    if dff_override is None:
        if scoped_source_df is None:
            base_df = _scope_df(
                load_issues_df(settings.DATA_PATH), country=country, source_id=source_id
            )
        else:
            base_df = scoped_source_df.copy(deep=False)
        base_df = apply_analysis_depth_filter(base_df, settings=settings)
        dff = _apply_report_filters(base_df, filters)
        open_df, closed_df = _open_closed(dff)
    else:
        dff = apply_analysis_depth_filter(dff_override.copy(deep=False), settings=settings)
        if open_df_override is None:
            open_df, closed_df = _open_closed(dff)
        else:
            open_df = apply_analysis_depth_filter(
                open_df_override.copy(deep=False), settings=settings
            )
            if not dff.empty and not open_df.empty:
                open_index = set(open_df.index.tolist())
                closed_df = dff.loc[[idx for idx in dff.index if idx not in open_index]].copy(
                    deep=False
                )
            else:
                closed_df = pd.DataFrame()

    if dff.empty:
        raise ValueError(
            "No hay incidencias para el país/origen y filtros seleccionados. Ajusta scope o filtros."
        )

    sections = _build_sections(settings, dff=dff, open_df=open_df)
    return _ScopeContext(
        country=str(country or "").strip(),
        source_id=str(source_id or "").strip(),
        source_label=_resolve_source_label(settings, source_id),
        filters=filters,
        dff=dff,
        open_df=open_df,
        closed_df=closed_df,
        generated_at=datetime.now(timezone.utc),
        sections=sections,
    )


def _fig_to_png(fig: Optional[go.Figure]) -> Optional[bytes]:
    if fig is None:
        return None
    export_width = 1280
    export_height = 820
    try:

        def _trace_values(trace: object, attr: str) -> List[object]:
            raw = getattr(trace, attr, None)
            if raw is None:
                return []
            try:
                return list(raw)
            except Exception:
                return [raw]

        def _fmt_total(value: float) -> str:
            if abs(value - round(value)) < 1e-9:
                return f"{int(round(value)):,}".replace(",", ".")
            txt = f"{value:,.1f}"
            return txt.replace(",", "X").replace(".", ",").replace("X", ".")

        export_fig = go.Figure(fig)
        traces = list(getattr(export_fig, "data", []) or [])
        bar_category_keys: set[str] = set()
        for trace in traces:
            trace_type = str(getattr(trace, "type", "") or "").strip().lower()
            is_vertical_bar = (
                trace_type == "bar"
                and str(getattr(trace, "orientation", "v") or "v").strip().lower() != "h"
            )
            if not is_vertical_bar:
                continue
            for raw_x in _trace_values(trace, "x"):
                key = str(raw_x or "").strip()
                if key:
                    bar_category_keys.add(key)

        named_traces = [t for t in traces if str(getattr(t, "name", "") or "").strip()]
        pie_labels: List[str] = []
        for trace in traces:
            trace_type = str(getattr(trace, "type", "") or "").strip().lower()
            if trace_type != "pie":
                continue
            for raw_label in _trace_values(trace, "labels"):
                label = str(raw_label or "").strip()
                if label:
                    pie_labels.append(label)

        showlegend = bool(getattr(getattr(export_fig, "layout", None), "showlegend", False))
        if not showlegend and (len(named_traces) > 1 or len(set(pie_labels)) > 1):
            showlegend = True

        legend_ref_size = max(len(named_traces), len(set(pie_labels)))
        legend_font, max_name_len = _legend_profile(export_fig, legend_ref_size)
        legend_render_font = max(14.2, legend_font + 2.4)
        legend_item_count = 0
        if showlegend:
            legend_names: List[str] = []
            for trace in traces:
                trace_type = str(getattr(trace, "type", "") or "").strip().lower()
                if trace_type == "pie":
                    labels = [
                        _clip_chart_text(str(lbl or "").strip(), max_len=max_name_len)
                        for lbl in _trace_values(trace, "labels")
                    ]
                    if labels:
                        trace.labels = labels
                        legend_names.extend(labels)
                    continue
                name = str(getattr(trace, "name", "") or "").strip()
                if not name:
                    continue
                short_name = _clip_chart_text(name, max_len=max_name_len)
                trace.name = short_name
                legend_names.append(short_name)

            # Preserve order while removing duplicates (common in pie labels).
            seen_legend: set[str] = set()
            unique_legend_names: List[str] = []
            for name in legend_names:
                if not name or name in seen_legend:
                    continue
                seen_legend.add(name)
                unique_legend_names.append(name)
            legend_names = unique_legend_names
            legend_item_count = len(legend_names)
            rows = _estimate_legend_rows(
                legend_names, font_size=legend_render_font, width_px=export_width
            )
            while rows > 2 and legend_render_font > 10.0:
                legend_render_font = max(10.0, legend_render_font - 0.5)
                rows = _estimate_legend_rows(
                    legend_names, font_size=legend_render_font, width_px=export_width
                )
        else:
            rows = 1

        dense_chart = max(len(bar_category_keys), legend_item_count, len(set(pie_labels))) >= 8
        if dense_chart:
            legend_render_font = min(17.4, legend_render_font + 1.5)

        export_fig.update_layout(
            template="plotly_white",
            showlegend=showlegend,
            uniformtext=dict(minsize=18 if dense_chart else 16, mode="hide"),
            font=dict(
                family=FONT_BODY_BOOK,
                size=18,
                color=f"#{PALETTE['ink']}",
            ),
            paper_bgcolor="#FFFFFF",
            plot_bgcolor="#FFFFFF",
            legend=dict(
                orientation="h",
                yanchor="top",
                y=-(0.16 + (max(1, rows) - 1) * 0.10),
                xanchor="left",
                x=0.0,
                font=dict(
                    family=FONT_BODY_MEDIUM,
                    size=legend_render_font,
                    color=f"#{PALETTE['ink']}",
                ),
                bgcolor="rgba(255,255,255,0.98)",
                bordercolor=f"#{PALETTE['line']}",
                borderwidth=1,
                title=dict(text=""),
            ),
            margin=dict(
                l=max(24, int(getattr(getattr(export_fig.layout, "margin", None), "l", 24) or 24)),
                r=max(24, int(getattr(getattr(export_fig.layout, "margin", None), "r", 24) or 24)),
                t=max(52, int(getattr(getattr(export_fig.layout, "margin", None), "t", 52) or 52)),
                b=max(
                    (106 + (max(1, rows) * 28)) if showlegend else 24,
                    int(
                        getattr(
                            getattr(export_fig.layout, "margin", None),
                            "b",
                            (106 + (max(1, rows) * 28)) if showlegend else 24,
                        )
                        or ((106 + (max(1, rows) * 28)) if showlegend else 24)
                    ),
                ),
            ),
        )
        # Force light-mode readability for exported charts regardless of UI theme.
        export_fig.update_annotations(font=dict(family=FONT_BODY_BOOK, color=f"#{PALETTE['ink']}"))

        # For every column chart in the report, show a clear total per column
        # in bold on top and avoid truncated per-segment labels.
        bar_traces = [
            tr
            for tr in traces
            if str(getattr(tr, "type", "") or "").strip().lower() == "bar"
            and str(getattr(tr, "orientation", "v") or "v").strip().lower() != "h"
        ]
        if bar_traces:
            totals: Dict[str, float] = {}
            x_by_key: Dict[str, object] = {}
            x_order: List[str] = []
            for trace in bar_traces:
                xs = _trace_values(trace, "x")
                ys = _trace_values(trace, "y")
                points = min(len(xs), len(ys))
                for idx in range(points):
                    x_raw = xs[idx]
                    y_raw = ys[idx]
                    if y_raw is None:
                        continue
                    if isinstance(y_raw, (int, float)):
                        y_val = float(y_raw)
                    else:
                        y_txt = str(y_raw).strip()
                        if not y_txt:
                            continue
                        try:
                            y_val = float(y_txt)
                        except ValueError:
                            continue
                    if pd.isna(y_val):
                        continue
                    key = str(x_raw)
                    if key not in totals:
                        totals[key] = 0.0
                        x_by_key[key] = x_raw
                        x_order.append(key)
                    totals[key] += max(0.0, y_val)

            if totals:
                for trace in bar_traces:
                    if hasattr(trace, "text"):
                        trace.text = None
                    if hasattr(trace, "texttemplate"):
                        trace.texttemplate = None
                    if hasattr(trace, "textposition"):
                        trace.textposition = "none"
                    if hasattr(trace, "cliponaxis"):
                        trace.cliponaxis = False

                ordered_totals = [float(totals[key]) for key in x_order]
                max_total = max(ordered_totals) if ordered_totals else 0.0
                if max_total > 0.0:
                    top_pad = max(1.0, max_total * 0.10)
                    label_y = [val + (top_pad * 0.38) for val in ordered_totals]
                    export_fig.add_trace(
                        go.Scatter(
                            x=[x_by_key[key] for key in x_order],
                            y=label_y,
                            mode="text",
                            text=[f"<b>{_fmt_total(val)}</b>" for val in ordered_totals],
                            textposition="top center",
                            textfont=dict(
                                size=28 if dense_chart else 24, color=f"#{PALETTE['ink']}"
                            ),
                            showlegend=False,
                            hoverinfo="skip",
                        )
                    )

        axis_tick_size = 24 if dense_chart else 22
        axis_title_size = 26 if dense_chart else 24
        export_fig.update_xaxes(
            tickfont=dict(family=FONT_BODY_BOOK, size=axis_tick_size, color=f"#{PALETTE['ink']}"),
            title_font=dict(
                family=FONT_BODY_MEDIUM, size=axis_title_size, color=f"#{PALETTE['ink']}"
            ),
            color=f"#{PALETTE['ink']}",
            showline=True,
            linecolor=f"#{PALETTE['line']}",
            gridcolor="#E8EDF4",
            zerolinecolor="#E8EDF4",
        )
        export_fig.update_yaxes(
            tickfont=dict(family=FONT_BODY_BOOK, size=axis_tick_size, color=f"#{PALETTE['ink']}"),
            title_font=dict(
                family=FONT_BODY_MEDIUM, size=axis_title_size, color=f"#{PALETTE['ink']}"
            ),
            color=f"#{PALETTE['ink']}",
            showline=True,
            linecolor=f"#{PALETTE['line']}",
            gridcolor="#E8EDF4",
            zerolinecolor="#E8EDF4",
        )
        for trace in traces:
            trace_type = str(getattr(trace, "type", "") or "").strip().lower()
            if trace_type == "pie":
                trace.showlegend = bool(showlegend)
                if hasattr(trace, "textposition"):
                    trace.textposition = "inside"
                if hasattr(trace, "insidetextorientation"):
                    trace.insidetextorientation = "horizontal"
                if hasattr(trace, "textinfo"):
                    trace.textinfo = "percent"
                if hasattr(trace, "textfont"):
                    trace.textfont = dict(
                        family=FONT_BODY_MEDIUM,
                        size=26 if dense_chart else 24,
                        color=f"#{PALETTE['ink']}",
                    )
                continue

            if str(getattr(trace, "name", "") or "").strip():
                trace.showlegend = bool(showlegend)
            if hasattr(trace, "textfont"):
                trace.textfont = dict(
                    family=FONT_BODY_MEDIUM,
                    size=24 if dense_chart else 22,
                    color=f"#{PALETTE['ink']}",
                )
            if trace_type == "bar" and hasattr(trace, "textposition"):
                trace.textposition = "auto"
                if hasattr(trace, "cliponaxis"):
                    trace.cliponaxis = False

        image = pio.to_image(
            export_fig, format="png", width=export_width, height=export_height, scale=3
        )
        return cast(bytes, image)
    except Exception:
        try:
            image = pio.to_image(
                fig, format="png", width=export_width, height=export_height, scale=2
            )
            return cast(bytes, image)
        except Exception:
            return None


def _add_bg(slide: Any, color_hex: str) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(color_hex)
    rect.line.fill.background()


def _add_header(slide: Any, *, title: str, subtitle: str, dark: bool = False) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        Inches(0.92),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(PALETTE["navy"] if dark else PALETTE["mist"])
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.14), Inches(11.9), Inches(0.30))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.bold = True
    run.font.size = Pt(23 if dark else 20)
    run.font.color.rgb = _rgb("FFFFFF" if dark else PALETTE["navy"])

    subtitle_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.50), Inches(11.9), Inches(0.20))
    tf2 = subtitle_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_BODY_BOOK
    run2.font.size = Pt(11.2)
    run2.font.color.rgb = _rgb("CFE2FF" if dark else PALETTE["muted"])


def _add_footer(slide: Any, *, context: _ScopeContext) -> None:
    foot = slide.shapes.add_textbox(Inches(0.52), Inches(7.16), Inches(12.2), Inches(0.20))
    tf = foot.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (
        f"{context.country} · {context.source_label}"
        f"  |  Filtros: {context.filters.summary()}"
        f"  |  {context.generated_at.strftime('%Y-%m-%d %H:%M UTC')}"
    )
    run.font.name = FONT_BODY_BOOK
    run.font.size = Pt(8.8)
    run.font.color.rgb = _rgb(PALETTE["muted"])


def _metric_card(
    slide: Any, *, x: float, y: float, w: float, h: float, label: str, value: str, tone: str
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    card.line.color.rgb = _rgb(PALETTE["line"])

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x + 0.14),
        Inches(y + 0.15),
        Inches(0.18),
        Inches(0.18),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(PALETTE.get(tone, PALETTE["blue"]))
    badge.line.fill.background()

    lbl = slide.shapes.add_textbox(
        Inches(x + 0.38), Inches(y + 0.12), Inches(w - 0.48), Inches(0.22)
    )
    tf = lbl.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = FONT_BODY_MEDIUM
    run.font.size = Pt(10.8)
    run.font.color.rgb = _rgb(PALETTE["muted"])

    val = slide.shapes.add_textbox(
        Inches(x + 0.14), Inches(y + 0.44), Inches(w - 0.24), Inches(0.48)
    )
    tf2 = val.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = value
    run2.font.name = FONT_HEAD
    run2.font.bold = True
    run2.font.size = Pt(23)
    run2.font.color.rgb = _rgb(PALETTE["ink"])


def _add_cover_slide(prs: Any, context: _ScopeContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, "001B4A")
    _add_header(
        slide,
        title="Radar de Resolución de Incidencias",
        subtitle="Lectura de negocio para priorizar decisiones y acelerar resultados",
        dark=True,
    )

    date_box = slide.shapes.add_textbox(Inches(10.42), Inches(1.12), Inches(2.20), Inches(0.34))
    dtf = date_box.text_frame
    dtf.clear()
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.RIGHT
    dr = dp.add_run()
    dr.text = context.generated_at.strftime("%B %Y").capitalize()
    dr.font.name = FONT_BODY_BOOK
    dr.font.size = Pt(11)
    dr.font.color.rgb = _rgb("CFE2FF")

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.58),
        Inches(1.34),
        Inches(4.6),
        Inches(0.08),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(PALETTE["teal"])
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.60), Inches(1.58), Inches(11.9), Inches(1.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{context.country}\n{context.source_label}"
    run.font.name = FONT_HEAD
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = _rgb("FFFFFF")

    _metric_card(
        slide,
        x=0.62,
        y=4.18,
        w=3.1,
        h=1.25,
        label="Casos analizados",
        value=_fmt_int(len(context.dff)),
        tone="sky",
    )
    _metric_card(
        slide,
        x=3.96,
        y=4.18,
        w=3.1,
        h=1.25,
        label="Abiertas hoy",
        value=_fmt_int(len(context.open_df)),
        tone="amber",
    )
    _metric_card(
        slide,
        x=7.30,
        y=4.18,
        w=3.1,
        h=1.25,
        label="Resueltas",
        value=_fmt_int(len(context.closed_df)),
        tone="green",
    )

    scope_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(5.64),
        Inches(12.00),
        Inches(1.26),
    )
    scope_box.fill.solid()
    scope_box.fill.fore_color.rgb = _rgb("0A2E67")
    scope_box.line.color.rgb = _rgb("2A66B8")
    tf2 = scope_box.text_frame
    tf2.clear()
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = f"Scope: {context.country} · {context.source_label}"
    run2.font.name = FONT_BODY_MEDIUM
    run2.font.size = Pt(12.5)
    run2.font.color.rgb = _rgb("FFFFFF")

    p3 = tf2.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = f"Filtros aplicados: {context.filters.summary()}"
    run3.font.name = FONT_BODY_BOOK
    run3.font.size = Pt(11.2)
    run3.font.color.rgb = _rgb("DDEBFF")

    foot = slide.shapes.add_textbox(Inches(0.62), Inches(6.97), Inches(12.0), Inches(0.20))
    ftf = foot.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "Objetivo del informe: concentrar esfuerzos donde el impacto en servicio y productividad es mayor."
    fr.font.name = FONT_BODY_BOOK
    fr.font.size = Pt(10)
    fr.font.color.rgb = _rgb("BDD8FF")


def _group_sections_by_theme(
    sections: Sequence[_ChartSection],
) -> List[Tuple[str, List[_ChartSection]]]:
    grouped: Dict[str, List[_ChartSection]] = {}
    order: List[str] = []
    for section in sections:
        theme = str(section.theme or "Señales operativas")
        if theme not in grouped:
            grouped[theme] = []
            order.append(theme)
        grouped[theme].append(section)
    return [(theme, grouped[theme]) for theme in order]


def _section_ribbon_items(section: _ChartSection) -> List[Tuple[str, str, str]]:
    cards = sorted(section.insight_pack.cards, key=lambda c: float(c.score), reverse=True)
    top_score = float(cards[0].score) if cards else 0.0
    urgency_label, urgency_tone = _urgency_from_score(top_score)

    metrics = [
        m
        for m in list(section.insight_pack.metrics or [])
        if str(m.label or "").strip() and str(m.value or "").strip() not in {"", "—"}
    ]
    items: List[Tuple[str, str, str]] = []
    if metrics:
        items.append(
            (
                _clip_text(metrics[0].label, max_len=24),
                _clip_text(metrics[0].value, max_len=20),
                "blue",
            )
        )
    if len(metrics) >= 2:
        items.append(
            (
                _clip_text(metrics[1].label, max_len=24),
                _clip_text(metrics[1].value, max_len=20),
                "sky",
            )
        )
    if cards:
        items.append(
            (
                "Acción foco",
                _clip_text(_soften_insight_tone(str(cards[0].title or "").strip()), max_len=30),
                "teal",
            )
        )
    items.append(("Prioridad de acción", urgency_label, urgency_tone))
    return items[:4]


def _add_index_row(
    tf: Any,
    *,
    code: str,
    title: str,
    detail: str,
    code_color: str,
) -> None:
    p_code = tf.add_paragraph()
    p_code.space_before = Pt(7)
    p_code.space_after = Pt(0)
    r_code = p_code.add_run()
    r_code.text = code
    r_code.font.name = FONT_HEAD
    r_code.font.bold = True
    r_code.font.size = Pt(16)
    r_code.font.color.rgb = _rgb(code_color)

    p_title = tf.add_paragraph()
    p_title.space_before = Pt(0)
    p_title.space_after = Pt(0)
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = FONT_HEAD
    r_title.font.bold = True
    r_title.font.size = Pt(14)
    r_title.font.color.rgb = _rgb(PALETTE["ink"])

    p_detail = tf.add_paragraph()
    p_detail.space_before = Pt(0)
    p_detail.space_after = Pt(2)
    r_detail = p_detail.add_run()
    r_detail.text = detail
    r_detail.font.name = FONT_BODY_BOOK
    r_detail.font.size = Pt(10.5)
    r_detail.font.color.rgb = _rgb(PALETTE["muted"])


def _add_index_slide(prs: Any, context: _ScopeContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE["bg"])
    _add_header(
        slide,
        title="Índice",
        subtitle="Historia del análisis y decisiones sugeridas",
    )
    rows: List[Tuple[str, str, str, str]] = [
        ("01", "Apertura", "Alcance y reglas de lectura del análisis", PALETTE["blue"]),
        ("02", "Panorama", "Tamaño del problema, foco y prioridades", PALETTE["teal"]),
    ]
    idx = 3
    for theme, themed_sections in _group_sections_by_theme(context.sections):
        detail = " · ".join([_clip_text(sec.title, max_len=45) for sec in themed_sections])
        rows.append((f"{idx:02d}", theme, detail, PALETTE["navy"]))
        idx += 1
    rows.append(
        (
            f"{idx:02d}",
            "Cierre y acciones",
            "Plan de trabajo priorizado por impacto",
            PALETTE["blue"],
        )
    )

    body = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(1.30),
        Inches(12.00),
        Inches(5.65),
    )
    body.fill.solid()
    body.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    body.line.color.rgb = _rgb(PALETTE["line"])

    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.paragraphs[0].clear()
    for row in rows:
        _add_index_row(
            tf,
            code=row[0],
            title=row[1],
            detail=row[2],
            code_color=row[3],
        )

    _add_footer(slide, context=context)


def _dominant_value(series: pd.Series, empty: str) -> str:
    if series.empty:
        return empty
    value = str(series.value_counts().index[0]).strip()
    return value or empty


def _add_exec_summary_slide(prs: Any, context: _ScopeContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE["bg"])
    _add_header(
        slide,
        title="Panorama de negocio",
        subtitle="Qué está pasando, por qué importa y dónde actuar primero",
    )
    total = int(len(context.dff))
    open_count = int(len(context.open_df))
    open_ratio = (open_count / total * 100.0) if total else 0.0

    top_status = _dominant_operational_status(context.open_df)
    if top_status == "—":
        top_status = "Sin foco operativo"

    top_priority = "-"
    if not context.open_df.empty and "priority" in context.open_df.columns:
        top_priority = _dominant_value(
            normalize_text_col(context.open_df["priority"], "(sin priority)").astype(str),
            "-",
        )

    themes = [theme for theme, _ in _group_sections_by_theme(context.sections)]
    _metric_card(
        slide,
        x=0.70,
        y=1.30,
        w=2.82,
        h=1.16,
        label="Fase operativa foco",
        value=top_status,
        tone="blue",
    )
    _metric_card(
        slide,
        x=3.64,
        y=1.30,
        w=2.82,
        h=1.16,
        label="Prioridad dominante",
        value=top_priority,
        tone="amber",
    )
    _metric_card(
        slide,
        x=6.58,
        y=1.30,
        w=2.82,
        h=1.16,
        label="Bloques de lectura",
        value=f"{len(themes)}",
        tone="green",
    )
    _metric_card(
        slide,
        x=9.52,
        y=1.30,
        w=2.82,
        h=1.16,
        label="% abiertas",
        value=f"{open_ratio:.1f}%",
        tone="teal",
    )

    left_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(2.70),
        Inches(7.20),
        Inches(4.24),
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    left_box.line.color.rgb = _rgb(PALETTE["line"])

    ltf = left_box.text_frame
    ltf.clear()
    lp0 = ltf.paragraphs[0]
    lr0 = lp0.add_run()
    lr0.text = "Ruta de lectura"
    lr0.font.name = FONT_HEAD
    lr0.font.bold = True
    lr0.font.size = Pt(20)
    lr0.font.color.rgb = _rgb(PALETTE["navy"])

    for idx, theme in enumerate(themes, start=1):
        p = ltf.add_paragraph()
        p.text = f"{idx}. {theme}"
        p.space_after = Pt(7)
        p.font.name = FONT_BODY_BOOK
        p.font.size = Pt(13.0)
        p.font.color.rgb = _rgb(PALETTE["ink"])

    right_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.04),
        Inches(2.70),
        Inches(4.66),
        Inches(4.24),
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    right_box.line.color.rgb = _rgb(PALETTE["line"])

    rtf = right_box.text_frame
    rtf.clear()
    rp0 = rtf.paragraphs[0]
    rr0 = rp0.add_run()
    rr0.text = "Lectura inicial"
    rr0.font.name = FONT_HEAD
    rr0.font.bold = True
    rr0.font.size = Pt(20)
    rr0.font.color.rgb = _rgb(PALETTE["navy"])

    status_line = f"**Fase operativa foco (abiertas):** {top_status}"

    opening = [
        f"**Alcance activo:** {context.country} · {context.source_label}",
        f"**Filtros:** {context.filters.summary()}",
        status_line,
        f"**Prioridad dominante (abiertas):** {top_priority}",
    ]
    for line in opening:
        p = rtf.add_paragraph()
        p.space_after = Pt(5)
        _set_rich_text(
            p,
            line,
            font_name=FONT_BODY_BOOK,
            font_size=12.0,
            color_hex=PALETTE["ink"],
        )

    _add_footer(slide, context=context)


def _top_cards(pack: TrendInsightPack, *, limit: int = 3) -> List[ActionInsight]:
    cards = sorted(pack.cards, key=lambda c: float(c.score), reverse=True)
    return cards[:limit]


def _scaled_font(base: float, *, scale: float, min_size: float) -> float:
    return max(min_size, float(base) * max(0.70, min(1.0, float(scale))))


def _chart_insights_font_scale(
    *,
    tip_text: str,
    context_text: str,
    to_be_context_text: str,
    cards: Sequence[ActionInsight],
) -> float:
    # Side panel available text height (~5.58 in - paddings) in points.
    max_height_pt = 355.0
    estimate = 20.0
    estimate += 30.0  # As-is heading
    estimate += (
        _estimate_wrapped_lines(f"Mensaje clave: {tip_text}", max_chars_per_line=43) * 16.6 + 10.0
    )
    estimate += _estimate_wrapped_lines(context_text, max_chars_per_line=47) * 14.6 + 8.0
    estimate += 30.0  # To-Be heading
    estimate += _estimate_wrapped_lines(to_be_context_text, max_chars_per_line=47) * 14.0 + 8.0

    for idx_card, card in enumerate(list(cards or []), start=1):
        urgency_label, _ = _urgency_from_score(float(card.score))
        title_line = (
            f"{idx_card}. [{urgency_label.upper()}] "
            f"{_clip_text(_soften_insight_tone(card.title), max_len=68)}"
        )
        body_line = _clip_text(
            _soften_insight_tone(f"Recomendación de acción: {str(card.body or '').strip()}"),
            max_len=178,
        )
        estimate += _estimate_wrapped_lines(title_line, max_chars_per_line=44) * 15.2 + 4.0
        estimate += _estimate_wrapped_lines(body_line, max_chars_per_line=49) * 13.8 + 6.0

    if estimate <= max_height_pt:
        return 1.0
    return max(0.72, max_height_pt / estimate)


def _add_chart_insight_slide(
    prs: Any,
    *,
    context: _ScopeContext,
    section: _ChartSection,
    index: int,
    total: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE["bg"])
    _add_header(
        slide,
        title=f"{index:02d}/{total:02d} · {section.title}",
        subtitle=f"{section.theme} · {CHART_BUSINESS_SUBTITLES.get(section.chart_id, section.subtitle)}",
    )
    _add_mini_kpi_ribbons(
        slide,
        items=_section_ribbon_items(section),
        y=1.00,
        x=2.92,
        total_width=9.70,
        height=0.30,
    )

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.56),
        Inches(1.02),
        Inches(2.20),
        Inches(0.30),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = _rgb("EAF2FF")
    chip.line.color.rgb = _rgb("B8CCE8")
    ctf = chip.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = section.chart_id.replace("_", " ").upper()
    cr.font.name = FONT_BODY_MEDIUM
    cr.font.bold = True
    cr.font.size = Pt(8.5)
    cr.font.color.rgb = _rgb(PALETTE["navy"])

    chart_frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.52),
        Inches(1.40),
        Inches(8.16),
        Inches(5.58),
    )
    chart_frame.fill.solid()
    chart_frame.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    chart_frame.line.color.rgb = _rgb(PALETTE["line"])

    image = _fig_to_png(section.figure)
    if image is None:
        tf = chart_frame.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if section.figure is None:
            run.text = "No hay datos suficientes para este gráfico con el filtro actual."
        else:
            run.text = "No se pudo exportar este gráfico para la presentación."
        run.font.name = FONT_BODY_BOOK
        run.font.size = Pt(13)
        run.font.color.rgb = _rgb(PALETTE["muted"])
    else:
        slide.shapes.add_picture(
            BytesIO(image),
            Inches(0.70),
            Inches(1.58),
            width=Inches(7.56),
            height=Inches(5.20),
        )

    side_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.78),
        Inches(1.40),
        Inches(3.98),
        Inches(5.58),
    )
    side_panel.fill.solid()
    side_panel.fill.fore_color.rgb = _rgb(PALETTE["panel"])
    side_panel.line.color.rgb = _rgb(PALETTE["line"])

    itf = side_panel.text_frame
    itf.clear()
    itf.word_wrap = True

    metric_items = [
        m
        for m in list(section.insight_pack.metrics or [])
        if str(getattr(m, "label", "") or "").strip()
        and str(getattr(m, "value", "") or "").strip() not in {"", "—"}
    ]
    cards = _top_cards(section.insight_pack, limit=2)

    p1 = itf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "As-is"
    p1.alignment = PP_ALIGN.CENTER
    run1.font.name = FONT_HEAD
    run1.font.bold = True
    run1.font.color.rgb = _rgb(PALETTE["navy"])

    tip = _clip_text(
        _soften_insight_tone(
            section.insight_pack.executive_tip or "Sin señal concluyente con el filtro actual."
        ),
        max_len=170,
    )

    if len(metric_items) >= 2:
        to_be_ctx_text = _clip_text(
            f"Dato de contexto: {metric_items[1].label}: {metric_items[1].value}.",
            max_len=120,
        )
    else:
        to_be_ctx_text = f"Dato de contexto: {len(cards)} acciones priorizadas para esta vista."

    context_detail = _clip_text(
        f"{CHART_BUSINESS_SUBTITLES.get(section.chart_id, section.subtitle)}. "
        + (
            f"{metric_items[0].label}: {metric_items[0].value}."
            if metric_items
            else "Las cifras más relevantes están en la banda KPI superior para facilitar la lectura."
        ),
        max_len=188,
    )
    panel_scale = _chart_insights_font_scale(
        tip_text=tip,
        context_text=context_detail,
        to_be_context_text=to_be_ctx_text,
        cards=cards,
    )
    run1.font.size = Pt(_scaled_font(24.0, scale=panel_scale, min_size=18.0))

    p_tip = itf.add_paragraph()
    p_tip.space_before = Pt(2)
    p_tip.space_after = Pt(8)
    _set_rich_text(
        p_tip,
        f"**Mensaje clave:** {tip}",
        font_name=FONT_BODY_BOOK,
        font_size=_scaled_font(14.4, scale=panel_scale, min_size=11.0),
        color_hex=PALETTE["ink"],
    )

    p_context = itf.add_paragraph()
    p_context.space_before = Pt(2)
    p_context.space_after = Pt(6)
    _set_rich_text(
        p_context,
        f"**Por qué importa:** {context_detail}",
        font_name=FONT_BODY_BOOK,
        font_size=_scaled_font(12.8, scale=panel_scale, min_size=10.0),
        color_hex=PALETTE["muted"],
    )

    p_actions = itf.add_paragraph()
    p_actions.text = "To-Be"
    p_actions.alignment = PP_ALIGN.CENTER
    p_actions.space_before = Pt(4)
    p_actions.space_after = Pt(2)
    p_actions.font.name = FONT_HEAD
    p_actions.font.bold = True
    p_actions.font.size = Pt(_scaled_font(24.0, scale=panel_scale, min_size=18.0))
    p_actions.font.color.rgb = _rgb(PALETTE["navy"])

    p_to_be_ctx = itf.add_paragraph()
    p_to_be_ctx.space_before = Pt(0)
    p_to_be_ctx.space_after = Pt(4)
    if len(metric_items) >= 2:
        _set_rich_text(
            p_to_be_ctx,
            f"**Dato de contexto:** {_clip_text(to_be_ctx_text.replace('Dato de contexto: ', ''), max_len=112)}",
            font_name=FONT_BODY_BOOK,
            font_size=_scaled_font(12.6, scale=panel_scale, min_size=10.0),
            color_hex=PALETTE["muted"],
        )
    else:
        _set_rich_text(
            p_to_be_ctx,
            f"**Dato de contexto:** {_clip_text(to_be_ctx_text.replace('Dato de contexto: ', ''), max_len=112)}",
            font_name=FONT_BODY_BOOK,
            font_size=_scaled_font(12.6, scale=panel_scale, min_size=10.0),
            color_hex=PALETTE["muted"],
        )

    if not cards:
        p_none = itf.add_paragraph()
        p_none.text = "No se detectaron acciones concretas en este corte."
        p_none.space_after = Pt(4)
        p_none.font.name = FONT_BODY_BOOK
        p_none.font.size = Pt(_scaled_font(12.0, scale=panel_scale, min_size=10.0))
        p_none.font.color.rgb = _rgb(PALETTE["muted"])

    for idx_card, card in enumerate(cards, start=1):
        urgency_label, urgency_tone = _urgency_from_score(float(card.score))
        _, _, urgency_txt = _tone_style(urgency_tone)
        p_title = itf.add_paragraph()
        p_title.space_before = Pt(3)
        p_title.space_after = Pt(0)
        p_title.clear()
        run_num = p_title.add_run()
        run_num.text = f"{idx_card}. "
        run_num.font.name = FONT_BODY_MEDIUM
        run_num.font.bold = True
        run_num.font.size = Pt(_scaled_font(13.2, scale=panel_scale, min_size=10.8))
        run_num.font.color.rgb = _rgb(PALETTE["ink"])

        run_tag = p_title.add_run()
        run_tag.text = f"[{urgency_label.upper()}] "
        run_tag.font.name = FONT_BODY_MEDIUM
        run_tag.font.bold = True
        run_tag.font.size = Pt(_scaled_font(11.8, scale=panel_scale, min_size=9.8))
        run_tag.font.color.rgb = _rgb(urgency_txt)

        run_title = p_title.add_run()
        run_title.text = _clip_text(_soften_insight_tone(card.title), max_len=68)
        run_title.font.name = FONT_BODY_MEDIUM
        run_title.font.bold = True
        run_title.font.size = Pt(_scaled_font(14.2, scale=panel_scale, min_size=11.2))
        run_title.font.color.rgb = _rgb(PALETTE["ink"])

        p = itf.add_paragraph()
        p.space_after = Pt(3)
        _set_rich_text(
            p,
            _clip_text(
                _soften_insight_tone(f"**Recomendación de acción:** {card.body}"), max_len=178
            ),
            font_name=FONT_BODY_BOOK,
            font_size=_scaled_font(13.0, scale=panel_scale, min_size=10.2),
            color_hex=PALETTE["muted"],
        )

    _add_footer(slide, context=context)


def _norm_token(value: object) -> str:
    return str(value or "").strip().lower()


def _target_impact(action: ActionInsight, *, open_df: pd.DataFrame) -> Tuple[int, float, int]:
    total_open = int(len(open_df))
    if total_open <= 0:
        return 0, 0.0, 0

    status_filters = [_norm_token(v) for v in list(action.status_filters or []) if _norm_token(v)]
    priority_filters = [
        _norm_token(v) for v in list(action.priority_filters or []) if _norm_token(v)
    ]
    assignee_filters = [
        _norm_token(v) for v in list(action.assignee_filters or []) if _norm_token(v)
    ]
    specificity = (
        int(bool(status_filters)) + int(bool(priority_filters)) + int(bool(assignee_filters))
    )

    mask = pd.Series(True, index=open_df.index)
    applied_filters = 0

    if status_filters and "status" in open_df.columns:
        status_series = (
            normalize_text_col(open_df["status"], "(sin estado)")
            .astype(str)
            .str.lower()
            .str.strip()
        )
        mask &= status_series.isin(status_filters)
        applied_filters += 1
    if priority_filters and "priority" in open_df.columns:
        prio_series = (
            normalize_text_col(open_df["priority"], "(sin priority)")
            .astype(str)
            .str.lower()
            .str.strip()
        )
        mask &= prio_series.isin(priority_filters)
        applied_filters += 1
    if assignee_filters and "assignee" in open_df.columns:
        assignee_series = (
            normalize_text_col(open_df["assignee"], "(sin asignar)")
            .astype(str)
            .str.lower()
            .str.strip()
        )
        mask &= assignee_series.isin(assignee_filters)
        applied_filters += 1

    if applied_filters == 0:
        # Generic actions are still useful, but they should not dominate focused actions.
        fallback_ratio = 0.18
        fallback_count = max(1, int(round(total_open * fallback_ratio)))
        return fallback_count, fallback_ratio, specificity

    impacted = int(mask.sum())
    ratio = float(impacted) / float(max(total_open, 1))
    return impacted, ratio, specificity


def _actionability_strength(action: ActionInsight) -> float:
    text = f"{str(action.title or '')} {str(action.body or '')}".lower()
    cues = (
        "prioriza",
        "desbloq",
        "escal",
        "asign",
        "triage",
        "forzar",
        "cerr",
        "reduc",
        "war room",
        "ata",
    )
    hits = sum(1 for cue in cues if cue in text)
    return min(1.0, hits / 4.0)


def _is_finalist_centered_action(action: ActionInsight) -> bool:
    status_filters = [
        str(s or "").strip() for s in list(action.status_filters or []) if str(s or "").strip()
    ]
    if status_filters and all(_is_finalist_status(s) for s in status_filters):
        return True

    text = f"{str(action.title or '').lower()} {str(action.body or '').lower()}"
    if "accepted" in text or "ready to deploy" in text or "deployed" in text:
        if not any(
            token in text
            for token in (
                "blocked",
                "new",
                "analys",
                "progress",
                "rework",
                "test",
                "ready to verify",
                "en progreso",
            )
        ):
            return True
    return False


def _resolution_value_score(action: ActionInsight, *, open_df: pd.DataFrame) -> Tuple[float, int]:
    base_score = max(0.0, float(action.score))
    base_strength = min(1.0, math.log1p(base_score) / math.log1p(90.0))
    impacted, impact_ratio, specificity = _target_impact(action, open_df=open_df)

    urgency_strength = 1.0 if base_score >= 30.0 else 0.72 if base_score >= 18.0 else 0.42
    specificity_strength = min(1.0, float(specificity) / 3.0)
    actionability = _actionability_strength(action)

    value = 100.0 * (
        (0.30 * base_strength)
        + (0.44 * impact_ratio)
        + (0.12 * specificity_strength)
        + (0.08 * actionability)
        + (0.06 * urgency_strength)
    )

    text = f"{str(action.title or '').lower()} {str(action.body or '').lower()}"
    if any(tkn in text for tkn in ("datos insuficientes", "señal insuficiente", "sin insights")):
        value -= 22.0
    if specificity > 0 and impacted <= 0:
        value -= 26.0
    return value, impacted


def _best_actions(
    sections: Sequence[_ChartSection], *, open_df: pd.DataFrame, limit: int = 5
) -> List[ActionInsight]:
    all_cards: List[ActionInsight] = []
    for section in sections:
        all_cards.extend(section.insight_pack.cards)

    ranked_cards: List[Tuple[ActionInsight, float, int]] = []
    for card in all_cards:
        value_score, impacted = _resolution_value_score(card, open_df=open_df)
        ranked_cards.append((card, value_score, impacted))
    ranked_cards.sort(key=lambda item: (item[1], item[2], float(item[0].score)), reverse=True)

    out: List[ActionInsight] = []
    seen: set[str] = set()
    for card, _, _ in ranked_cards:
        if _is_finalist_centered_action(card):
            continue
        key = _soften_insight_tone(str(card.title or "").strip()).lower()
        if not key or key in seen:
            continue
        seen.add(key)
        title = _clip_text(_soften_insight_tone(str(card.title or "").strip()), max_len=58)
        body = _clip_text(_soften_insight_tone(str(card.body or "").strip()), max_len=126)
        out.append(
            ActionInsight(
                title=title,
                body=body,
                score=float(card.score),
                status_filters=card.status_filters,
                priority_filters=card.priority_filters,
                assignee_filters=card.assignee_filters,
            )
        )
        if len(out) >= limit:
            break
    if not out:
        out.append(
            ActionInsight(
                title="Señal insuficiente",
                body="No hay acciones sugeridas por falta de señal suficiente en el scope actual.",
                score=0.0,
            )
        )
    return out


def _estimate_wrapped_lines(text: str, *, max_chars_per_line: int) -> int:
    clean = re.sub(r"\s+", " ", str(text or "").strip())
    if not clean:
        return 1
    width = max(8, int(max_chars_per_line))
    return max(1, (len(clean) + width - 1) // width)


def _estimated_actions_panel_height_pt(actions: Sequence[ActionInsight]) -> float:
    # Estimated text height for the current final-slide typography.
    base_height = 48.0  # section title + top spacing
    for action in list(actions or []):
        title_lines = _estimate_wrapped_lines(str(action.title or ""), max_chars_per_line=48)
        body_lines = _estimate_wrapped_lines(
            f"Acción: {str(action.body or '')}", max_chars_per_line=62
        )
        base_height += (title_lines * 20.0) + (body_lines * 18.0) + 8.0
    return base_height


def _select_actions_for_final_slide(actions: Sequence[ActionInsight]) -> List[ActionInsight]:
    candidates = list(actions or [])[:4]
    if len(candidates) <= 3:
        return candidates
    max_height_pt = 368.0
    if _estimated_actions_panel_height_pt(candidates) <= max_height_pt:
        return candidates
    return candidates[:3]


def _add_final_summary_slide(prs: Any, context: _ScopeContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, "001C4A")
    _add_header(
        slide,
        title="Plan de acción",
        subtitle="Prioridades concretas para reducir backlog y acelerar el cierre",
        dark=True,
    )
    actions_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(1.20),
        Inches(12.00),
        Inches(5.64),
    )
    actions_panel.fill.solid()
    actions_panel.fill.fore_color.rgb = _rgb("FFFFFF")
    actions_panel.line.color.rgb = _rgb("D3D8E1")

    atf = actions_panel.text_frame
    atf.clear()
    p0 = atf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = "Acciones recomendadas (ordenadas por valor)"
    run0.font.name = FONT_HEAD
    run0.font.bold = True
    run0.font.size = Pt(29)
    run0.font.color.rgb = _rgb(PALETTE["navy"])

    final_actions = _select_actions_for_final_slide(
        _best_actions(context.sections, open_df=context.open_df, limit=4)
    )
    for idx, action in enumerate(final_actions, start=1):
        urgency_label, urgency_tone = _urgency_from_score(float(action.score))
        _, _, urgency_txt = _tone_style(urgency_tone)
        p_line = atf.add_paragraph()
        p_line.level = 0
        p_line.space_after = Pt(0)
        p_line.clear()

        run_num = p_line.add_run()
        run_num.text = f"{idx}. "
        run_num.font.name = FONT_BODY_MEDIUM
        run_num.font.bold = True
        run_num.font.size = Pt(17.2)
        run_num.font.color.rgb = _rgb(PALETTE["ink"])

        run_tag = p_line.add_run()
        run_tag.text = f"[{urgency_label.upper()}] "
        run_tag.font.name = FONT_BODY_MEDIUM
        run_tag.font.bold = True
        run_tag.font.size = Pt(13.4)
        run_tag.font.color.rgb = _rgb(urgency_txt)

        run_title = p_line.add_run()
        run_title.text = str(action.title or "").strip()
        run_title.font.name = FONT_BODY_MEDIUM
        run_title.font.bold = True
        run_title.font.size = Pt(17.2)
        run_title.font.color.rgb = _rgb(PALETTE["ink"])

        p_body = atf.add_paragraph()
        p_body.level = 0
        p_body.space_after = Pt(6)
        _set_rich_text(
            p_body,
            f"**Acción recomendada:** {str(action.body or '').strip()}",
            font_name=FONT_BODY_BOOK,
            font_size=15.0,
            color_hex=PALETTE["muted"],
        )

    foot = slide.shapes.add_textbox(Inches(0.62), Inches(6.96), Inches(12.0), Inches(0.24))
    ftf = foot.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    frun = fp.add_run()
    frun.text = f"{context.country} · {context.source_label} · {context.generated_at.strftime('%Y-%m-%d %H:%M UTC')}"
    frun.font.name = FONT_BODY_BOOK
    frun.font.size = Pt(10)
    frun.font.color.rgb = _rgb("BDD8FF")


def _compose_presentation(context: _ScopeContext) -> Any:
    prs: Any = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_cover_slide(prs, context)
    _add_index_slide(prs, context)
    _add_exec_summary_slide(prs, context)

    total = len(context.sections)
    for idx, section in enumerate(context.sections, start=1):
        _add_chart_insight_slide(prs, context=context, section=section, index=idx, total=total)

    _add_final_summary_slide(prs, context)
    return prs


def generate_scope_executive_ppt(
    settings: Settings,
    *,
    country: str,
    source_id: str,
    status_filters: Sequence[str] | None = None,
    priority_filters: Sequence[str] | None = None,
    assignee_filters: Sequence[str] | None = None,
    dff_override: pd.DataFrame | None = None,
    open_df_override: pd.DataFrame | None = None,
) -> ExecutiveReportResult:
    """Generate an executive PPT for selected scope using dashboard-equivalent visuals."""
    source_txt = str(source_id or "").strip()
    if not source_txt:
        raise ValueError("No se ha seleccionado un origen válido para generar el informe.")

    scoped_source_df = _scope_df(
        load_issues_df(settings.DATA_PATH),
        country=str(country or "").strip(),
        source_id=source_txt,
    )
    available_months = max_available_backlog_months(scoped_source_df)
    lookback_months = effective_analysis_lookback_months(settings, df=scoped_source_df)

    filters = _FilterSnapshot(
        status=_normalize_filter_values(status_filters),
        priority=_normalize_filter_values(priority_filters),
        assignee=_normalize_filter_values(assignee_filters),
        analysis_lookback_months=int(lookback_months),
        analysis_available_months=int(available_months),
    )

    context = _build_context(
        settings,
        country=str(country or "").strip(),
        source_id=source_txt,
        filters=filters,
        dff_override=dff_override,
        open_df_override=open_df_override,
        scoped_source_df=scoped_source_df if dff_override is None else None,
    )
    prs = _compose_presentation(context)

    buff = BytesIO()
    prs.save(buff)
    content = buff.getvalue()

    stamp = context.generated_at.strftime("%Y%m%d-%H%M")
    file_name = f"radar-{_slug(context.country)}-{_slug(context.source_id)}-{stamp}.pptx"

    return ExecutiveReportResult(
        file_name=file_name,
        content=content,
        slide_count=len(prs.slides),
        total_issues=int(len(context.dff)),
        open_issues=int(len(context.open_df)),
        closed_issues=int(len(context.closed_df)),
        country=context.country,
        source_id=context.source_id,
        source_label=context.source_label,
        applied_filter_summary=context.filters.summary(),
    )
