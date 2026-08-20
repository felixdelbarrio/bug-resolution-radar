"""Role-aware BBVA corporate branding for generated presentations."""

from __future__ import annotations

from enum import StrEnum
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from bug_resolution_radar.theme.brand_identity import (
    CORPORATE_DESCRIPTOR_LINES,
    CORPORATE_WORDMARK,
)
from bug_resolution_radar.theme.design_tokens import (
    BBVA_ELECTRIC,
    BBVA_FONT_SANS_MEDIUM_PPT,
    BBVA_MIDNIGHT,
    BBVA_SERENE,
    BBVA_WHITE,
    hex_to_rgb,
)

_EMU_PER_INCH = 914_400
_CORPORATE_PREFIX = "BBVA Corporate"
_LOCKUP_PREFIX = "BBVA Corporate Lockup"


class SlideBrandRole(StrEnum):
    COVER = "cover"
    SECTION = "section"
    CONTENT = "content"


def _rgb(color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(color))


def _fill_rgb(fill: Any) -> tuple[int, int, int] | None:
    try:
        rgb = fill.fore_color.rgb
        if rgb is None:
            return None
        return int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return None


def _is_dark_rgb(rgb: tuple[int, int, int] | None) -> bool:
    if rgb is None:
        return False
    return (0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]) < 125


def _dominant_dark_shape(slide: Any, *, slide_width: int, slide_height: int) -> Any | None:
    candidates: list[tuple[int, Any]] = []
    slide_area = max(int(slide_width) * int(slide_height), 1)
    for shape in slide.shapes:
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        area = width * height
        if area < int(slide_area * 0.60):
            continue
        if _is_dark_rgb(_fill_rgb(getattr(shape, "fill", None))):
            candidates.append((area, shape))
    return max(candidates, key=lambda item: item[0])[1] if candidates else None


def _slide_is_dark(slide: Any, *, slide_width: int, slide_height: int) -> bool:
    if _dominant_dark_shape(slide, slide_width=slide_width, slide_height=slide_height):
        return True
    rgb = _fill_rgb(getattr(getattr(slide, "background", None), "fill", None))
    return _is_dark_rgb(rgb)


def _is_brand_text(text: str) -> bool:
    normalized = " ".join(str(text or "").split()).casefold()
    descriptor = " ".join(CORPORATE_DESCRIPTOR_LINES).casefold()
    return normalized in {
        CORPORATE_WORDMARK.casefold(),
        descriptor,
        CORPORATE_DESCRIPTOR_LINES[0].casefold(),
        CORPORATE_DESCRIPTOR_LINES[1].casefold(),
    }


def _remove_shape(shape: Any) -> None:
    element = shape.element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _remove_existing_branding(slide: Any) -> None:
    for shape in list(slide.shapes):
        name = str(getattr(shape, "name", "") or "")
        text = (
            str(getattr(shape, "text", "") or "") if getattr(shape, "has_text_frame", False) else ""
        )
        if name.startswith(_CORPORATE_PREFIX) or _is_brand_text(text):
            _remove_shape(shape)


def _has_visual_or_table(slide: Any) -> bool:
    for shape in slide.shapes:
        if bool(getattr(shape, "has_table", False)):
            return True
        if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
            return True
    return False


def _non_brand_text(slide: Any) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join(str(getattr(shape, "text", "") or "").split())
        if text and not _is_brand_text(text):
            chunks.append(text)
    return " ".join(chunks)


def _role_for_slide(
    slide: Any,
    *,
    index: int,
    slide_width: int,
    slide_height: int,
) -> SlideBrandRole:
    if index == 0:
        return SlideBrandRole.COVER
    dominant = _dominant_dark_shape(
        slide,
        slide_width=slide_width,
        slide_height=slide_height,
    )
    if (
        dominant is not None
        and not _has_visual_or_table(slide)
        and len(_non_brand_text(slide)) <= 180
    ):
        return SlideBrandRole.SECTION
    return SlideBrandRole.CONTENT


def _set_full_bleed_background(
    slide: Any,
    *,
    role: SlideBrandRole,
    slide_width: int,
    slide_height: int,
) -> None:
    dominant = _dominant_dark_shape(
        slide,
        slide_width=slide_width,
        slide_height=slide_height,
    )
    color = (
        _fill_rgb(getattr(dominant, "fill", None)) if role == SlideBrandRole.SECTION else None
    ) or hex_to_rgb(BBVA_ELECTRIC)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _remove_inherited_brand_artifacts(slide: Any) -> None:
    """Remove obsolete BBVA-only artwork inherited from cover/divider layouts."""
    layout = getattr(slide, "slide_layout", None)
    for shape in list(getattr(layout, "shapes", ())):
        text = (
            str(getattr(shape, "text", "") or "") if getattr(shape, "has_text_frame", False) else ""
        )
        is_small_top_left_picture = (
            getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
            and int(getattr(shape, "left", 0) or 0) < int(Inches(2.0))
            and int(getattr(shape, "top", 0) or 0) < int(Inches(1.0))
            and int(getattr(shape, "width", 0) or 0) < int(Inches(2.0))
            and int(getattr(shape, "height", 0) or 0) < int(Inches(1.0))
        )
        if is_small_top_left_picture or _is_brand_text(text):
            _remove_shape(shape)


def _lockup_geometry(
    role: SlideBrandRole,
    *,
    slide_width: int,
) -> tuple[float, float, float, float, float]:
    width_in = float(slide_width) / _EMU_PER_INCH
    if role in {SlideBrandRole.COVER, SlideBrandRole.SECTION}:
        return 0.48, 0.32, 2.75, 0.42, 0.94
    lockup_width = 2.32
    return width_in - lockup_width - 0.22, 0.14, lockup_width, 0.31, 0.72


def _reserve_content_title_space(
    slide: Any,
    *,
    lockup_left: float,
    lockup_top: float,
    lockup_height: float,
) -> None:
    left_limit = int(Inches(lockup_left - 0.30))
    lockup_top_emu = int(Inches(lockup_top))
    lockup_bottom_emu = int(Inches(lockup_top + lockup_height))
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join(str(getattr(shape, "text", "") or "").split())
        if not text or _is_brand_text(text):
            continue
        top = int(getattr(shape, "top", 0) or 0)
        bottom = top + int(getattr(shape, "height", 0) or 0)
        right = int(getattr(shape, "left", 0) or 0) + int(getattr(shape, "width", 0) or 0)
        if top < lockup_bottom_emu and bottom > lockup_top_emu and right > left_limit:
            new_width = left_limit - int(getattr(shape, "left", 0) or 0)
            if new_width > int(Inches(2.0)):
                shape.width = new_width


def _add_lockup(
    slide: Any,
    *,
    role: SlideBrandRole,
    slide_width: int,
    slide_height: int,
) -> None:
    left, top, lockup_width, lockup_height, wordmark_width = _lockup_geometry(
        role,
        slide_width=slide_width,
    )
    if role == SlideBrandRole.CONTENT:
        _reserve_content_title_space(
            slide,
            lockup_left=left,
            lockup_top=top,
            lockup_height=lockup_height,
        )
    divider_x = left + wordmark_width + 0.09
    descriptor_x = divider_x + 0.12
    descriptor_width = lockup_width - (descriptor_x - left)
    dark = _slide_is_dark(slide, slide_width=slide_width, slide_height=slide_height)
    wordmark_color = BBVA_WHITE if dark else BBVA_ELECTRIC
    descriptor_color = BBVA_SERENE if dark else BBVA_MIDNIGHT
    wordmark_size = 16.0 if role in {SlideBrandRole.COVER, SlideBrandRole.SECTION} else 11.5
    descriptor_size = 7.4 if role in {SlideBrandRole.COVER, SlideBrandRole.SECTION} else 6.0

    wordmark = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(wordmark_width), Inches(lockup_height)
    )
    wordmark.name = f"{_LOCKUP_PREFIX} {role.value} Wordmark"
    wordmark.text_frame.clear()
    wordmark.text_frame.margin_left = 0
    wordmark.text_frame.margin_right = 0
    wordmark.text_frame.margin_top = 0
    wordmark.text_frame.margin_bottom = 0
    wordmark.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = wordmark.text_frame.paragraphs[0].add_run()
    run.text = CORPORATE_WORDMARK
    run.font.name = BBVA_FONT_SANS_MEDIUM_PPT
    run.font.bold = True
    run.font.size = Pt(wordmark_size)
    run.font.color.rgb = _rgb(wordmark_color)

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(divider_x),
        Inches(top + 0.04),
        Inches(0.014),
        Inches(lockup_height - 0.08),
    )
    divider.name = f"{_LOCKUP_PREFIX} {role.value} Divider"
    divider.fill.solid()
    divider.fill.fore_color.rgb = _rgb(descriptor_color)
    divider.line.fill.background()

    descriptor = slide.shapes.add_textbox(
        Inches(descriptor_x),
        Inches(top),
        Inches(descriptor_width),
        Inches(lockup_height),
    )
    descriptor.name = f"{_LOCKUP_PREFIX} {role.value} Descriptor"
    descriptor.text_frame.clear()
    descriptor.text_frame.margin_left = 0
    descriptor.text_frame.margin_right = 0
    descriptor.text_frame.margin_top = 0
    descriptor.text_frame.margin_bottom = 0
    descriptor.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    descriptor_run = descriptor.text_frame.paragraphs[0].add_run()
    descriptor_run.text = "\n".join(CORPORATE_DESCRIPTOR_LINES)
    descriptor_run.font.name = BBVA_FONT_SANS_MEDIUM_PPT
    descriptor_run.font.size = Pt(descriptor_size)
    descriptor_run.font.color.rgb = _rgb(descriptor_color)


def apply_corporate_branding(presentation: Any) -> None:
    """Apply one canonical lockup and geometry policy to the complete deck."""
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    for index, slide in enumerate(presentation.slides):
        _remove_existing_branding(slide)
        role = _role_for_slide(
            slide,
            index=index,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        if role in {SlideBrandRole.COVER, SlideBrandRole.SECTION}:
            _set_full_bleed_background(
                slide,
                role=role,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            _remove_inherited_brand_artifacts(slide)
        _add_lockup(
            slide,
            role=role,
            slide_width=slide_width,
            slide_height=slide_height,
        )
