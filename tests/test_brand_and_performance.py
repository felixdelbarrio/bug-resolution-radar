from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from bug_resolution_radar.config import Settings
from bug_resolution_radar.models.schema import IssuesDocument, NormalizedIssue
from bug_resolution_radar.reports.branding import apply_corporate_branding
from bug_resolution_radar.repositories.issues_store import load_issues_meta, save_issues_doc
from bug_resolution_radar.services.issue_enrichment import enrich_issue_dataframe_with_helix
from bug_resolution_radar.theme.brand_identity import frontend_brand_contract


def test_brand_contract_and_presentation_lockup_are_canonical() -> None:
    contract = frontend_brand_contract()
    assert contract == {
        "name": "BBVA Banca de Empresas e Instituciones",
        "wordmark": "BBVA",
        "descriptorLines": ["Banca de Empresas", "e Instituciones"],
    }

    presentation = Presentation()
    cover_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    legacy = cover_slide.shapes.add_textbox(0, 0, 1_000_000, 300_000)
    legacy.text = "BBVA"
    section_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = section_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        300_000,
        200_000,
        presentation.slide_width - 600_000,
        presentation.slide_height - 400_000,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(7, 14, 70)
    section_title = section_slide.shapes.add_textbox(700_000, 2_000_000, 5_000_000, 800_000)
    section_title.text = "Dashboard de KPIs"
    content_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    content_title = content_slide.shapes.add_textbox(300_000, 120_000, 10_500_000, 500_000)
    content_title.text = "Seguimiento de incidencias"
    content_body = content_slide.shapes.add_textbox(300_000, 1_000_000, 9_000_000, 1_000_000)
    content_body.text = "Detalle ejecutivo " * 20

    apply_corporate_branding(presentation)
    apply_corporate_branding(presentation)

    for slide in (cover_slide, section_slide, content_slide):
        brand_shapes = [
            shape for shape in slide.shapes if shape.name.startswith("BBVA Corporate Lockup")
        ]
        assert len(brand_shapes) == 3
        assert any(getattr(shape, "text", "") == "BBVA" for shape in brand_shapes)
        assert any("Banca de Empresas" in getattr(shape, "text", "") for shape in brand_shapes)

    cover_wordmark = next(
        shape for shape in cover_slide.shapes if shape.name.endswith("cover Wordmark")
    )
    section_wordmark = next(
        shape for shape in section_slide.shapes if shape.name.endswith("section Wordmark")
    )
    content_wordmark = next(
        shape for shape in content_slide.shapes if shape.name.endswith("content Wordmark")
    )
    assert cover_wordmark.left == section_wordmark.left
    assert cover_wordmark.top == section_wordmark.top
    assert content_wordmark.top < int(presentation.slide_height * 0.10)
    content_descriptor = next(
        shape for shape in content_slide.shapes if shape.name.endswith("content Descriptor")
    )
    assert int(content_descriptor.left) + int(content_descriptor.width) > int(
        presentation.slide_width * 0.96
    )
    assert int(content_title.left) + int(content_title.width) < int(content_wordmark.left)

    for slide in (cover_slide, section_slide):
        full_bleed = next(
            shape for shape in slide.shapes if shape.name == "BBVA Corporate Full Bleed Background"
        )
        assert (full_bleed.left, full_bleed.top) == (0, 0)
        assert (full_bleed.width, full_bleed.height) == (
            presentation.slide_width,
            presentation.slide_height,
        )
    assert (
        sum(
            1
            for shape in cover_slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text == "BBVA"
        )
        == 1
    )


def test_issue_metadata_sidecar_avoids_loading_full_document(tmp_path: Path) -> None:
    path = tmp_path / "issues.json"
    save_issues_doc(
        str(path),
        IssuesDocument(
            schema_version="1.0",
            ingested_at="2026-08-20T09:00:00+00:00",
            jira_base_url="https://jira.example.com",
            query="project = RADAR",
            issues=[
                NormalizedIssue(
                    key="RAD-1",
                    summary="Incidencia",
                    status="Open",
                    type="Bug",
                    priority="Medium",
                    source_id="jira:espana:core",
                    source_type="jira",
                )
            ],
        ),
    )

    metadata = load_issues_meta(str(path))
    assert metadata["issues_count"] == 1
    assert metadata["jira_source_count"] == 1
    assert metadata["query"] == "project = RADAR"


def test_helix_enrichment_prefers_columnar_sidecar(tmp_path: Path) -> None:
    helix_path = tmp_path / "helix.json"
    helix_path.write_text("not-json", encoding="utf-8")
    pd.DataFrame(
        [
            {
                "id": "INC001",
                "source_id": "helix:mexico:gema",
                "BBVA_ExecutiveDescription": "Servicio recuperado",
            }
        ]
    ).to_parquet(helix_path.with_suffix(".raw.parquet"), index=False)
    settings = Settings(HELIX_DATA_PATH=str(helix_path))
    source = pd.DataFrame(
        [
            {
                "key": "INC001",
                "source_id": "helix:mexico:gema",
                "helix_executive_description": "",
            }
        ]
    )

    enriched = enrich_issue_dataframe_with_helix(source, settings=settings)
    assert enriched.loc[0, "helix_executive_description"] == "Servicio recuperado"
