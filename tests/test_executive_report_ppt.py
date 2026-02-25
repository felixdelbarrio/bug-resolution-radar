from __future__ import annotations

from datetime import datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
import time

import pandas as pd
import pytest
from pptx import Presentation

from bug_resolution_radar.config import Settings
from bug_resolution_radar.reports import generate_scope_executive_ppt
from bug_resolution_radar.reports.executive_ppt import (
    _best_actions,
    _build_sections,
    _call_in_subprocess_with_timeout,
    _ChartSection,
    _fig_to_png,
    _is_finalist_status,
    _open_closed,
    _select_actions_for_final_slide,
    _soften_insight_tone,
    _urgency_from_score,
)
from bug_resolution_radar.models.schema import IssuesDocument, NormalizedIssue
from bug_resolution_radar.ui.common import save_issues_doc
from bug_resolution_radar.ui.dashboard.registry import ChartContext, _render_open_priority_pie
from bug_resolution_radar.ui.insights.engine import ActionInsight, TrendInsightPack


def _seed_issues(path: Path) -> None:
    doc = IssuesDocument(
        schema_version="1.0",
        ingested_at="2026-02-20T10:00:00+00:00",
        jira_base_url="https://jira.example.com",
        query="multi-source",
        issues=[
            NormalizedIssue(
                key="MX-1",
                summary="Softoken no valida",
                status="New",
                type="Bug",
                priority="High",
                created="2026-02-01T10:00:00+00:00",
                updated="2026-02-02T09:00:00+00:00",
                resolved=None,
                assignee="Ana",
                country="México",
                source_type="jira",
                source_alias="Core MX",
                source_id="jira:mexico:core-mx",
            ),
            NormalizedIssue(
                key="MX-2",
                summary="Error en transferencias SPEI",
                status="Blocked",
                type="Bug",
                priority="Highest",
                created="2026-01-25T10:00:00+00:00",
                updated="2026-02-03T11:00:00+00:00",
                resolved=None,
                assignee="Luis",
                country="México",
                source_type="jira",
                source_alias="Core MX",
                source_id="jira:mexico:core-mx",
            ),
            NormalizedIssue(
                key="MX-3",
                summary="Ajuste de dashboard",
                status="Deployed",
                type="Task",
                priority="Medium",
                created="2026-01-10T10:00:00+00:00",
                updated="2026-01-20T10:00:00+00:00",
                resolved="2026-01-20T10:00:00+00:00",
                assignee="Ana",
                country="México",
                source_type="jira",
                source_alias="Core MX",
                source_id="jira:mexico:core-mx",
            ),
            NormalizedIssue(
                key="MX-9",
                summary="Otro origen MX",
                status="In Progress",
                type="Bug",
                priority="Low",
                created="2026-02-03T10:00:00+00:00",
                updated="2026-02-04T10:00:00+00:00",
                resolved=None,
                assignee="Otro",
                country="México",
                source_type="jira",
                source_alias="Retail MX",
                source_id="jira:mexico:retail-mx",
            ),
            NormalizedIssue(
                key="ES-1",
                summary="Issue España",
                status="In Progress",
                type="Bug",
                priority="Low",
                created="2026-02-03T10:00:00+00:00",
                updated="2026-02-04T10:00:00+00:00",
                resolved=None,
                assignee="Marta",
                country="España",
                source_type="jira",
                source_alias="Core ES",
                source_id="jira:espana:core-es",
            ),
        ],
    )
    save_issues_doc(str(path), doc)


def _slide_text(slide: object) -> str:
    lines: list[str] = []
    for shape in getattr(slide, "shapes", []):
        txt = str(getattr(shape, "text", "") or "").strip()
        if txt:
            lines.append(txt)
    return "\n".join(lines)


def test_call_in_subprocess_with_timeout_returns_result() -> None:
    out = _call_in_subprocess_with_timeout(sum, [2, 3, 5], hard_timeout_s=5)
    assert out == 10


def test_call_in_subprocess_with_timeout_raises_timeout() -> None:
    started = time.monotonic()
    with pytest.raises(TimeoutError):
        _call_in_subprocess_with_timeout(time.sleep, 5, hard_timeout_s=1)
    elapsed = time.monotonic() - started
    assert elapsed < 4.0


def test_generate_scope_executive_ppt_is_scoped_and_valid_ppt(tmp_path: Path) -> None:
    issues_path = tmp_path / "issues.json"
    _seed_issues(issues_path)

    settings = Settings(
        DATA_PATH=str(issues_path),
        SUPPORTED_COUNTRIES="México,España",
        JIRA_SOURCES_JSON=(
            '[{"country":"México","alias":"Core MX","jql":"project = MX"},'
            '{"country":"México","alias":"Retail MX","jql":"project = RETAIL"},'
            '{"country":"España","alias":"Core ES","jql":"project = ES"}]'
        ),
    )

    out = generate_scope_executive_ppt(
        settings,
        country="México",
        source_id="jira:mexico:core-mx",
    )

    assert out.total_issues == 3
    assert out.open_issues == 2
    assert out.closed_issues == 1
    assert out.slide_count >= 8
    assert out.file_name.endswith(".pptx")
    assert out.content
    assert "Ventana temporal" in out.applied_filter_summary
    assert "Estado: Todos" in out.applied_filter_summary

    prs = Presentation(BytesIO(out.content))
    assert len(prs.slides) == out.slide_count
    cover_txt = _slide_text(prs.slides[0])
    index_txt = _slide_text(prs.slides[1])
    assert "México" in cover_txt
    assert "Índice" in index_txt


def test_generate_scope_executive_ppt_applies_filters(tmp_path: Path) -> None:
    issues_path = tmp_path / "issues.json"
    _seed_issues(issues_path)

    settings = Settings(DATA_PATH=str(issues_path))

    out = generate_scope_executive_ppt(
        settings,
        country="México",
        source_id="jira:mexico:core-mx",
        status_filters=["New"],
        priority_filters=["High"],
        assignee_filters=["Ana"],
    )

    assert out.total_issues == 1
    assert out.open_issues == 1
    assert out.closed_issues == 0
    assert "Estado: New" in out.applied_filter_summary
    assert "Prioridad: High" in out.applied_filter_summary
    assert "Responsable: Ana" in out.applied_filter_summary


def test_generate_scope_executive_ppt_applies_analysis_depth_window(tmp_path: Path) -> None:
    now = datetime.now(timezone.utc)
    issues_path = tmp_path / "issues.json"
    doc = IssuesDocument(
        schema_version="1.0",
        ingested_at=now.isoformat(),
        jira_base_url="https://jira.example.com",
        query="window",
        issues=[
            NormalizedIssue(
                key="MX-W1",
                summary="Issue reciente",
                status="New",
                type="Bug",
                priority="High",
                created=(now - timedelta(days=15)).isoformat(),
                updated=(now - timedelta(days=5)).isoformat(),
                resolved=None,
                assignee="Ana",
                country="México",
                source_type="jira",
                source_alias="Core MX",
                source_id="jira:mexico:core-mx",
            ),
            NormalizedIssue(
                key="MX-W2",
                summary="Issue antigua",
                status="Blocked",
                type="Bug",
                priority="Highest",
                created=(now - timedelta(days=130)).isoformat(),
                updated=(now - timedelta(days=90)).isoformat(),
                resolved=None,
                assignee="Luis",
                country="México",
                source_type="jira",
                source_alias="Core MX",
                source_id="jira:mexico:core-mx",
            ),
        ],
    )
    save_issues_doc(str(issues_path), doc)

    settings = Settings(DATA_PATH=str(issues_path), ANALYSIS_LOOKBACK_MONTHS=2)

    out = generate_scope_executive_ppt(
        settings,
        country="México",
        source_id="jira:mexico:core-mx",
    )

    assert out.total_issues == 1
    assert out.open_issues == 1
    assert "últimos 2 de" in out.applied_filter_summary


def test_generate_scope_executive_ppt_raises_when_scope_has_no_data(tmp_path: Path) -> None:
    issues_path = tmp_path / "issues.json"
    _seed_issues(issues_path)

    settings = Settings(DATA_PATH=str(issues_path))

    with pytest.raises(ValueError):
        generate_scope_executive_ppt(
            settings,
            country="México",
            source_id="jira:mexico:does-not-exist",
        )


def test_fig_to_png_renders_open_priority_pie() -> None:
    df = pd.DataFrame(
        [
            {"priority": "Highest", "status": "New", "created": "2026-02-01T00:00:00+00:00"},
            {"priority": "High", "status": "Analysing", "created": "2026-02-02T00:00:00+00:00"},
            {"priority": "Medium", "status": "Blocked", "created": "2026-02-03T00:00:00+00:00"},
            {"priority": "Low", "status": "En progreso", "created": "2026-02-04T00:00:00+00:00"},
            {"priority": "Low", "status": "Deployed", "created": "2026-02-05T00:00:00+00:00"},
        ]
    )
    ctx = ChartContext(dff=df, open_df=df, kpis={})
    fig = _render_open_priority_pie(ctx)
    assert fig is not None

    image = _fig_to_png(fig)
    if image is None:
        pytest.skip("Exportador de imágenes de Plotly no disponible en este entorno.")
    assert image is not None
    assert len(image) > 1_000


def test_build_sections_skips_resolution_chart_when_no_closed_data() -> None:
    dff = pd.DataFrame(
        [
            {
                "key": "MX-100",
                "status": "Analysing",
                "priority": "High",
                "created": "2026-02-01T00:00:00+00:00",
                "resolved": None,
            },
            {
                "key": "MX-101",
                "status": "Blocked",
                "priority": "Medium",
                "created": "2026-02-02T00:00:00+00:00",
                "resolved": None,
            },
        ]
    )
    dff["created"] = pd.to_datetime(dff["created"], utc=True, errors="coerce")
    dff["resolved"] = pd.to_datetime(dff["resolved"], utc=True, errors="coerce")
    open_df = dff.copy()

    sections = _build_sections(Settings(DATA_PATH="unused.json"), dff=dff, open_df=open_df)
    by_id = {sec.chart_id: sec for sec in sections}
    assert "resolution_hist" not in by_id


def test_select_actions_for_final_slide_returns_4_when_text_fits() -> None:
    actions = [
        ActionInsight(
            title=f"Acción {idx}", body="Texto breve para ejecutar pronto.", score=10.0 - idx
        )
        for idx in range(1, 5)
    ]
    selected = _select_actions_for_final_slide(actions)
    assert len(selected) == 4


def test_select_actions_for_final_slide_returns_3_when_text_is_long() -> None:
    long_body = (
        "Esta acción requiere coordinación entre equipos, validación funcional y técnica, "
        "seguimiento diario de bloqueos, decisión de alcance y cierre con evidencias operativas."
    )
    actions = [
        ActionInsight(
            title=(
                f"Acción estratégica {idx} para estabilizar el flujo y reducir fricción acumulada"
            ),
            body=long_body,
            score=20.0 - idx,
        )
        for idx in range(1, 5)
    ]
    selected = _select_actions_for_final_slide(actions)
    assert len(selected) == 3


def test_best_actions_prioritizes_resolution_value_over_raw_score() -> None:
    open_df = pd.DataFrame(
        [{"status": "New", "priority": "High", "assignee": "A"} for _ in range(70)]
        + [{"status": "Deployed", "priority": "Low", "assignee": "B"} for _ in range(30)]
    )
    high_impact = ActionInsight(
        title="Desbloqueo triage crítico",
        body="Prioriza el triage y la asignación de críticos en estados iniciales.",
        score=21.0,
        status_filters=["New"],
        priority_filters=["High"],
    )
    generic_high_score = ActionInsight(
        title="Seguimiento general",
        body="Mantener seguimiento transversal semanal.",
        score=36.0,
    )
    section = _ChartSection(
        chart_id="open_status_bar",
        theme="Flujo",
        title="Backlog por estado",
        subtitle="",
        figure=None,
        insight_pack=TrendInsightPack(
            metrics=[], cards=[generic_high_score, high_impact], executive_tip=None
        ),
    )

    ranked = _best_actions([section], open_df=open_df, limit=2)
    assert "desbloqueo triage" in ranked[0].title.lower()
    # The deck already shows urgency as a separate [MUST]/[SHOULD]/[NICE] tag, so titles
    # should not repeat "Must" and must stay business-readable.
    assert "must" not in ranked[0].title.lower()


def test_urgency_uses_must_should_nice_to_have() -> None:
    assert _urgency_from_score(24.0)[0] == "Must"
    assert _urgency_from_score(12.0)[0] == "Should"
    assert _urgency_from_score(4.0)[0] == "Nice to have"


def test_soften_insight_tone_reduces_catastrophic_language() -> None:
    txt = "Urgencia crítica: Criticidad atrapada en entrada."
    out = _soften_insight_tone(txt)
    assert "Urgencia" not in out
    assert "crítica" not in out.lower()
    assert "priorización" in out.lower()
    assert "mayor impacto" in out.lower()


def test_is_finalist_status_detects_terminal_flow_states() -> None:
    assert _is_finalist_status("Accepted")
    assert _is_finalist_status("Ready to deploy")
    assert _is_finalist_status("Deployed")
    assert not _is_finalist_status("Analysing")


def test_open_closed_treats_accepted_without_resolved_as_closed() -> None:
    df = pd.DataFrame(
        {
            "key": ["A-1", "A-2", "A-3"],
            "status": ["New", "Accepted", "Blocked"],
            "resolved": [pd.NaT, pd.NaT, "2026-02-01T00:00:00+00:00"],
        }
    )
    open_df, closed_df = _open_closed(df)
    assert set(open_df["key"].astype(str).tolist()) == {"A-1"}
    assert set(closed_df["key"].astype(str).tolist()) == {"A-2", "A-3"}
