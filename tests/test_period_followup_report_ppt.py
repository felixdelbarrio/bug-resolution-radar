from __future__ import annotations

import base64
import re
from io import BytesIO
from pathlib import Path
from typing import Any
from zipfile import ZipFile

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

from bug_resolution_radar.config import Settings, bundled_period_ppt_template_path
from bug_resolution_radar.models.schema import IssuesDocument, NormalizedIssue
from bug_resolution_radar.reports import generate_country_period_followup_ppt
from bug_resolution_radar.reports import period_followup_ppt as period_ppt_mod
from bug_resolution_radar.reports.period_followup_layout import metric_card_typography
from bug_resolution_radar.repositories.issues_store import save_issues_doc
from bug_resolution_radar.theme.design_tokens import (
    BBVA_REPORT_AMBER_BG,
    BBVA_REPORT_RED_BG,
    EXEC_CHART_AXIS_FONT_PT,
    EXEC_CHART_EXPORT_HEIGHT,
    EXEC_CHART_EXPORT_WIDTH,
    EXEC_CHART_INSIDE_VALUE_FONT_PT,
    EXEC_CHART_LEGEND_FONT_PT,
    EXEC_CHART_TOTAL_FONT_PT,
    EXEC_CHART_TREND_EXPORT_HEIGHT,
    hex_to_rgb,
)


def _build_minimal_template(path: Path) -> None:
    prs = Presentation()
    for idx in range(9):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if idx == 0:
            tb = slide.shapes.add_textbox(0, 0, 4_000_000, 600_000)
            tb.text = "Periodo dd/mm - dd/mm 2026"
        else:
            tb = slide.shapes.add_textbox(0, 0, 4_000_000, 600_000)
            tb.text = f"Slide {idx + 1}"
    prs.save(str(path))


def _build_compact_template(path: Path) -> None:
    prs = Presentation()
    for idx in range(7):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(0, 0, 4_000_000, 600_000)
        if idx == 0:
            tb.text = "Periodo dd/mm - dd/mm 2026"
        elif idx == 1:
            tb.text = "Dashboard de KPIs"
        elif idx in (2, 3, 4):
            tb.text = "Seguimiento de incidencias - Resumen ejecutivo"
        elif idx == 5:
            tb.text = "Gráficos de evolución"
        else:
            tb.text = "Seguimiento de KPIs - Gráficos"
    prs.save(str(path))


def _slide_text(slide: Any) -> str:
    return " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def _slide_table_text(slide: Any) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                chunks.append(str(cell.text or ""))
    return " ".join(chunks)


def _slide_all_text(slide: Any) -> str:
    return f"{_slide_text(slide)} {_slide_table_text(slide)}".strip()


def _find_slide_index(prs: Presentation, needle: str) -> int:
    for idx, slide in enumerate(prs.slides):
        if needle in _slide_all_text(slide):
            return idx
    raise AssertionError(f"No slide contains {needle!r}")


def _native_tables(slide: Any) -> list[Any]:
    return [shape for shape in slide.shapes if getattr(shape, "has_table", False)]


def _table_intersects_picture(slide: Any, table_shape: Any) -> bool:
    table_left = int(getattr(table_shape, "left", 0) or 0)
    table_top = int(getattr(table_shape, "top", 0) or 0)
    table_right = table_left + int(getattr(table_shape, "width", 0) or 0)
    table_bottom = table_top + int(getattr(table_shape, "height", 0) or 0)
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            continue
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        right = left + int(getattr(shape, "width", 0) or 0)
        bottom = top + int(getattr(shape, "height", 0) or 0)
        if left < table_right and right > table_left and top < table_bottom and bottom > table_top:
            return True
    return False


def _assert_native_table_font_floor(table_shape: Any) -> None:
    table = table_shape.table
    for ridx, row in enumerate(table.rows):
        min_size = 8.0 if ridx == 0 else 9.0
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not str(getattr(run, "text", "") or "").strip():
                        continue
                    assert run.font.name in {"Arial", "Aptos"}
                    assert run.font.size is not None
                    assert float(run.font.size.pt) >= min_size


def test_generate_country_period_followup_ppt_with_minimal_template(tmp_path: Path) -> None:
    template = tmp_path / "template.pptx"
    _build_minimal_template(template)

    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(template))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )

    assert out.slide_count == 15
    assert out.total_issues == 2
    assert out.open_issues == 1
    assert out.closed_issues == 1
    assert out.content
    prs = Presentation(BytesIO(out.content))
    assert len(prs.slides) == 15
    deck_text = " ".join(_slide_text(slide) for slide in prs.slides)
    assert "Incidencias abiertas por criticidad alta" in deck_text
    assert "Incidencias abiertas con más de 30 días" in deck_text
    dashboard_idx = _find_slide_index(
        prs, "Seguimiento de KPIs - Incidencias abiertas por funcionalidad"
    )
    # Functional follow-up slides must preserve light background from source template.
    bg_fill = prs.slides[dashboard_idx].background.fill
    assert int(bg_fill.type or 0) == 1
    assert bg_fill.fore_color.rgb == RGBColor(247, 248, 248)


def test_generate_country_period_followup_ppt_with_compact_template(tmp_path: Path) -> None:
    template = tmp_path / "compact-template.pptx"
    _build_compact_template(template)

    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(template))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )

    assert out.slide_count == 15
    assert out.total_issues == 2
    assert out.open_issues == 1
    assert out.closed_issues == 1
    assert out.content


def test_period_followup_ppt_finalist_discrepancies_section_is_always_included(
    tmp_path: Path,
) -> None:
    template = tmp_path / "template.pptx"
    _build_minimal_template(template)
    now = pd.Timestamp("2026-05-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "EAM-94000",
                "summary": "Jira pendiente",
                "description": "Helix INC000104154954",
                "status": "To Rework",
                "priority": "High",
                "assignee": "Ana",
                "created": (now - pd.Timedelta(days=20)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
                "url": "https://jira.example.com/browse/EAM-94000",
            },
            {
                "key": "EAM-93998",
                "summary": "Jira pendiente adicional",
                "description": "Tambien ligada a INC000104154954",
                "status": "To Rework",
                "priority": "High",
                "assignee": "Bea",
                "created": (now - pd.Timedelta(days=25)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
                "url": "https://jira.example.com/browse/EAM-93998",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )
    discrepancies = pd.DataFrame(
        [
            {
                "helix_id": "INC000104154954",
                "helix_summary": "Helix cerrado",
                "helix_description": "Resolución de INC000104154954 validada en Helix",
                "helix_status": "Closed",
                "helix_url": "https://helix.example.com/INC000104154954",
                "jira_key": "EAM-94000",
                "jira_summary": "Jira pendiente",
                "jira_status": "To Rework",
                "jira_priority": "High",
                "jira_assignee": "Ana",
                "jira_open_days": 20,
                "jira_url": "https://jira.example.com/browse/EAM-94000",
                "source_alias": "Senda",
            },
            {
                "helix_id": "INC000104154954",
                "helix_summary": "Helix cerrado",
                "helix_description": "Resolución de INC000104154954 validada en Helix",
                "helix_status": "Closed",
                "helix_url": "https://helix.example.com/INC000104154954",
                "jira_key": "EAM-93998",
                "jira_summary": "Jira pendiente adicional",
                "jira_status": "To Rework",
                "jira_priority": "High",
                "jira_assignee": "Bea",
                "jira_open_days": 25,
                "jira_url": "https://jira.example.com/browse/EAM-93998",
                "source_alias": "Senda",
            },
        ]
    )

    result = generate_country_period_followup_ppt(
        Settings(PERIOD_PPT_TEMPLATE_PATH=str(template)),
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        finalist_discrepancies_override=discrepancies,
        reference_day=now,
    )
    prs = Presentation(BytesIO(result.content))
    text = " ".join(_slide_all_text(slide) for slide in prs.slides)

    assert "Incidencias con discrepancias en estado finalista" in text
    assert "JIRA: To Rework" in text
    assert "Helix: Closed" in text
    assert "EAM-93998" in text
    assert any(
        "Incidencias abiertas por criticidad alta" in _slide_all_text(slide)
        and "EAM-93998" in _slide_all_text(slide)
        for slide in prs.slides
    )

    finalist_slide = next(
        slide
        for slide in prs.slides
        if "Incidencias con discrepancias en estado finalista" in _slide_all_text(slide)
        and "EAM-94000" in _slide_all_text(slide)
    )
    table = _native_tables(finalist_slide)[0].table
    id_values = [str(table.cell(row_idx, 0).text or "") for row_idx in range(1, len(table.rows))]
    assert "EAM-94000" in id_values
    assert "EAM-93998" in id_values
    assert all("INC000104154954" not in value for value in id_values)

    description_values = [
        str(table.cell(row_idx, 1).text or "") for row_idx in range(1, len(table.rows))
    ]
    assert any("INC000104154954" in value for value in description_values)
    first_id_runs = list(table.cell(1, 0).text_frame.paragraphs[0].runs)
    assert first_id_runs
    assert str(first_id_runs[0].hyperlink.address or "").startswith("https://jira.example.com")
    description_runs = [
        run for paragraph in table.cell(1, 1).text_frame.paragraphs for run in paragraph.runs
    ]
    assert any(
        "INC000104154954" in str(run.text or "")
        and str(run.hyperlink.address or "").startswith("https://helix.example.com")
        for run in description_runs
    )


def test_period_followup_ppt_splits_root_cause_evolutives_and_renders_notes(
    tmp_path: Path,
) -> None:
    template = tmp_path / "template.pptx"
    notes_path = tmp_path / "notes.json"
    _build_minimal_template(template)
    notes_path.write_text(
        '{"EAM-ROOT":"Cliente confirma evolutivo priorizado para resolver la causa raíz."}',
        encoding="utf-8",
    )
    now = pd.Timestamp("2026-05-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "EAM-ROOT",
                "summary": "Evolutivo causa raíz",
                "description": "Helix INC000104154954",
                "status": "To Rework",
                "priority": "High",
                "assignee": "Ana",
                "created": (now - pd.Timedelta(days=20)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
                "url": "https://jira.example.com/browse/EAM-ROOT",
            },
            {
                "key": "EAM-REG",
                "summary": "Discrepancia regular",
                "description": "Helix INC000104154955",
                "status": "To Rework",
                "priority": "High",
                "assignee": "Bea",
                "created": (now - pd.Timedelta(days=25)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
                "url": "https://jira.example.com/browse/EAM-REG",
            },
        ]
    )
    discrepancies = pd.DataFrame(
        [
            {
                "helix_id": "INC000104154954",
                "helix_summary": "Helix cerrado root",
                "helix_description": "Resolución root",
                "helix_status": "Closed",
                "helix_url": "https://helix.example.com/INC000104154954",
                "jira_key": "EAM-ROOT",
                "jira_summary": "Evolutivo causa raíz",
                "jira_status": "To Rework",
                "jira_priority": "High",
                "jira_assignee": "Ana",
                "jira_open_days": 20,
                "jira_url": "https://jira.example.com/browse/EAM-ROOT",
                "jira_labels": ("CAUSA_RAIZ",),
                "source_alias": "Senda",
            },
            {
                "helix_id": "INC000104154955",
                "helix_summary": "Helix cerrado regular",
                "helix_description": "Resolución regular",
                "helix_status": "Closed",
                "helix_url": "https://helix.example.com/INC000104154955",
                "jira_key": "EAM-REG",
                "jira_summary": "Discrepancia regular",
                "jira_status": "To Rework",
                "jira_priority": "High",
                "jira_assignee": "Bea",
                "jira_open_days": 25,
                "jira_url": "https://jira.example.com/browse/EAM-REG",
                "jira_labels": (),
                "source_alias": "Senda",
            },
        ]
    )

    result = generate_country_period_followup_ppt(
        Settings(
            PERIOD_PPT_TEMPLATE_PATH=str(template),
            NOTES_PATH=str(notes_path),
            JIRA_ROOT_CAUSE_LABELS_BY_COUNTRY_JSON=(
                '[{"country":"México","labels":["causa_raiz"]}]'
            ),
        ),
        country="México",
        source_ids=["jira:mexico:senda"],
        dff_override=dff,
        finalist_discrepancies_override=discrepancies,
        reference_day=now,
    )

    prs = Presentation(BytesIO(result.content))
    slide_texts = [_slide_all_text(slide) for slide in prs.slides]
    root_idx = next(
        idx
        for idx, text in enumerate(slide_texts)
        if "Evolutivos para solucionar causas raíces" in text
    )
    finalist_idx = next(
        idx
        for idx, text in enumerate(slide_texts)
        if "Incidencias con discrepancias en estado finalista" in text
    )
    assert root_idx < finalist_idx
    assert "EAM-ROOT" in " ".join(slide_texts)

    finalist_slide = next(
        slide
        for slide in prs.slides
        if "Incidencias con discrepancias en estado finalista" in _slide_all_text(slide)
        and "EAM-REG" in _slide_all_text(slide)
    )
    assert "EAM-ROOT" not in _slide_all_text(finalist_slide)

    root_slide = next(
        slide
        for slide in prs.slides
        if "Evolutivos para solucionar causas raíces" in _slide_all_text(slide)
        and "EAM-ROOT" in _slide_all_text(slide)
    )
    root_text = _slide_all_text(root_slide)
    assert "Comentarios registrados" in root_text
    assert "Cliente confirma evolutivo priorizado" in root_text


def test_period_followup_ppt_links_helix_ids_in_risk_tables_from_full_dataset(
    tmp_path: Path,
) -> None:
    template = tmp_path / "template.pptx"
    data_path = tmp_path / "issues.json"
    _build_minimal_template(template)
    save_issues_doc(
        str(data_path),
        IssuesDocument(
            issues=[
                NormalizedIssue(
                    key="MEX-AGED-CLOSED",
                    summary=(
                        "INC000102885426 - liquidez / bbva net cash / cuenta control / "
                        "caso con mas de 30 dias"
                    ),
                    description="",
                    status="En progreso",
                    type="Historia",
                    priority="Medium",
                    assignee="Ana",
                    created="2026-03-01T00:00:00Z",
                    updated="2026-05-20T00:00:00Z",
                    country="México",
                    source_type="jira",
                    source_id="jira:mexico:core",
                    source_alias="Core",
                    url="https://jira.example.com/browse/MEX-AGED-CLOSED",
                ),
                NormalizedIssue(
                    key="INC000102885426",
                    summary="Cerrado en Helix",
                    description="",
                    status="Closed",
                    type="Helix",
                    priority="Medium",
                    created="2026-03-01T00:00:00Z",
                    updated="2026-05-20T00:00:00Z",
                    resolved="2026-05-20T00:00:00Z",
                    country="México",
                    source_type="helix",
                    source_id="helix:mexico:lookup-estados-finalistas-jira",
                    source_alias="Lookup estados finalistas Jira",
                    helix_lookup_kind="post_jql_inc_lookup",
                    url="https://helix.example.com/smartit/app/#/incidentPV/IDG102885426",
                ),
                NormalizedIssue(
                    key="MEX-AGED-OPEN",
                    summary="INC000104451980 - caso Helix localizado y aun abierto",
                    description="",
                    status="En progreso",
                    type="Historia",
                    priority="Medium",
                    assignee="Ana",
                    created="2026-03-01T00:00:00Z",
                    updated="2026-05-20T00:00:00Z",
                    country="México",
                    source_type="jira",
                    source_id="jira:mexico:core",
                    source_alias="Core",
                    url="https://jira.example.com/browse/MEX-AGED-OPEN",
                ),
                NormalizedIssue(
                    key="INC000104451980",
                    summary="En curso en Helix",
                    description="",
                    status="Assigned",
                    type="Helix",
                    priority="Medium",
                    created="2026-03-01T00:00:00Z",
                    updated="2026-05-20T00:00:00Z",
                    country="México",
                    source_type="helix",
                    source_id="helix:mexico:lookup-estados-finalistas-jira",
                    source_alias="Lookup estados finalistas Jira",
                    helix_lookup_kind="post_jql_inc_lookup",
                    url="https://helix.example.com/smartit/app/#/incident/IDG104451980",
                ),
            ]
        ),
    )

    result = generate_country_period_followup_ppt(
        Settings(DATA_PATH=str(data_path), PERIOD_PPT_TEMPLATE_PATH=str(template)),
        country="México",
        source_ids=["jira:mexico:core"],
        reference_day=pd.Timestamp("2026-05-24T00:00:00Z"),
    )
    prs = Presentation(BytesIO(result.content))
    aged_detail_text = "\n".join(
        _slide_all_text(slide)
        for slide in prs.slides
        if "Incidencias abiertas con más de 30 días (" in _slide_all_text(slide)
    )
    assert "MEX-AGED-CLOSED" not in aged_detail_text
    aged_slide = next(
        slide
        for slide in prs.slides
        if "Incidencias abiertas con más de 30 días" in _slide_all_text(slide)
        and "MEX-AGED-OPEN" in _slide_all_text(slide)
    )
    table = _native_tables(aged_slide)[0].table
    description_runs = [
        run for paragraph in table.cell(1, 1).text_frame.paragraphs for run in paragraph.runs
    ]

    assert any(
        "INC000104451980" in str(run.text or "")
        and str(run.hyperlink.address or "")
        == "https://helix.example.com/smartit/app/#/incident/IDG104451980"
        for run in description_runs
    )


def test_generate_country_period_followup_ppt_uses_open_focus_label_from_settings() -> None:
    template = bundled_period_ppt_template_path()
    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )

    settings_critical = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(template),
        OPEN_ISSUES_FOCUS_MODE="criticidad_alta",
    )
    out_critical = generate_country_period_followup_ppt(
        settings_critical,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs_critical = Presentation(BytesIO(out_critical.content))
    critical_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs_critical.slides[2].shapes
        if getattr(shape, "has_text_frame", False)
    ).upper()
    assert "CRITICIDAD ALTA" in critical_blob
    assert "CRITICIDADES ALTAS" in critical_blob
    assert "RESTO" in critical_blob

    settings_maestras = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(template),
        OPEN_ISSUES_FOCUS_MODE="maestras",
    )
    out_maestras = generate_country_period_followup_ppt(
        settings_maestras,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs_maestras = Presentation(BytesIO(out_maestras.content))
    maestras_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs_maestras.slides[2].shapes
        if getattr(shape, "has_text_frame", False)
    ).upper()
    assert "INCIDENCIAS MAESTRAS" in maestras_blob
    assert "MAESTRAS" in maestras_blob
    assert "RESTO" in maestras_blob


def test_generate_country_period_followup_ppt_bundled_template_layout_regression() -> None:
    template = bundled_period_ppt_template_path()
    assert template.exists()

    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "A-2",
                "summary": "Issue A2",
                "status": "Resolved",
                "priority": "Low",
                "created": (now - pd.Timedelta(days=8)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=2)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(template))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now - pd.Timedelta(days=1),
    )

    prs = Presentation(BytesIO(out.content))

    s3_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs.slides[2].shapes
        if getattr(shape, "has_text_frame", False)
    ).lower()
    assert "seguimiento de incidencias - méxico (vista agregada)" in s3_blob
    for slide_idx in (2, 3, 4):
        summary_slide = prs.slides[slide_idx]
        assert summary_slide.shapes[4].fill.fore_color.rgb == RGBColor(
            *hex_to_rgb(BBVA_REPORT_RED_BG)
        )
        assert summary_slide.shapes[5].fill.fore_color.rgb == RGBColor(
            *hex_to_rgb(BBVA_REPORT_AMBER_BG)
        )

    # Regression guard: redesigned slides 7/8 keep a single hero chart panel.
    for slide_idx in (6, 7):  # slides 7 and 8 (0-based indexes)
        slide = prs.slides[slide_idx]
        pic_shapes = []
        for shape in slide.shapes:
            try:
                _ = shape.image
            except Exception:
                continue
            area_in2 = float(shape.width) * float(shape.height) / (914400.0 * 914400.0)
            if area_in2 >= 1.0:
                pic_shapes.append(shape)
        assert len(pic_shapes) == 1

    s7_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs.slides[6].shapes
        if getattr(shape, "has_text_frame", False)
    ).lower()
    assert "visión agregada de incidencias abiertas : rango de días por prioridad" in s7_blob
    assert "insights accionables" not in s7_blob

    s8_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs.slides[7].shapes
        if getattr(shape, "has_text_frame", False)
    ).lower()
    assert "visión agregada de incidencias abiertas por prioridad" in s8_blob
    assert "total abiertas" not in s8_blob
    assert "prioridad dominante" not in s8_blob
    assert "riesgo ponderado" not in s8_blob
    assert "insights accionables" not in s8_blob

    # Regression guard: summary metrics should not concatenate duplicated labels.
    s4 = prs.slides[3]
    s4_closed = str(s4.shapes[8].text or "").upper().replace(" ", "")
    s4_days = str(s4.shapes[11].text or "").upper().replace(" ", "")
    s4_blob = " ".join(
        str(getattr(shape, "text", "") or "")
        for shape in s4.shapes
        if getattr(shape, "has_text_frame", False)
    ).upper()
    assert "CERRADASINCIDENCIA" not in s4_closed
    assert "RESOLUCIÓNDÍASDERESOLUCIÓN" not in s4_days
    assert s4_blob.count("CRITICIDADES ALTAS") == 1
    assert "0 CRITICIDADES ALTAS" in s4_blob
    assert "1 RESTO" in s4_blob
    assert "7 DÍAS MAX" in s4_blob
    assert "7 DÍAS MIN" in s4_blob

    # Long titles should be marked for in-shape fit in PowerPoint.
    s5_title = prs.slides[4].shapes[2]
    assert s5_title.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Cover period must be rendered only in the dedicated period placeholder.
    cover_blob = " || ".join(
        str(getattr(shape, "text", "") or "")
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    cover_blob_lower = cover_blob.lower()
    assert "periodo dd/mm - dd/mm 2026" not in cover_blob_lower
    period_matches = re.findall(
        r"periodo\s+\d{2}/\d{2}\s*-\s*\d{2}/\d{2}/\d{4}",
        cover_blob_lower,
    )
    assert len(period_matches) == 1
    assert "seguimiento incidencias" in cover_blob_lower
    assert "kpis, evolución y análisis del periodo" not in cover_blob_lower

    period_shape = next(
        (
            shape
            for shape in prs.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
            and re.search(
                r"periodo\s+\d{2}/\d{2}\s*-\s*\d{2}/\d{2}/\d{4}",
                str(getattr(shape, "text", "") or "").lower(),
            )
        ),
        None,
    )
    assert period_shape is not None
    period_run_size = None
    for paragraph in period_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if str(getattr(run, "text", "") or "").strip():
                period_run_size = run.font.size
                break
        if period_run_size is not None:
            break
    assert period_run_size is not None
    assert float(period_run_size.pt) >= 10.5
    assert period_shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert period_shape.text_frame.word_wrap is False
    assert int(period_shape.text_frame.margin_left) > 0
    assert int(period_shape.text_frame.margin_right) > 0


def test_resolution_chart_uses_executive_fonts_and_column_totals(monkeypatch: Any) -> None:
    now = pd.Timestamp("2026-04-30T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=1)).isoformat(),
            },
            {
                "key": "A-2",
                "summary": "Issue B",
                "status": "New",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=1)).isoformat(),
            },
            {
                "key": "A-3",
                "summary": "Issue C",
                "status": "New",
                "priority": "Highest",
                "created": (now - pd.Timedelta(days=45)).isoformat(),
            },
        ]
    )
    captured: dict[str, Any] = {}

    def _capture_fig(fig: Any, *, width: int, height: int, scale: float = 1.0) -> bytes:
        captured["fig"] = fig
        captured["width"] = width
        captured["height"] = height
        captured["scale"] = scale
        return b"png"

    monkeypatch.setattr(period_ppt_mod, "_fig_to_png_exact", _capture_fig)

    payload = period_ppt_mod._resolution_chart_png_executive(
        Settings(),
        dff=dff,
        open_df=dff,
        reference_now=now,
    )

    assert payload == b"png"
    assert captured["width"] == EXEC_CHART_EXPORT_WIDTH
    assert captured["height"] == EXEC_CHART_EXPORT_HEIGHT
    fig = captured["fig"]
    assert int(fig.layout.xaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.yaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.legend.font.size) >= EXEC_CHART_LEGEND_FONT_PT
    bars = [trace for trace in fig.data if str(trace.type) == "bar"]
    assert bars
    assert all(int(trace.textfont.size) >= EXEC_CHART_INSIDE_VALUE_FONT_PT for trace in bars)
    totals_by_x: dict[str, int] = {}
    for trace in bars:
        for x_val, y_val in zip(list(trace.x), list(trace.y)):
            totals_by_x[str(x_val)] = totals_by_x.get(str(x_val), 0) + int(y_val)
    total_trace = [
        trace
        for trace in fig.data
        if str(trace.type) == "scatter" and list(getattr(trace, "text", []) or [])
    ][-1]
    assert int(total_trace.textfont.size) >= EXEC_CHART_TOTAL_FONT_PT
    assert {
        str(x_val): int(text)
        for x_val, text in zip(list(total_trace.x), list(total_trace.text))
        if str(text).strip()
    } == {label: total for label, total in totals_by_x.items() if total > 0}


def test_priority_chart_uses_executive_fonts_and_safe_y_range(monkeypatch: Any) -> None:
    dff = pd.DataFrame(
        [
            {"key": "A-1", "status": "New", "priority": "High"},
            {"key": "A-2", "status": "New", "priority": "High"},
            {"key": "A-3", "status": "New", "priority": "Medium"},
        ]
    )
    captured: dict[str, Any] = {}

    def _capture_fig(fig: Any, *, width: int, height: int, scale: float = 1.0) -> bytes:
        captured["fig"] = fig
        captured["width"] = width
        captured["height"] = height
        return b"png"

    monkeypatch.setattr(period_ppt_mod, "_fig_to_png_exact", _capture_fig)

    payload = period_ppt_mod._priority_chart_png_executive(Settings(), dff=dff, open_df=dff)

    assert payload == b"png"
    assert captured["width"] == EXEC_CHART_EXPORT_WIDTH
    assert captured["height"] == EXEC_CHART_EXPORT_HEIGHT
    fig = captured["fig"]
    bars = [trace for trace in fig.data if str(trace.type) == "bar"]
    assert bars
    assert int(fig.layout.xaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.yaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert all(int(trace.textfont.size) >= EXEC_CHART_INSIDE_VALUE_FONT_PT for trace in bars)
    pct_trace = [trace for trace in fig.data if str(trace.type) == "scatter"][-1]
    assert int(pct_trace.textfont.size) >= EXEC_CHART_TOTAL_FONT_PT
    max_bar = max(int(trace.y[0]) for trace in bars)
    assert float(fig.layout.yaxis.range[1]) > float(max_bar)


def test_functionality_trend_filters_last_six_months_and_adds_totals(
    monkeypatch: Any,
) -> None:
    rows: list[dict[str, object]] = []
    for idx, created in enumerate(
        [
            "2026-01-05T00:00:00+00:00",
            "2026-02-05T00:00:00+00:00",
            "2026-03-05T00:00:00+00:00",
            "2026-04-05T00:00:00+00:00",
            "2026-05-05T00:00:00+00:00",
            "2026-06-05T00:00:00+00:00",
            "2026-07-05T00:00:00+00:00",
            "2026-08-05T00:00:00+00:00",
        ],
        start=1,
    ):
        rows.append(
            {
                "key": f"A-{idx}",
                "summary": "Pagos no refleja saldo" if idx % 2 else "Login falla",
                "status": "New",
                "priority": "High",
                "created": created,
            }
        )
    captured: dict[str, Any] = {}

    def _capture_fig(fig: Any, *, width: int, height: int, scale: float = 1.0) -> bytes:
        captured["fig"] = fig
        captured["width"] = width
        captured["height"] = height
        return b"png"

    monkeypatch.setattr(period_ppt_mod, "_fig_to_png_exact", _capture_fig)

    payload = period_ppt_mod._functionality_fortnight_trend_png(open_df=pd.DataFrame(rows))

    assert payload == b"png"
    assert captured["width"] == EXEC_CHART_EXPORT_WIDTH
    assert captured["height"] == EXEC_CHART_TREND_EXPORT_HEIGHT
    fig = captured["fig"]
    bars = [trace for trace in fig.data if str(trace.type) == "bar"]
    assert bars
    axis_labels = [str(label) for label in list(bars[0].x)]
    assert all(not label.startswith(("01 |", "02 |")) for label in axis_labels)
    assert any(label.startswith("03 |") for label in axis_labels)
    assert int(fig.layout.xaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.yaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.legend.font.size) >= EXEC_CHART_LEGEND_FONT_PT
    assert all(int(trace.textfont.size) >= EXEC_CHART_INSIDE_VALUE_FONT_PT for trace in bars)
    totals_by_x: dict[str, int] = {}
    for trace in bars:
        for x_val, y_val in zip(list(trace.x), list(trace.y)):
            totals_by_x[str(x_val)] = totals_by_x.get(str(x_val), 0) + int(y_val)
    total_trace = [
        trace
        for trace in fig.data
        if str(trace.type) == "scatter" and list(getattr(trace, "text", []) or [])
    ][-1]
    assert int(total_trace.textfont.size) >= EXEC_CHART_TOTAL_FONT_PT
    assert {
        str(x_val): int(text)
        for x_val, text in zip(list(total_trace.x), list(total_trace.text))
        if str(text).strip()
    } == {label: total for label, total in totals_by_x.items() if total > 0}


def test_filter_last_six_months_trend_keeps_month_window() -> None:
    trend = pd.DataFrame(
        {
            "quincena_start": pd.to_datetime(
                [
                    "2026-01-01",
                    "2026-02-01",
                    "2026-03-01",
                    "2026-04-01",
                    "2026-05-01",
                    "2026-06-01",
                    "2026-07-01",
                    "2026-08-01",
                ]
            ),
            "quincena_end": pd.to_datetime(
                [
                    "2026-01-15",
                    "2026-02-15",
                    "2026-03-15",
                    "2026-04-15",
                    "2026-05-15",
                    "2026-06-15",
                    "2026-07-15",
                    "2026-08-15",
                ]
            ),
            "quincena_label": [f"2026-{month:02d} · 1-15" for month in range(1, 9)],
            "tema": ["Pagos"] * 8,
            "issues": [1] * 8,
            "issues_cumulative": list(range(1, 9)),
            "issues_value": list(range(1, 9)),
        }
    )

    filtered = period_ppt_mod._filter_last_six_months_trend(trend)

    assert filtered["quincena_start"].min() == pd.Timestamp("2026-03-01")
    assert filtered["quincena_start"].max() == pd.Timestamp("2026-08-01")
    assert filtered["quincena_start"].tolist() == sorted(filtered["quincena_start"].tolist())


def test_summary_timeseries_chart_uses_executive_export_tokens(monkeypatch: Any) -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=20)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
            },
            {
                "key": "A-2",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=18)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=2)).isoformat(),
            },
            {
                "key": "A-3",
                "summary": "Issue C",
                "status": "New",
                "priority": "Low",
                "created": (now - pd.Timedelta(days=5)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
            },
        ]
    )
    captured: dict[str, Any] = {}

    def _capture_fig(fig: Any, *, width: int, height: int, scale: float = 1.0) -> bytes:
        captured["fig"] = fig
        captured["width"] = width
        captured["height"] = height
        captured["scale"] = scale
        return b"png"

    monkeypatch.setattr(period_ppt_mod, "_fig_to_png_exact", _capture_fig)

    payload = period_ppt_mod._chart_png(
        Settings(),
        dff=dff,
        open_df=dff[dff["resolved"].isna()],
        chart_id="timeseries",
    )

    assert payload == b"png"
    assert captured["width"] == EXEC_CHART_EXPORT_WIDTH
    assert captured["height"] == EXEC_CHART_TREND_EXPORT_HEIGHT
    fig = captured["fig"]
    assert int(fig.layout.xaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.yaxis.tickfont.size) >= EXEC_CHART_AXIS_FONT_PT
    assert int(fig.layout.legend.font.size) >= EXEC_CHART_LEGEND_FONT_PT
    assert int(fig.layout.margin.b) >= 190
    for trace in fig.data:
        if str(getattr(trace, "type", "") or "").lower() in {"scatter", "scattergl"}:
            assert float(trace.marker.size) >= 8.0
            assert float(trace.line.width) >= 4.2


def test_summary_timeseries_chart_uses_slide_native_ratio_and_readable_fonts(
    monkeypatch: Any,
) -> None:
    now = pd.Timestamp("2026-05-30T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": f"A-{idx}",
                "summary": "Issue",
                "status": "Resolved" if idx % 3 == 0 else "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=idx)).isoformat(),
                "resolved": (
                    (now - pd.Timedelta(days=max(idx - 2, 0))).isoformat() if idx % 3 == 0 else None
                ),
            }
            for idx in range(1, 45)
        ]
    )
    captured: dict[str, Any] = {}

    def _capture_fig(fig: Any, *, width: int, height: int, scale: float = 1.0) -> bytes:
        captured["fig"] = fig
        captured["width"] = width
        captured["height"] = height
        captured["scale"] = scale
        return b"png"

    monkeypatch.setattr(period_ppt_mod, "_fig_to_png_exact", _capture_fig)

    payload = period_ppt_mod._chart_png(
        Settings(),
        dff=dff,
        open_df=dff[dff["resolved"].isna()],
        chart_id="timeseries",
        width=1200,
        height=838,
        slide_optimized=True,
    )

    assert payload == b"png"
    assert captured["width"] == 1200
    assert captured["height"] == 838
    assert captured["scale"] == 1.0
    fig = captured["fig"]
    assert int(fig.layout.xaxis.tickfont.size) >= int(EXEC_CHART_AXIS_FONT_PT * 1.35)
    assert int(fig.layout.yaxis.tickfont.size) >= int(EXEC_CHART_AXIS_FONT_PT * 1.35)
    assert int(fig.layout.legend.font.size) >= int(EXEC_CHART_LEGEND_FONT_PT * 1.35)
    assert int(fig.layout.xaxis.nticks) <= 5
    assert int(fig.layout.xaxis.tickangle) <= -30
    assert int(fig.layout.margin.b) >= 260
    for trace in fig.data:
        if str(getattr(trace, "type", "") or "").lower() in {"scatter", "scattergl"}:
            assert float(trace.marker.size) >= 11.0
            assert float(trace.line.width) >= 6.4


def test_generate_country_period_followup_ppt_uses_timeseries_for_summary(
    monkeypatch: Any,
) -> None:
    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "Medium",
                "created": (now - pd.Timedelta(days=10)).isoformat(),
                "updated": now.isoformat(),
                "resolved": (now - pd.Timedelta(days=1)).isoformat(),
                "country": "México",
                "source_id": "jira:mexico:gema",
                "source_type": "jira",
            },
        ]
    )
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        JIRA_BASE_URL="https://jira.example",
        PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED="true",
    )

    called_chart_ids: list[str] = []
    tiny_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO5N0nQAAAAASUVORK5CYII="
    )

    def _fake_chart_png(
        _settings: Settings,
        *,
        dff: pd.DataFrame,
        open_df: pd.DataFrame,
        chart_id: str,
        **_kwargs: object,
    ) -> bytes:
        _ = (dff, open_df)
        called_chart_ids.append(str(chart_id))
        return tiny_png

    monkeypatch.setattr(period_ppt_mod, "_chart_png", _fake_chart_png)

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    assert out.content
    assert called_chart_ids[:3] == ["timeseries", "timeseries", "timeseries"]
    # Redesigned slides 7/8 render with dedicated executive chart builders and
    # no longer call the generic _chart_png pipeline.
    assert called_chart_ids == ["timeseries", "timeseries", "timeseries"]


def test_period_followup_ppt_handles_three_rollup_sources(monkeypatch: Any, tmp_path: Path) -> None:
    _ = tmp_path
    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    rows = []
    for idx, (source_id, alias) in enumerate(
        (
            ("jira:mexico:senda", "Senda"),
            ("jira:mexico:gema", "Gema"),
            ("jira:mexico:core", "Core"),
        ),
        start=1,
    ):
        rows.append(
            {
                "key": f"MEX-{idx}",
                "summary": f"Issue {alias}",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=idx)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": source_id,
                "source_type": "jira",
            }
        )
    dff = pd.DataFrame(rows)
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        JIRA_SOURCES_JSON=(
            '[{"country":"México","alias":"Senda","jql":"project = SENDA"},'
            '{"country":"México","alias":"Gema","jql":"project = GEMA"},'
            '{"country":"México","alias":"Core","jql":"project = CORE"}]'
        ),
    )
    chart_scopes: list[tuple[str, ...]] = []

    def _fake_chart_png(
        _settings: Settings,
        *,
        dff: pd.DataFrame,
        open_df: pd.DataFrame,
        chart_id: str,
        **_kwargs: object,
    ) -> bytes:
        _ = (open_df, chart_id)
        chart_scopes.append(tuple(sorted(dff["source_id"].dropna().astype(str).unique())))
        return b""

    monkeypatch.setattr(period_ppt_mod, "_chart_png", _fake_chart_png)
    monkeypatch.setattr(
        period_ppt_mod, "_append_functionality_followup_slides", lambda *a, **k: None
    )
    monkeypatch.setattr(period_ppt_mod, "_append_period_risk_issue_sections", lambda *a, **k: None)
    monkeypatch.setattr(
        period_ppt_mod, "_populate_open_aging_executive_slide", lambda *a, **k: None
    )
    monkeypatch.setattr(
        period_ppt_mod, "_populate_open_priority_executive_slide", lambda *a, **k: None
    )
    monkeypatch.setattr(period_ppt_mod, "validate_shapes_inside_slide", lambda *a, **k: None)

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema", "jira:mexico:core"],
        dff_override=dff,
        reference_day=now,
    )

    prs = Presentation(BytesIO(out.content))
    full_text = " ".join(_slide_all_text(slide) for slide in prs.slides)
    assert out.source_ids == ("jira:mexico:senda", "jira:mexico:gema", "jira:mexico:core")
    assert out.slide_count == 9
    assert "SENDA" in full_text
    assert "GEMA" in full_text
    assert "CORE" in full_text
    assert chart_scopes == [
        ("jira:mexico:core", "jira:mexico:gema", "jira:mexico:senda"),
        ("jira:mexico:senda",),
        ("jira:mexico:gema",),
        ("jira:mexico:core",),
    ]


def test_period_followup_ppt_without_rollups_omits_aggregate_and_source_summary_slides(
    monkeypatch: Any, tmp_path: Path
) -> None:
    template = tmp_path / "template.pptx"
    _build_compact_template(template)
    now = pd.Timestamp("2026-03-15T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "MEX-1",
                "summary": "Issue país",
                "status": "New",
                "priority": "High",
                "created": (now - pd.Timedelta(days=2)).isoformat(),
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "source_type": "jira",
            }
        ]
    )
    monkeypatch.setattr(
        period_ppt_mod, "_append_functionality_followup_slides", lambda *a, **k: None
    )
    monkeypatch.setattr(period_ppt_mod, "_append_period_risk_issue_sections", lambda *a, **k: None)
    monkeypatch.setattr(
        period_ppt_mod, "_populate_open_aging_executive_slide", lambda *a, **k: None
    )
    monkeypatch.setattr(
        period_ppt_mod, "_populate_open_priority_executive_slide", lambda *a, **k: None
    )

    out = generate_country_period_followup_ppt(
        Settings(PERIOD_PPT_TEMPLATE_PATH=str(template)),
        country="México",
        source_ids=[],
        dff_override=dff,
        reference_day=now,
    )

    prs = Presentation(BytesIO(out.content))
    full_text = " ".join(_slide_all_text(slide) for slide in prs.slides)
    assert out.source_ids == ()
    assert out.slide_count == 5
    assert "vista agregada" not in full_text
    assert "Seguimiento de incidencias - SENDA" not in full_text


def test_period_followup_ppt_renders_po_under_assignee() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    table_shape = period_ppt_mod._populate_issue_native_table(
        slide,
        table_shape_index=0,
        headers=period_ppt_mod._RISK_ASSIGNEE_TABLE_HEADERS,
        rows=[
            [
                "MEX-1",
                "Issue",
                period_ppt_mod._assignee_with_po_text(
                    "MARCELA FONSECA MONTEALEGRE", "Víctor Expósito"
                ),
                "Open",
                "High",
                "5 días",
            ]
        ],
    )
    assignee_cell = table_shape.table.cell(1, 2)

    assert assignee_cell.text == "MARCELA FONSECA MONTEALEGRE\n(Víctor Expósito)"
    assert len(assignee_cell.text_frame.paragraphs) >= 2
    first_run = assignee_cell.text_frame.paragraphs[0].runs[0]
    po_run = assignee_cell.text_frame.paragraphs[1].runs[0]
    assert po_run.font.size.pt < first_run.font.size.pt
    assert period_ppt_mod._assignee_with_po_text("MARCELA", "") == "MARCELA"
    long_po = period_ppt_mod._assignee_with_po_text(
        "RODRIGO GALLEGOS SUAREZ",
        "Juan Vicente Guerrero con un nombre extremadamente largo para tabla",
        po_max_chars=24,
    )
    assert long_po.startswith("RODRIGO GALLEGOS SUAREZ\n(")
    assert long_po.endswith("…)")


def test_period_followup_risk_issue_table_renders_notes_as_grouped_row() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows, row_links, comment_by_row = period_ppt_mod._risk_issue_rows_for_table(
        [
            period_ppt_mod.PeriodRiskIssueRow(
                key="MEXBMI1-283305",
                summary="Tras login la app queda con spinner",
                functionality="Acceso",
                assignee="Rodrigo",
                status="Ready To Verify",
                priority="High",
                open_days=171,
                url="https://jira.example.com/browse/MEXBMI1-283305",
                po_team_leader="Juan Vicente",
            )
        ],
        empty_message="Sin incidencias",
        notes_by_key={
            "MEXBMI1-283305": (
                "Juan Vicente indica que debería haber sido descartada y que la va a revisar"
            )
        },
    )
    table_shape = period_ppt_mod._populate_issue_native_table(
        slide,
        table_shape_index=0,
        headers=period_ppt_mod._RISK_ASSIGNEE_TABLE_HEADERS,
        rows=rows,
        hyperlink_by_row=row_links,
    )
    period_ppt_mod._style_issue_comment_rows(table_shape, comment_by_row=comment_by_row)

    table = table_shape.table
    assert len(table.rows) == 3
    assert table.cell(1, 0).text == "MEXBMI1-283305"
    assert "Comentarios registrados" in table.cell(2, 1).text
    assert "debería haber sido descartada" in table.cell(2, 1).text
    assert str(table.cell(1, 0).text_frame.paragraphs[0].runs[0].hyperlink.address).startswith(
        "https://jira.example.com"
    )


def test_period_followup_ppt_enriches_po_from_source_config() -> None:
    df = pd.DataFrame(
        [
            {
                "key": "MEX-1",
                "summary": "Issue",
                "status": "Open",
                "priority": "High",
                "created": "2026-05-01T00:00:00Z",
                "source_id": "jira:mexico:core",
            }
        ]
    )
    out = period_ppt_mod._enrich_po_team_leader_from_sources(
        df,
        Settings(
            JIRA_SOURCES_JSON=(
                '[{"source_id":"jira:mexico:core","country":"México","alias":"Core",'
                '"po_team_leader":"Juan Vicente Guerrero","jql":"project = CORE"}]'
            )
        ),
    )

    assert out.loc[0, "po_team_leader"] == "Juan Vicente Guerrero"


def test_generate_country_period_followup_ppt_zoom_table_matches_issue_count() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Login falla en acceso de usuario",
                "status": "New",
                "priority": "High",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Login falla en acceso biometrico",
                "status": "Ready To Verify",
                "priority": "Highest",
                "created": "2026-04-03T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "B-1",
                "summary": "TAREAS PENDIENTES - No se visualiza dashboard",
                "status": "New",
                "priority": "High",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        JIRA_BASE_URL="https://jira.example",
        PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED="true",
    )
    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))
    dashboard_slide = prs.slides[
        _find_slide_index(prs, "Seguimiento de KPIs - Incidencias abiertas por funcionalidad")
    ]
    dashboard_tables = _native_tables(dashboard_slide)
    assert len(dashboard_tables) == 1
    assert "Resto incidencias abiertas" in _slide_table_text(dashboard_slide)
    assert not _table_intersects_picture(dashboard_slide, dashboard_tables[0])
    _assert_native_table_font_floor(dashboard_tables[0])

    first_zoom_idx = _find_slide_index(prs, "Incidencias, en Login y acceso, abiertas")
    for slide_idx in (first_zoom_idx, first_zoom_idx + 1, first_zoom_idx + 2):
        slide = prs.slides[slide_idx]
        zoom_tables = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
        assert len(zoom_tables) == 1

    zoom_table = [
        shape for shape in prs.slides[first_zoom_idx].shapes if getattr(shape, "has_table", False)
    ][0].table
    first_data_key_cell = zoom_table.cell(1, 0)
    runs = list(first_data_key_cell.text_frame.paragraphs[0].runs)
    assert runs
    assert str(runs[0].hyperlink.address or "").startswith("https://")
    assert first_data_key_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT
    assert first_data_key_cell.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    header_criticity_cell = zoom_table.cell(0, 4)
    header_runs = list(header_criticity_cell.text_frame.paragraphs[0].runs)
    assert header_runs
    assert float(header_runs[0].font.size.pt) >= 8.0
    tc_pr = first_data_key_cell._tc.tcPr
    assert tc_pr.find(qn("a:lnL")) is not None
    assert tc_pr.find(qn("a:lnR")) is not None
    assert tc_pr.find(qn("a:lnT")) is not None
    assert tc_pr.find(qn("a:lnB")) is not None


def test_generate_country_period_followup_ppt_top3_lines_include_avg_days() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Pagos no refleja saldo",
                "status": "New",
                "priority": "High",
                "created": "2026-03-11T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Pagos timeout intermitente",
                "status": "New",
                "priority": "Medium",
                "created": "2026-03-21T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "B-1",
                "summary": "Transferencias fallan",
                "status": "Ready To Verify",
                "priority": "High",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED="true",
    )
    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))
    dashboard_slide = prs.slides[
        _find_slide_index(prs, "Seguimiento de KPIs - Incidencias abiertas por funcionalidad")
    ]
    top_three_blob = _slide_text(dashboard_slide).lower()
    assert "en total" in top_three_blob
    assert "días promedio" in top_three_blob
    assert "acum." not in top_three_blob
    assert "d. p." not in top_three_blob


def test_generate_country_period_followup_ppt_functionality_color_contrast_is_readable() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Pago no refleja movimiento",
                "status": "New",
                "priority": "High",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Transferencias en tiempo real fallan",
                "status": "Blocked",
                "priority": "Medium",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED="true",
    )

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))

    dashboard_idx = _find_slide_index(
        prs, "Seguimiento de KPIs - Incidencias abiertas por funcionalidad"
    )
    dashboard_blob_shape = next(
        shape
        for shape in prs.slides[dashboard_idx].shapes
        if getattr(shape, "has_text_frame", False)
        and "INCIDENCIAS" in str(getattr(shape, "text", "") or "")
        and "ABIERTAS" in str(getattr(shape, "text", "") or "")
    )
    dashboard_blob_text = str(getattr(dashboard_blob_shape, "text", "") or "")
    assert "|" not in dashboard_blob_text
    assert "INCIDENCIAS ABIERTAS" in dashboard_blob_text.replace("\n", " ")
    dashboard_run = dashboard_blob_shape.text_frame.paragraphs[0].runs[0]
    assert dashboard_run.font.color.rgb == RGBColor(255, 255, 255)

    dashboard_tables = _native_tables(prs.slides[dashboard_idx])
    assert len(dashboard_tables) == 1
    dashboard_table = dashboard_tables[0]
    assert not _table_intersects_picture(prs.slides[dashboard_idx], dashboard_table)
    _assert_native_table_font_floor(dashboard_table)
    assert int(dashboard_table.top) > int(dashboard_blob_shape.top + dashboard_blob_shape.height)
    mitigation_panel = next(
        (
            shape
            for shape in prs.slides[dashboard_idx].shapes
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE
            and int(getattr(shape, "left", 0)) > 5_400_000
            and int(getattr(shape, "width", 0)) >= 3_000_000
            and int(getattr(shape, "height", 0)) >= 2_000_000
        ),
        None,
    )
    assert mitigation_panel is not None
    dashboard_text = _slide_all_text(prs.slides[dashboard_idx])
    assert "Estado Ready to Verify" in dashboard_text
    assert "Estado New" in dashboard_text
    assert "Estado bloqueadas" in dashboard_text
    assert "Resto:" in dashboard_text
    assert "d. prom." not in dashboard_text
    assert "Incidencias en New:" not in dashboard_text
    assert "Resto de incidencias:" not in dashboard_text
    assert int(dashboard_table.left + dashboard_table.width) < int(mitigation_panel.left)
    assert int(dashboard_table.top + dashboard_table.height) <= int(
        mitigation_panel.top + mitigation_panel.height
    )

    first_zoom_idx = _find_slide_index(prs, "Incidencias, en Pagos, abiertas")
    root_cause_shape = next(
        shape
        for shape in prs.slides[first_zoom_idx].shapes
        if getattr(shape, "has_text_frame", False)
        and any(
            token in str(getattr(shape, "text", "") or "").lower()
            for token in ("causa", "sin incidencias")
        )
    )
    root_run = root_cause_shape.text_frame.paragraphs[0].runs[0]
    assert root_run.font.color.rgb == RGBColor(*period_ppt_mod._TABLE_BODY_FG_RGB)


def test_generate_country_period_followup_ppt_zoom_paginates_when_overflow() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    rows: list[dict[str, object]] = []
    for idx in range(7):
        rows.append(
            {
                "key": f"P-{idx + 1}",
                "summary": (
                    f"INC0001 - PAGOS / SENDA BNC / TRANSFERENCIAS EN TIEMPO REAL / CASO {idx + 1}"
                ),
                "status": "New",
                "priority": "Medium",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            }
        )
    rows.extend(
        [
            {
                "key": "M-1",
                "summary": "Saldo monetarias no actualizado",
                "status": "Analysing",
                "priority": "Medium",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "T-1",
                "summary": "Transferencias con timeout intermitente",
                "status": "Blocked",
                "priority": "Medium",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    dff = pd.DataFrame(rows)
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        PERIOD_REPORT_FUNCTIONALITY_DETAIL_ENABLED="true",
    )
    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))
    assert len(prs.slides) == 19
    first_zoom_idx = _find_slide_index(prs, "Incidencias, en Pagos, abiertas en la quincena (I)")
    zoom_titles = [
        str(getattr(shape, "text", "") or "").strip()
        for slide_idx in (first_zoom_idx, first_zoom_idx + 1)
        for shape in prs.slides[slide_idx].shapes
        if getattr(shape, "has_text_frame", False)
    ]
    joined_titles = " | ".join(zoom_titles)
    deck_text = " ".join(_slide_text(slide) for slide in prs.slides)
    assert "Incidencias abiertas por criticidad alta" in deck_text
    assert "Incidencias abiertas con más de 30 días" in deck_text
    assert "Incidencias, en Pagos, abiertas en la quincena (I)" in joined_titles
    assert "Incidencias, en Pagos, abiertas en la quincena (II)" in joined_titles


def test_period_followup_risk_sections_use_native_tables_after_functionality() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "MEXBMI1-101",
                "summary": "Alta criticidad con bloqueo de operativa",
                "status": "New",
                "priority": "High",
                "assignee": "Luis Pérez",
                "created": "2026-02-20T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
                "url": "https://jira.example/browse/MEXBMI1-101",
            },
            {
                "key": "EAM-77",
                "summary": "Impedimento en autenticación para clientes",
                "status": "Ready To Verify",
                "priority": "Supone un impedimento",
                "assignee": "Ana López",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "SKSEMEX-9",
                "summary": "Incidencia abierta antigua de prioridad media",
                "status": "Analysing",
                "priority": "Medium",
                "assignee": "Marta Ruiz",
                "created": "2026-02-01T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "MEXBMI1-102",
                "summary": "Alta criticidad adicional 1",
                "status": "New",
                "priority": "High",
                "assignee": "Equipo Core",
                "created": "2026-02-21T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "MEXBMI1-103",
                "summary": "Alta criticidad adicional 2",
                "status": "New",
                "priority": "High",
                "assignee": "Equipo Core",
                "created": "2026-02-22T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "MEXBMI1-104",
                "summary": "Alta criticidad adicional 3",
                "status": "New",
                "priority": "High",
                "assignee": "Equipo Core",
                "created": "2026-02-23T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "MEXBMI1-105",
                "summary": "Alta criticidad adicional 4",
                "status": "New",
                "priority": "High",
                "assignee": "Equipo Core",
                "created": "2026-02-24T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
        ]
    )
    settings = Settings(
        PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()),
        JIRA_BASE_URL="https://jira.example",
    )

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))

    trend_title = "Tendencia por funcionalidad : vista agregada ultimo semestre"
    trend_idx = _find_slide_index(prs, trend_title)
    dashboard_idx = _find_slide_index(
        prs,
        "Seguimiento de KPIs - Incidencias abiertas por funcionalidad",
    )
    high_cover_idx = _find_slide_index(prs, "Incidencias abiertas por criticidad alta")
    high_detail_idx = _find_slide_index(prs, "Incidencias abiertas por criticidad alta (I)")
    aged_cover_idx = _find_slide_index(prs, "Incidencias abiertas con más de 30 días")
    aged_detail_idx = _find_slide_index(prs, "Incidencias abiertas con más de 30 días (I)")
    high_detail_slides = [
        idx
        for idx, slide in enumerate(prs.slides)
        if "Incidencias abiertas por criticidad alta (" in _slide_text(slide)
    ]
    aged_detail_slides = [
        idx
        for idx, slide in enumerate(prs.slides)
        if "Incidencias abiertas con más de 30 días (" in _slide_text(slide)
    ]

    assert trend_idx < dashboard_idx < high_cover_idx < high_detail_idx < aged_cover_idx
    assert aged_cover_idx < aged_detail_idx
    assert len(high_detail_slides) >= 2
    assert len(aged_detail_slides) >= 2
    trend_title_shapes = [
        shape
        for shape in prs.slides[trend_idx].shapes
        if getattr(shape, "has_text_frame", False) and trend_title in str(shape.text or "")
    ]
    assert trend_title_shapes
    assert all("\n" not in str(shape.text or "") for shape in trend_title_shapes)
    assert not _native_tables(prs.slides[high_cover_idx])
    assert not _native_tables(prs.slides[aged_cover_idx])

    template_prs = Presentation(str(bundled_period_ppt_template_path()))
    section_cover_template = template_prs.slides[1]
    section_cover_title = section_cover_template.shapes[1]
    for cover_idx, expected_title in (
        (high_cover_idx, "Incidencias abiertas por criticidad alta"),
        (aged_cover_idx, "Incidencias abiertas con más de 30 días"),
    ):
        cover = prs.slides[cover_idx]
        cover_text = _slide_text(cover)
        assert "Haga clic para agregar" not in cover_text
        assert (
            cover.shapes[0].fill.fore_color.rgb
            == section_cover_template.shapes[0].fill.fore_color.rgb
        )
        assert cover.shapes[1].left == section_cover_title.left
        assert cover.shapes[1].top == section_cover_title.top
        assert cover.shapes[1].width == section_cover_title.width
        assert cover.shapes[1].height == section_cover_title.height
        assert cover.shapes[1].text == expected_title
        assert cover.shapes[1].text_frame.paragraphs[0].runs[0].font.name == (
            section_cover_title.text_frame.paragraphs[0].runs[0].font.name
        )

    expected_headers = {
        "ID",
        "Descripción",
        "Responsable",
        "Estado",
        "Criticidad",
        "Días abierta",
    }
    old_high_note = "Detalle ordenado por 1º : Criticidad, 2º: Días abierta y 3º: Estado"
    old_aged_note = "Detalle ordenado por 1º : Días abierta, 2º: Criticidad y 3º: Estado"
    for slide_idx in high_detail_slides:
        slide_text = _slide_all_text(prs.slides[slide_idx])
        assert period_ppt_mod._RISK_HIGH_PRIORITY_ORDER_NOTE in slide_text
        assert old_high_note not in slide_text
    for slide_idx in aged_detail_slides:
        slide_text = _slide_all_text(prs.slides[slide_idx])
        assert period_ppt_mod._RISK_AGED_ORDER_NOTE in slide_text
        assert old_aged_note not in slide_text
        assert "2ºCrriticidad" not in slide_text
    for slide_idx, expected_ids, expected_assignees, order_note, forbidden in (
        (
            high_detail_idx,
            ["EAM-77", "MEXBMI1-101"],
            ["Ana López", "Luis Pérez"],
            period_ppt_mod._RISK_HIGH_PRIORITY_ORDER_NOTE,
            "SKSEMEX-9",
        ),
        (
            aged_detail_idx,
            ["SKSEMEX-9", "MEXBMI1-101"],
            ["Marta Ruiz", "Luis Pérez"],
            period_ppt_mod._RISK_AGED_ORDER_NOTE,
            "",
        ),
    ):
        slide = prs.slides[slide_idx]
        slide_text = _slide_all_text(slide)
        assert "Detalle de incidencias abiertas:" not in slide_text
        assert order_note in slide_text
        tables = _native_tables(slide)
        assert len(tables) == 1
        table_shape = tables[0]
        assert not _table_intersects_picture(slide, table_shape)
        _assert_native_table_font_floor(table_shape)
        headers = {table_shape.table.cell(0, cidx).text for cidx in range(6)}
        assert headers == expected_headers
        table_text = _slide_table_text(slide)
        row_ids = [
            table_shape.table.cell(row_idx, 0).text for row_idx in range(1, len(expected_ids) + 1)
        ]
        assert row_ids == expected_ids
        for expected_assignee in expected_assignees:
            assert expected_assignee in table_text
        if forbidden:
            assert forbidden not in table_text
        assert "Descripción" in table_text
        assert "Criticidad" in table_text
        assert "días" in table_text


def test_period_followup_functionality_detail_toggle_off_omits_zoom_slides() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Pagos no refleja saldo",
                "status": "New",
                "priority": "High",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "B-1",
                "summary": "Login falla en acceso",
                "status": "Analysing",
                "priority": "Medium",
                "created": "2026-03-01T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))
    deck_text = " ".join(_slide_text(slide) for slide in prs.slides)

    assert len(prs.slides) == 15
    assert "Incidencias abiertas por criticidad alta" in deck_text
    assert "Incidencias abiertas con más de 30 días" in deck_text
    assert "Seguimiento de KPIs - Incidencias abiertas por funcionalidad" in deck_text
    assert "Incidencias, en Pagos, abiertas en la quincena" not in deck_text


def test_functionality_dashboard_table_headers_include_business_wording() -> None:
    assert period_ppt_mod._FUNCTIONALITY_DASHBOARD_TABLE_HEADERS == (
        "#",
        "Resto incidencias abiertas",
        "Nuevas",
        "Agregadas",
        "Días promedio abiertas",
    )
    weights = period_ppt_mod._FUNCTIONALITY_DASHBOARD_TABLE_COLUMN_WEIGHTS
    assert weights[1] < 3.05
    assert weights[3] > 1.12


def test_generate_country_period_followup_ppt_functionality_wording_depends_on_priority_filter() -> (
    None
):
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Pagos no refleja saldo",
                "status": "New",
                "priority": "High",
                "created": "2026-04-05T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Transferencias con timeout",
                "status": "Blocked",
                "priority": "Medium",
                "created": "2026-04-06T09:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))

    out_default = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs_default = Presentation(BytesIO(out_default.content))
    default_dashboard_idx = _find_slide_index(
        prs_default,
        "Seguimiento de KPIs - Incidencias abiertas por funcionalidad",
    )
    default_risk_idx = _find_slide_index(prs_default, "Incidencias abiertas por criticidad alta")
    blob_default = " ".join(
        str(getattr(shape, "text", "") or "")
        for idx in range(default_dashboard_idx, default_risk_idx)
        for shape in prs_default.slides[idx].shapes
        if getattr(shape, "has_text_frame", False)
    ).lower()
    assert "seguimiento de kpis - incidencias abiertas por funcionalidad" in blob_default
    assert "incidencias críticas" not in blob_default

    out_critical = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        functionality_priority_filters=["High", "Highest", "Supone un impedimento"],
        reference_day=now,
    )
    prs_critical = Presentation(BytesIO(out_critical.content))
    critical_dashboard_idx = _find_slide_index(
        prs_critical,
        "Seguimiento de KPIs - Incidencias críticas abiertas por funcionalidad",
    )
    critical_risk_idx = _find_slide_index(prs_critical, "Incidencias abiertas por criticidad alta")
    blob_critical = " ".join(
        str(getattr(shape, "text", "") or "")
        for idx in range(critical_dashboard_idx, critical_risk_idx)
        for shape in prs_critical.slides[idx].shapes
        if getattr(shape, "has_text_frame", False)
    ).lower()
    assert "seguimiento de kpis - incidencias críticas abiertas por funcionalidad" in blob_critical
    assert "incidencias críticas" in blob_critical


def test_period_followup_ppt_resolution_min_max_matches_closed_in_selected_fortnight() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "S-1",
                "summary": "Senda cerrada rápida",
                "status": "Resolved",
                "priority": "High",
                "created": "2026-04-07T08:00:00+00:00",  # 1 día
                "updated": now.isoformat(),
                "resolved": "2026-04-08T08:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "S-2",
                "summary": "Senda cerrada lenta",
                "status": "Resolved",
                "priority": "Medium",
                "created": "2026-03-20T08:00:00+00:00",  # 21 días
                "updated": now.isoformat(),
                "resolved": "2026-04-10T08:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "G-1",
                "summary": "Gema cerrada media",
                "status": "Resolved",
                "priority": "Low",
                "created": "2026-03-31T08:00:00+00:00",  # 10 días
                "updated": now.isoformat(),
                "resolved": "2026-04-10T08:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "S-OUT",
                "summary": "Senda cerrada fuera de quincena",
                "status": "Resolved",
                "priority": "High",
                "created": "2026-03-01T08:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": "2026-03-05T08:00:00+00:00",  # fuera de quincena actual
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "OPEN-1",
                "summary": "Incidencia abierta no debe contar",
                "status": "New",
                "priority": "High",
                "created": "2026-04-09T08:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))

    quincenal = period_ppt_mod.build_country_quincenal_result(
        df=dff,
        settings=settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        source_label_by_id=period_ppt_mod.source_label_map(
            settings,
            country="México",
            source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        ),
        reference_day=now,
    )
    expected_by_slide = {
        2: quincenal.aggregate.summary,
        3: quincenal.by_source["jira:mexico:senda"].summary,
        4: quincenal.by_source["jira:mexico:gema"].summary,
    }

    for slide_idx, summary in expected_by_slide.items():
        slide_blob = " ".join(
            str(getattr(shape, "text", "") or "")
            for shape in prs.slides[slide_idx].shapes
            if getattr(shape, "has_text_frame", False)
        )
        max_match = re.search(r"(\d+)\s*d[ií]as\s*MAX", slide_blob, flags=re.IGNORECASE)
        min_match = re.search(r"(\d+)\s*d[ií]as\s*MIN", slide_blob, flags=re.IGNORECASE)
        assert max_match is not None
        assert min_match is not None
        assert int(max_match.group(1)) == int(round(float(summary.resolution_days_max_now or 0.0)))
        assert int(min_match.group(1)) == int(round(float(summary.resolution_days_min_now or 0.0)))


def test_period_followup_summary_metric_cards_keep_template_blue_text_color() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Issue A",
                "status": "New",
                "priority": "High",
                "created": "2026-04-08T10:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Issue A2",
                "status": "Resolved",
                "priority": "Medium",
                "created": "2026-04-01T10:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": "2026-04-08T10:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "B-1",
                "summary": "Issue B",
                "status": "Resolved",
                "priority": "High",
                "created": "2026-04-02T10:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": "2026-04-10T10:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))
    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))

    def _first_run_color(shape: Any) -> RGBColor | None:
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        for paragraph in list(shape.text_frame.paragraphs):
            for run in list(paragraph.runs):
                if not str(getattr(run, "text", "") or "").strip():
                    continue
                color = getattr(getattr(run, "font", None), "color", None)
                rgb = getattr(color, "rgb", None)
                if rgb is not None:
                    return rgb
        return None

    for slide_idx in (2, 3, 4):
        slide = prs.slides[slide_idx]
        detail_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "en TOTAL" in str(shape.text or "")
        )
        closed_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "CERRADA" in str(shape.text or "")
        )
        resolution_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and "DÍAS DE RESOLUCIÓN" in str(shape.text or "")
        )
        created_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "CREADA" in str(shape.text or "")
        )
        reference_color = _first_run_color(detail_shape)
        assert reference_color == RGBColor(4, 19, 139)
        assert _first_run_color(closed_shape) == reference_color
        assert _first_run_color(resolution_shape) == reference_color
        assert _first_run_color(created_shape) == reference_color


def test_period_followup_summary_uses_quincenal_flow_wording_and_removes_artifacts(
    monkeypatch: Any,
) -> None:
    monkeypatch.setattr(period_ppt_mod, "_chart_png", lambda *args, **kwargs: b"")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Actual abierta",
                "status": "New",
                "priority": "High",
                "created": "2026-03-10T10:00:00+00:00",
                "updated": "2026-03-12T10:00:00+00:00",
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "A-2",
                "summary": "Previa cerrada",
                "status": "Resolved",
                "priority": "Medium",
                "created": "2026-02-20T10:00:00+00:00",
                "updated": "2026-03-11T10:00:00+00:00",
                "resolved": "2026-03-11T10:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "B-1",
                "summary": "Previa abierta",
                "status": "New",
                "priority": "High",
                "created": "2026-02-15T10:00:00+00:00",
                "updated": "2026-03-12T10:00:00+00:00",
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
            {
                "key": "B-2",
                "summary": "Actual cerrada",
                "status": "Resolved",
                "priority": "Highest",
                "created": "2026-03-01T10:00:00+00:00",
                "updated": "2026-03-12T10:00:00+00:00",
                "resolved": "2026-03-12T10:00:00+00:00",
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=pd.Timestamp("2026-03-12T00:00:00+00:00"),
    )

    prs = Presentation(BytesIO(out.content))
    deck_text = " ".join(_slide_all_text(slide) for slide in prs.slides)
    assert "CREADAS DEL 01 AL 12 MAR" in deck_text
    assert "2 del 15-28 FEB" in deck_text
    assert "4 en TOTAL" in deck_text
    assert "CERRADAS DEL 01-12 MAR" in deck_text
    assert "AHORA" not in deck_text
    assert "ACUMULADO" not in deck_text
    assert "NUEVAS INCIDENCIAS" not in deck_text
    assert "Ver detalle" not in deck_text

    for slide in prs.slides:
        for shape in slide.shapes:
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            right = left + int(getattr(shape, "width", 0) or 0)
            bottom = top + int(getattr(shape, "height", 0) or 0)
            assert left >= 0
            assert top >= 0
            assert right <= int(prs.slide_width)
            assert bottom <= int(prs.slide_height)

    with ZipFile(BytesIO(out.content)) as pptx:
        xml_payload = "\n".join(
            pptx.read(name).decode("utf-8", errors="ignore")
            for name in pptx.namelist()
            if name.endswith(".xml")
        )
    assert 'type="slidenum"' not in xml_payload
    assert "p. </a:t>" not in xml_payload


def test_period_followup_summary_uses_safe_central_delta_badges() -> None:
    reference_day = pd.Timestamp("2026-03-20T00:00:00+00:00")
    rows: list[dict[str, object]] = [
        {
            "key": "PREV-1",
            "summary": "Referencia mínima",
            "status": "Resolved",
            "priority": "High",
            "created": "2026-03-02T00:00:00+00:00",
            "updated": "2026-03-02T00:30:00+00:00",
            "resolved": "2026-03-02T00:30:00+00:00",
            "country": "México",
            "source_id": "jira:mexico:senda",
            "source_type": "jira",
        }
    ]
    for idx in range(15):
        created_day = 15 + (idx % 6)
        row: dict[str, object] = {
            "key": f"CUR-{idx + 1}",
            "summary": "Actual",
            "status": "New",
            "priority": "Medium",
            "created": f"2026-03-{created_day:02d}T00:00:00+00:00",
            "updated": f"2026-03-{created_day:02d}T00:00:00+00:00",
            "resolved": None,
            "country": "México",
            "source_id": "jira:mexico:senda" if idx % 2 == 0 else "jira:mexico:gema",
            "source_type": "jira",
        }
        if idx < 2:
            row["status"] = "Resolved"
            row["created"] = f"2026-03-{15 + idx:02d}T00:00:00+00:00"
            row["resolved"] = f"2026-03-{18 + idx:02d}T00:00:00+00:00"
        rows.append(row)
    dff = pd.DataFrame(rows)
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))

    quincenal = period_ppt_mod.build_country_quincenal_result(
        df=dff,
        settings=settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        reference_day=reference_day,
    )
    summary = quincenal.aggregate.summary
    assert summary.new_delta.display_kind == "absolute"
    assert period_ppt_mod._summary_delta_badge(summary.new_delta)[0] == "▲>100%"
    assert period_ppt_mod._summary_delta_badge(summary.resolution_delta)[0].endswith("%")

    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=reference_day,
    )
    prs = Presentation(BytesIO(out.content))
    deck_text = " ".join(_slide_all_text(slide) for slide in prs.slides)
    assert "▲>100%" in deck_text
    assert "1400%" not in deck_text
    assert "14673%" not in deck_text
    badge_pattern = re.compile(r"^[▲▼•](?:>?\d+%)$")
    for slide_idx in (2, 3, 4):
        slide = prs.slides[slide_idx]
        badge_texts = [
            str(getattr(shape, "text", "") or "").strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and badge_pattern.fullmatch(str(getattr(shape, "text", "") or "").strip())
        ]
        assert len(badge_texts) >= 3
        chart_pics = []
        for shape in slide.shapes:
            try:
                image = shape.image
            except Exception:
                continue
            if image.size[0] >= 1000:
                chart_pics.append((shape, image))
        assert chart_pics
        shape, image = max(chart_pics, key=lambda item: int(item[0].width) * int(item[0].height))
        shape_ratio = float(shape.width) / float(shape.height)
        image_ratio = float(image.size[0]) / float(image.size[1])
        assert abs(shape_ratio - image_ratio) < 0.03


def test_period_followup_metric_typography_handles_one_to_four_digit_values() -> None:
    values = ["9", "69", "107", "138", "1024"]
    sizes = [
        metric_card_typography(value, "CREADAS DEL 01 AL 15 MAR").value_size_pt for value in values
    ]

    assert min(sizes) >= 17.0
    assert sizes[0] >= sizes[-1]
    for value in values:
        typography = metric_card_typography(value, "CERRADAS DEL 01 AL 15 MAR")
        assert typography.label_size_pt >= 8.6
        assert typography.detail_size_pt >= 8.8


def test_period_followup_functionality_trend_title_matches_template_style() -> None:
    now = pd.Timestamp("2026-04-10T00:00:00+00:00")
    dff = pd.DataFrame(
        [
            {
                "key": "A-1",
                "summary": "Login falla en alta de usuario",
                "status": "New",
                "priority": "High",
                "created": "2026-04-08T10:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:senda",
            },
            {
                "key": "B-1",
                "summary": "Transferencias con timeout",
                "status": "Blocked",
                "priority": "Medium",
                "created": "2026-04-02T10:00:00+00:00",
                "updated": now.isoformat(),
                "resolved": None,
                "country": "México",
                "source_id": "jira:mexico:gema",
            },
        ]
    )
    settings = Settings(PERIOD_PPT_TEMPLATE_PATH=str(bundled_period_ppt_template_path()))
    out = generate_country_period_followup_ppt(
        settings,
        country="México",
        source_ids=["jira:mexico:senda", "jira:mexico:gema"],
        dff_override=dff,
        reference_day=now,
    )
    prs = Presentation(BytesIO(out.content))
    trend_slide = prs.slides[_find_slide_index(prs, "Tendencia por funcionalidad")]

    title_shape = next(
        (
            shape
            for shape in trend_slide.shapes
            if getattr(shape, "has_text_frame", False)
            and "Tendencia por funcionalidad" in str(getattr(shape, "text", "") or "")
        ),
        None,
    )
    assert title_shape is not None

    first_run = None
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if str(getattr(run, "text", "") or "").strip():
                first_run = run
                break
        if first_run is not None:
            break
    assert first_run is not None
    assert first_run.font.name == "Source Serif 4"
    assert first_run.font.color.rgb == RGBColor(4, 19, 139)
